"""
Phase 1: 마스터 템플릿 슬라이드 구조 분석.

대상: C:\\Users\\00\\Desktop\\제안서_마스터템플릿\\*.pptx
작업:
  - 슬라이드별 shape 타입 분포 (텍스트박스/표/차트/이미지/그룹/SmartArt 등)
  - AUTO 모드로 거버닝/소제목/본문 영역 자동 식별 (폰트 크기 기반)
  - 섹션 자동 분류 (텍스트 키워드 기반)
  - 결과: master_template_analysis.json + 사람이 읽을 보고서
"""
from __future__ import annotations
import json
import re
import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MASTER_DIR = Path(r"C:\Users\00\Desktop\제안서_마스터템플릿")
OUT_JSON = Path("master_template_analysis.json")
OUT_REPORT = Path("master_template_analysis_report.txt")

# 섹션 키워드
SECTION_PATTERNS = [
    ("표지",       [r"제\s*안\s*서", r"용역\s*제안", r"정성\s*제안", r"^contents$", r"^cover$"]),
    ("목차",       [r"^\s*목\s*차", r"^\s*contents", r"^\s*인덱스"]),
    ("사업이해",   [r"사업\s*이해", r"제안\s*배경", r"사업\s*개요", r"과업\s*이해", r"사업의\s*이해"]),
    ("추진전략",   [r"추진\s*전략", r"핵심\s*전략", r"차별화", r"제안\s*전략", r"전략\s*과제", r"수행\s*전략"]),
    ("수행조직",   [r"수행\s*조직", r"인력\s*투입", r"조직도", r"운영\s*인력", r"인력\s*구성"]),
    ("일정",       [r"추진\s*일정", r"운영\s*일정", r"로드맵", r"마일스톤", r"timeline", r"세부\s*일정"]),
    ("예산",       [r"소요\s*예산", r"산출\s*내역", r"예산\s*집행", r"단가표", r"예산\s*편성"]),
    ("프로그램",   [r"세부\s*프로그램", r"프로그램\s*구성", r"콘텐츠\s*기획", r"운영\s*프로그램"]),
    ("홍보",       [r"홍보\s*계획", r"홍보\s*전략", r"마케팅", r"확산"]),
    ("안전관리",   [r"안전\s*관리", r"위기\s*대응", r"비상\s*매뉴얼", r"안전\s*대책"]),
    ("기대효과",   [r"기대\s*효과", r"성과\s*확산", r"파급\s*효과", r"기대\s*성과"]),
    ("실적",       [r"수행\s*실적", r"유사\s*실적", r"주요\s*실적", r"제안사\s*소개"]),
]


def classify_section(slide_text: str, slide_idx: int) -> str:
    """섹션 자동 분류."""
    if slide_idx == 0:
        return "표지"
    for sec, pats in SECTION_PATTERNS:
        for pat in pats:
            if re.search(pat, slide_text, flags=re.IGNORECASE):
                return sec
    return "기타"


def shape_type_label(shape) -> str:
    """shape 타입 → 사람이 읽기 좋은 라벨."""
    try:
        st = shape.shape_type
    except Exception:
        return "unknown"
    if st == MSO_SHAPE_TYPE.PICTURE:        return "이미지"
    if st == MSO_SHAPE_TYPE.TABLE:          return "표"
    if st == MSO_SHAPE_TYPE.CHART:          return "차트"
    if st == MSO_SHAPE_TYPE.GROUP:          return "그룹"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:    return "placeholder"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:     return "도형"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:       return "텍스트박스"
    if st == MSO_SHAPE_TYPE.LINE:           return "선"
    if st == MSO_SHAPE_TYPE.FREEFORM:       return "freeform"
    if st == MSO_SHAPE_TYPE.MEDIA:          return "미디어"
    return f"기타({st})"


def extract_text_with_size(shape) -> list[dict]:
    """텍스트 영역 추출 — 폰트 크기 포함 (AUTO 모드용)."""
    out = []
    if not shape.has_text_frame:
        return out
    tf = shape.text_frame
    for p_idx, para in enumerate(tf.paragraphs):
        for run in para.runs:
            text = (run.text or "").strip()
            if not text:
                continue
            # 폰트 크기 (Pt 단위)
            sz = None
            try:
                if run.font.size is not None:
                    sz = int(run.font.size.pt)
            except Exception:
                pass
            # bold
            bold = False
            try:
                bold = bool(run.font.bold)
            except Exception:
                pass
            out.append({
                "text": text,
                "size_pt": sz,
                "bold": bold,
                "para": p_idx,
            })
    return out


def analyze_slide(slide, idx: int) -> dict:
    """슬라이드 1장 분석."""
    shapes_meta = []
    all_runs = []   # AUTO 모드 텍스트 분석용
    full_text_parts = []

    for shape in slide.shapes:
        meta = {
            "type": shape_type_label(shape),
            "name": shape.name,
            "has_text": shape.has_text_frame,
        }
        # 위치/크기 (EMU → cm)
        try:
            meta["pos"] = {
                "left_cm": round(shape.left / 360000, 2) if shape.left else 0,
                "top_cm":  round(shape.top / 360000, 2) if shape.top else 0,
                "w_cm":    round(shape.width / 360000, 2) if shape.width else 0,
                "h_cm":    round(shape.height / 360000, 2) if shape.height else 0,
            }
        except Exception:
            pass

        # 텍스트
        if shape.has_text_frame:
            runs = extract_text_with_size(shape)
            meta["runs"] = runs
            shape_text = " ".join(r["text"] for r in runs)
            if shape_text:
                full_text_parts.append(shape_text)
            for r in runs:
                if r["size_pt"]:
                    all_runs.append((r["size_pt"], r["text"], shape.name))

        # 표 — 행/열 수
        if meta["type"] == "표":
            try:
                meta["rows"] = len(shape.table.rows)
                meta["cols"] = len(shape.table.columns)
                # 표 안 모든 텍스트 (placeholder 마킹용)
                cells = []
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        ct = cell.text_frame.text if cell.text_frame else ""
                        if ct.strip():
                            cells.append({"r": r_idx, "c": c_idx, "text": ct.strip()[:80]})
                meta["cells"] = cells[:20]   # 최대 20개만
            except Exception:
                pass

        # 그룹 안 shape 수
        if meta["type"] == "그룹":
            try:
                meta["children"] = len(shape.shapes)
            except Exception:
                pass

        shapes_meta.append(meta)

    full_text = "\n".join(full_text_parts)
    section = classify_section(full_text, idx)

    # AUTO 모드 — 폰트 크기 기준 정렬해서 거버닝/소제목/본문 추정
    auto_text_zones = {"거버닝": None, "소제목": None, "본문": []}
    if all_runs:
        # 큰 순서로
        sorted_runs = sorted(all_runs, key=lambda x: -x[0])
        # 가장 큰 폰트
        largest_size = sorted_runs[0][0]
        # 거버닝 — 가장 큰 폰트 (단, 너무 짧은 건 제외 — 1글자 등)
        for sz, txt, name in sorted_runs:
            if sz == largest_size and len(txt) >= 3:
                auto_text_zones["거버닝"] = {"size": sz, "text": txt[:80], "shape": name}
                break
        # 소제목 — 두 번째 크기
        seen_sizes = set([largest_size])
        for sz, txt, name in sorted_runs:
            if sz < largest_size and sz not in seen_sizes:
                if len(txt) >= 3:
                    auto_text_zones["소제목"] = {"size": sz, "text": txt[:80], "shape": name}
                    break
        # 본문 — 그 외 작은 크기 텍스트 모음
        for sz, txt, name in sorted_runs[2:6]:
            auto_text_zones["본문"].append({"size": sz, "text": txt[:60], "shape": name})

    return {
        "idx": idx + 1,
        "section": section,
        "shape_count": len(shapes_meta),
        "shape_types": list({s["type"] for s in shapes_meta}),
        "type_breakdown": _count_types(shapes_meta),
        "auto_text_zones": auto_text_zones,
        "shapes": shapes_meta,
        "full_text_excerpt": full_text[:500],
    }


def _count_types(shapes_meta: list[dict]) -> dict:
    out = {}
    for s in shapes_meta:
        t = s["type"]
        out[t] = out.get(t, 0) + 1
    return out


def analyze_pptx(path: Path) -> dict:
    print(f"\n=== 분석: {path.name} ===", flush=True)
    print(f"  크기: {path.stat().st_size / 1024 / 1024:.1f} MB", flush=True)
    prs = Presentation(str(path))
    n_slides = len(prs.slides)
    print(f"  슬라이드: {n_slides}장", flush=True)

    # 페이지 사이즈 (EMU → cm)
    page_w_cm = round(prs.slide_width / 360000, 2)
    page_h_cm = round(prs.slide_height / 360000, 2)
    print(f"  페이지: {page_w_cm} × {page_h_cm} cm "
          f"({'A4 가로' if abs(page_w_cm/page_h_cm - 1.4142) < 0.1 else '기타'})", flush=True)

    slides_data = []
    for i, slide in enumerate(prs.slides):
        if i % 10 == 0:
            print(f"  분석 중 … {i+1}/{n_slides}", flush=True)
        try:
            slides_data.append(analyze_slide(slide, i))
        except Exception as e:
            print(f"  [슬라이드 {i+1} 분석 오류] {e}", flush=True)
            slides_data.append({"idx": i + 1, "error": str(e)[:200]})

    # 섹션별 집계
    section_counts = {}
    for s in slides_data:
        sec = s.get("section", "기타")
        section_counts[sec] = section_counts.get(sec, 0) + 1

    # 전체 shape 타입 집계
    all_types = {}
    for s in slides_data:
        for k, v in (s.get("type_breakdown") or {}).items():
            all_types[k] = all_types.get(k, 0) + v

    return {
        "filename": path.name,
        "size_mb": round(path.stat().st_size / 1024 / 1024, 1),
        "slide_count": n_slides,
        "page_w_cm": page_w_cm, "page_h_cm": page_h_cm,
        "section_counts": section_counts,
        "shape_type_totals": all_types,
        "slides": slides_data,
    }


def write_report(data: dict, out_path: Path):
    L = []
    L.append("마스터 템플릿 분석 보고서")
    L.append("=" * 70)
    L.append(f"파일: {data['filename']}")
    L.append(f"크기: {data['size_mb']} MB")
    L.append(f"슬라이드: {data['slide_count']}장")
    L.append(f"페이지 크기: {data['page_w_cm']} × {data['page_h_cm']} cm")
    L.append("")
    L.append("[섹션별 슬라이드 분포]")
    for sec, cnt in sorted(data["section_counts"].items(), key=lambda x: -x[1]):
        L.append(f"  {sec:10s}: {cnt}장")
    L.append("")
    L.append("[Shape 타입 총계]")
    for tp, cnt in sorted(data["shape_type_totals"].items(), key=lambda x: -x[1]):
        L.append(f"  {tp:14s}: {cnt}개")
    L.append("")
    L.append("[슬라이드별 상세 (요약)]")
    for s in data["slides"]:
        if "error" in s:
            L.append(f"  #{s['idx']:3d}  [오류] {s['error'][:100]}")
            continue
        types = s.get("type_breakdown") or {}
        types_str = ", ".join(f"{k}{v}" for k, v in sorted(types.items(), key=lambda x: -x[1]))
        L.append(f"  #{s['idx']:3d}  [{s.get('section','기타'):8s}] {s.get('shape_count',0):2d}개 shape  ({types_str})")
        # AUTO 모드 텍스트 영역
        atz = s.get("auto_text_zones") or {}
        if atz.get("거버닝"):
            g = atz["거버닝"]
            L.append(f"           거버닝({g['size']}pt): {g['text'][:60]}")
        if atz.get("소제목"):
            sub = atz["소제목"]
            L.append(f"           소제목({sub['size']}pt): {sub['text'][:60]}")
    out_path.write_text("\n".join(L), encoding="utf-8")
    print(f"\n  보고서: {out_path}", flush=True)


def main():
    if not MASTER_DIR.is_dir():
        print(f"[X] 폴더 없음: {MASTER_DIR}")
        sys.exit(1)
    pptx_files = sorted([f for f in MASTER_DIR.iterdir()
                         if f.is_file() and f.suffix.lower() == ".pptx"
                         and not f.name.startswith("~$")])
    if not pptx_files:
        print(f"[X] PPTX 파일 없음: {MASTER_DIR}")
        sys.exit(1)
    print(f"[OK] 발견: {len(pptx_files)}개 PPTX", flush=True)

    all_data = []
    for f in pptx_files:
        try:
            data = analyze_pptx(f)
            all_data.append(data)
        except Exception as e:
            print(f"[X] {f.name}: {e}", flush=True)
            all_data.append({"filename": f.name, "error": str(e)[:300]})

    # JSON 전체 저장
    OUT_JSON.write_text(
        json.dumps(all_data, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )
    print(f"\n  JSON: {OUT_JSON}")

    # 보고서 (사람용)
    if all_data:
        write_report(all_data[0], OUT_REPORT)


if __name__ == "__main__":
    main()
