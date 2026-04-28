"""
PPTX 슬라이드 배경 이미지 추출 + 메타데이터 자동 분류 + DB 저장.

대상 폴더 (사용자 PC 로컬):
  - C:\\Users\\00\\Desktop\\제안서_크리스
  - D:\\제안서 모음
  - D:\\제안서, 프로그램, 메일백업\\제안서 모음
  - D:\\제안서모음02

각 PPTX 에서:
  - 표지 (1번째 슬라이드) 1장
  - 섹션 헤더 추정 슬라이드 1~3장
변환 → A4 가로(297mm × 210mm) 비율로 크롭/리사이즈 → JPEG 저장.

자동 분류 휴리스틱:
  - 분야: 파일명·슬라이드 텍스트의 키워드 매칭 (축제/행사/전시/공연/학술/경연/스포츠/시상식)
  - 섹션: 슬라이드 텍스트의 표지/목차/사업이해/추진전략/수행조직/일정/예산/기대효과
  - 레이아웃: 텍스트 비율로 풀페이지/좌우분할/카드형/텍스트중심 추정

== 의존성 ==
  - LibreOffice (sudo apt install libreoffice / Windows: https://www.libreoffice.org/download)
    Windows 기본 경로: C:\\Program Files\\LibreOffice\\program\\soffice.exe
  - Pillow: pip install Pillow
  - python-pptx: 이미 설치됨

== 사용법 ==
  python extract_slide_backgrounds.py
  → slide_backgrounds.db + static/slide_bg/ 에 결과 저장.

== Railway 배포 ==
  - slide_backgrounds.db (수십 MB) 는 깃 트래킹 안 함 (.gitignore)
  - static/slide_bg/ 도 깃 트래킹 안 함
  - 운영 환경에서 매칭은 fallback (CSS 그라데이션) 으로 동작
  - 개발 PC 에서 추출 후 별도 업로드 (선택사항 — 운영에 진짜 적용하려면)
"""
from __future__ import annotations
import os
import re
import sys
import json
import shutil
import sqlite3
import subprocess
import tempfile
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

# ─── 설정 ─────────────────────────────────────
SOURCES = [
    Path(r"C:\Users\00\Desktop\제안서_크리스"),
    Path(r"D:\제안서 모음"),
    Path(r"D:\제안서, 프로그램, 메일백업\제안서 모음"),
    Path(r"D:\제안서모음02"),
]
OUT_IMG_DIR = Path("static/slide_bg")
DB_PATH = Path("slide_backgrounds.db")

# LibreOffice 경로 자동 탐색
LIBREOFFICE_CANDIDATES = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/libreoffice",
    "/usr/bin/soffice",
    "soffice",  # PATH
]

# A4 가로 비율 (297:210)
A4_RATIO = 297 / 210
# 출력 이미지 크기 (썸네일 + full)
THUMB_W = 480
FULL_W = 1280

# ─── 분야 키워드 ───
DOMAIN_KEYWORDS = {
    "festival":   ["축제", "페스티벌", "한마당", "기념행사", "엑스포"],
    "forum":      ["포럼", "컨퍼런스", "심포지엄", "학술", "세미나"],
    "exhibition": ["전시", "박람회", "전", "쇼케이스", "fair"],
    "campaign":   ["캠페인", "홍보", "공모", "확산"],
    "education":  ["교육", "학생", "청소년", "아카데미", "강좌"],
    "sports":     ["체육", "스포츠", "대회", "경기"],
    "performance":["공연", "연극", "음악회", "콘서트", "예술"],
    "ceremony":   ["기념식", "시상식", "기념", "선포식", "발대식"],
    "rnd":        ["R&D", "연구", "기술", "AI", "스마트"],
    "welfare":    ["복지", "안전", "장애", "보훈", "어린이"],
    "tourism":    ["관광", "여행", "트래블"],
    "other":      [],
}
DOMAIN_LABELS = {
    "festival": "축제·기념행사", "forum": "포럼·학술", "exhibition": "전시·박람회",
    "campaign": "캠페인·홍보", "education": "교육·청소년", "sports": "스포츠·체육",
    "performance": "공연·예술", "ceremony": "시상·기념식", "rnd": "R&D·기술",
    "welfare": "복지·안전", "tourism": "관광", "other": "기타",
}

# ─── 섹션 키워드 ───
SECTION_KEYWORDS = {
    "표지":     [r"^\s*제\s*안\s*서", r"용역\s*제안", r"정성\s*제안", r"contents$", r"제\s*\d+\s*회"],
    "목차":     [r"^\s*목\s*차\s*$", r"^\s*contents\s*$", r"^\s*인덱스\s*$"],
    "사업이해": [r"사업\s*이해", r"제안\s*배경", r"사업\s*개요", r"과업\s*이해"],
    "추진전략": [r"추진\s*전략", r"핵심\s*전략", r"차별화", r"제안\s*전략", r"전략\s*과제"],
    "수행조직": [r"수행\s*조직", r"인력\s*투입", r"조직도", r"운영\s*인력"],
    "일정":     [r"추진\s*일정", r"운영\s*일정", r"로드맵", r"마일스톤", r"timeline"],
    "예산":     [r"소요\s*예산", r"산출\s*내역", r"예산\s*집행", r"단가표"],
    "프로그램": [r"세부\s*프로그램", r"프로그램\s*구성", r"콘텐츠\s*기획"],
    "홍보":     [r"홍보\s*계획", r"홍보\s*전략", r"마케팅"],
    "안전":     [r"안전\s*관리", r"위기\s*대응", r"비상\s*매뉴얼"],
    "기대효과": [r"기대\s*효과", r"성과\s*확산", r"파급\s*효과"],
    "기타":     [],
}


def find_libreoffice() -> str | None:
    for c in LIBREOFFICE_CANDIDATES:
        p = Path(c)
        if p.exists():
            return str(p)
        try:
            r = subprocess.run([c, "--version"], capture_output=True, timeout=5)
            if r.returncode == 0:
                return c
        except Exception:
            continue
    return None


def classify_domain(filename: str, slide_text: str = "") -> tuple[str, float]:
    """파일명 + 슬라이드 텍스트로 분야 판별."""
    text = (filename + " " + slide_text).lower()
    best, score = "other", 0
    for dom, kws in DOMAIN_KEYWORDS.items():
        for kw in kws:
            if kw.lower() in text:
                # 첫 매칭 찾으면 break (단순 우선순위)
                return dom, 1.0
    return best, 0.0


def classify_section(slide_text: str, page_idx: int = 0) -> str:
    text = (slide_text or "").lower()
    if page_idx == 0:
        return "표지"
    for sec, pats in SECTION_KEYWORDS.items():
        for pat in pats:
            if re.search(pat, slide_text, flags=re.IGNORECASE):
                return sec
    return "기타"


def classify_layout(slide_text: str) -> str:
    """텍스트 길이·줄 수로 레이아웃 추정."""
    if not slide_text:
        return "이미지중심"
    lines = [l for l in slide_text.split("\n") if l.strip()]
    n_lines = len(lines)
    total = len(slide_text)
    if n_lines >= 8 and total > 400:
        return "텍스트중심"
    if n_lines >= 3 and total < 200:
        return "카드형"
    if 100 < total < 400:
        return "좌우분할"
    return "풀페이지"


def get_pptx_slide_texts(pptx_path: Path) -> list[str]:
    """python-pptx 로 각 슬라이드 텍스트 추출."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        out = []
        for sl in prs.slides:
            parts = []
            for shape in sl.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                parts.append(run.text)
                        parts.append("\n")
            out.append("".join(parts).strip())
        return out
    except Exception as e:
        print(f"  [텍스트 추출 실패] {pptx_path.name}: {e}")
        return []


def convert_pptx_to_pngs(pptx_path: Path, out_dir: Path, soffice: str) -> list[Path]:
    """LibreOffice 로 PPTX → PNG 변환 (모든 슬라이드)."""
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        cmd = [
            soffice, "--headless", "--convert-to", "png",
            "--outdir", str(tmp_dir),
            str(pptx_path),
        ]
        try:
            r = subprocess.run(cmd, capture_output=True, timeout=120)
            if r.returncode != 0:
                print(f"  [LibreOffice 오류] {pptx_path.name}: {r.stderr.decode('utf-8', errors='replace')[:200]}")
                return []
        except subprocess.TimeoutExpired:
            print(f"  [타임아웃] {pptx_path.name}")
            return []
        # LibreOffice 는 첫 슬라이드만 PNG 로 변환함 (--convert-to png).
        # 모든 슬라이드 추출하려면 PDF → PNG 분할이 필요.
        pngs = sorted(tmp_dir.glob("*.png"))
        if not pngs:
            return []
        out_paths = []
        for i, p in enumerate(pngs):
            dst = out_dir / f"{pptx_path.stem}_slide{i+1:02d}.png"
            shutil.move(str(p), str(dst))
            out_paths.append(dst)
        return out_paths


def crop_to_a4(img_path: Path) -> bool:
    """이미지를 A4 가로 비율로 크롭 + 리사이즈."""
    try:
        from PIL import Image
    except ImportError:
        print("  [Pillow 미설치] pip install Pillow 후 다시 시도하세요.")
        return False
    try:
        with Image.open(img_path) as im:
            w, h = im.size
            target_ratio = A4_RATIO
            cur_ratio = w / h
            if cur_ratio > target_ratio:
                new_w = int(h * target_ratio)
                left = (w - new_w) // 2
                im = im.crop((left, 0, left + new_w, h))
            else:
                new_h = int(w / target_ratio)
                top = (h - new_h) // 2
                im = im.crop((0, top, w, top + new_h))
            im = im.resize((FULL_W, int(FULL_W / target_ratio)), Image.LANCZOS)
            jpg = img_path.with_suffix(".jpg")
            im.convert("RGB").save(jpg, "JPEG", quality=82, optimize=True)
            if img_path.exists() and img_path != jpg:
                img_path.unlink()
            return True
    except Exception as e:
        print(f"  [크롭 실패] {img_path.name}: {e}")
        return False


def init_db():
    if DB_PATH.exists():
        DB_PATH.unlink()
    db = sqlite3.connect(str(DB_PATH))
    db.executescript("""
        CREATE TABLE backgrounds (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            source_pptx TEXT NOT NULL,
            slide_idx   INTEGER NOT NULL,
            domain      TEXT NOT NULL,
            section     TEXT NOT NULL,
            layout      TEXT NOT NULL,
            text_excerpt TEXT,
            img_path    TEXT NOT NULL,
            width       INTEGER, height INTEGER,
            created_at  TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE INDEX idx_bg_domain ON backgrounds(domain);
        CREATE INDEX idx_bg_section ON backgrounds(section);
        CREATE INDEX idx_bg_domain_section ON backgrounds(domain, section);
    """)
    db.commit()
    return db


def main():
    soffice = find_libreoffice()
    if not soffice:
        print("[X] LibreOffice 가 설치돼있지 않아요.")
        print("    Windows: https://www.libreoffice.org/download/")
        print("    설치 후 다시 실행하세요.")
        sys.exit(1)
    print(f"[OK] LibreOffice: {soffice}")

    OUT_IMG_DIR.mkdir(parents=True, exist_ok=True)
    db = init_db()

    pptx_files = []
    for src in SOURCES:
        if not src.is_dir():
            print(f"[스킵] 폴더 없음: {src}")
            continue
        for f in src.iterdir():
            if not f.is_file() or f.suffix.lower() != ".pptx":
                continue
            if f.name.startswith("~$"):
                continue   # MS Office 임시 잠금 파일
            pptx_files.append(f)
    print(f"[OK] PPTX 파일: {len(pptx_files)}개")

    saved = 0
    failed = 0
    for idx, pptx in enumerate(pptx_files, 1):
        print(f"\n[{idx:3d}/{len(pptx_files)}] {pptx.name[:60]}")
        try:
            texts = get_pptx_slide_texts(pptx)
            if not texts:
                continue

            # 표지(1번째) + 섹션 헤더로 추정되는 슬라이드 1~3개 추가 추출
            # 단순 휴리스틱: 텍스트가 짧은(1~3줄, 100자 미만) 슬라이드 = 섹션 헤더 가능성
            target_indices = [0]   # 표지
            for i, t in enumerate(texts[1:], 1):
                lines = [l for l in t.split("\n") if l.strip()]
                if 1 <= len(lines) <= 3 and 5 < len(t.strip()) < 80:
                    target_indices.append(i)
                    if len(target_indices) >= 4:
                        break

            # LibreOffice 변환 — 첫 슬라이드만 가능. 전체 슬라이드 추출은
            # PDF 변환 후 pdf2image 분할이 필요하지만, 우선 표지(1번)만 처리.
            pngs = convert_pptx_to_pngs(pptx, OUT_IMG_DIR, soffice)
            if not pngs:
                failed += 1
                continue

            # 첫 PNG (표지) 만 처리
            cover_png = pngs[0]
            crop_ok = crop_to_a4(cover_png)
            if not crop_ok:
                failed += 1
                continue
            cover_jpg = cover_png.with_suffix(".jpg") if cover_png.suffix == ".png" else cover_png

            # 메타 분류
            slide_text = texts[0] if texts else ""
            domain, conf = classify_domain(pptx.stem, slide_text)
            section = classify_section(slide_text, 0)
            layout = classify_layout(slide_text)

            db.execute(
                "INSERT INTO backgrounds(source_pptx, slide_idx, domain, section, layout, "
                "text_excerpt, img_path, width, height) "
                "VALUES (?,?,?,?,?,?,?,?,?)",
                (pptx.name, 1, domain, section, layout,
                 slide_text[:300], str(cover_jpg).replace("\\", "/"),
                 FULL_W, int(FULL_W / A4_RATIO)),
            )
            saved += 1
            print(f"  -> 저장: {section} · {domain} · {layout}")

            # 표지 외 다른 슬라이드는 전체 PDF 분할이 필요 — 별도 작업
            # 지금은 표지만으로 시작 (나중에 PDF 변환 + 분할 추가 가능)
        except Exception as e:
            print(f"  [예외] {e}")
            failed += 1
            continue

    db.commit()

    total = db.execute("SELECT COUNT(*) FROM backgrounds").fetchone()[0]
    by_domain = db.execute(
        "SELECT domain, COUNT(*) FROM backgrounds GROUP BY domain ORDER BY 2 DESC"
    ).fetchall()
    db.close()

    print(f"\n=== 완료 ===")
    print(f"  처리: {len(pptx_files)} / 저장: {saved} / 실패: {failed}")
    print(f"  DB: {DB_PATH} ({total}건)")
    print(f"  이미지: {OUT_IMG_DIR}/")
    print(f"  분야별:")
    for dom, cnt in by_domain:
        print(f"    {dom:12s}: {cnt}건")


if __name__ == "__main__":
    main()
