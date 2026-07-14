"""
NightOff PPTX Generator — 마스터 템플릿 기반 제안서 생성

핵심 철학:
  1. 마스터 PPTX (사람이 만든 잘 만든 제안서) 를 통째로 파일 복사
  2. 그 사본의 텍스트만 새 내용으로 치환 (배경/도형/표/차트/이미지/시각화 100% 보존)
  3. 사용 안 하는 슬라이드는 삭제
  4. 결과 = 사람이 만든 것처럼 보이는 PPTX

핵심 함수:
  - generate_from_master(master_path, content_per_slide, output_path, keep_indices)
  - extract_text_zones(slide) — AUTO 모드 텍스트 영역 식별
  - replace_text_in_slide(slide, content) — AUTO 모드 텍스트 치환

content_per_slide 형식:
  {
    0: {"거버닝": "...", "소제목": "...", "본문": ["...", "..."]},
    3: {"거버닝": "...", "본문": [...]},
    ...
  }
"""
from __future__ import annotations
import shutil
import logging
from copy import deepcopy
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# Spec D-Build-ThemeConnect (1-b) — 다크 테마 토대 연결 (theme.py 토큰 사용).
# 1-a(7c3ca11) 에서 정의만 한 토큰을 generate_from_shape_json 에서 처음 import.
# theme="light" 일 때 토큰값(#1A1A1A 등) = 현재 운영 색과 정확 일치 → 라이트 동작 무변경.
try:
    from theme import get_theme as _get_theme  # noqa: F401
except Exception:  # pragma: no cover — theme.py 누락 환경 (실험 폴더 등) 대비
    def _get_theme(name: str = "light") -> dict:
        # 안전 fallback — 라이트 팔레트만 인라인 (다크 호출되면 라이트로 떨어짐).
        return {
            "BG":          "#FFFFFF",
            "FG":          "#1A1A1A",
            "FG_SUB":      "#444444",
            "FG_META":     "#999999",
            "LINE":        "#DDDDDD",
            "PANEL_FILL":  "#1A1A1A",
            "PANEL_FG":    "#FFFFFF",
            "ACCENT":      "#1A1A1A",
        }


# ─── Spec D-Build-ThemeColorMap (1-c, 옵션 A) — render 단계 색 매핑 ────────────
# 배경: preset/inject/SLIDE_SYSTEM_PROMPT 가 전부 흑백 6색으로 그리고, 텍스트의 94%
#       가 color 를 명시. 1-b 의 "기본값만" 패치로는 다크 모드에서 글자 대부분이
#       검정 그대로 → 검정 배경에 묻혀 안 보임. 이 모듈은 명시된 색을 render 순간에
#       role(text/fill/stroke) 별로 다크 매핑한다. theme!='dark' 일 때 입력색 그대로.
#
# ★ light 무변경 보장: _map_color 는 theme != "dark" 면 hex_str 을 그대로 반환 →
#   라이트 경로는 패치 전후 비트 단위 동일 (split/asymmetric 의 검정면+흰글씨도
#   기존과 똑같이 출력).
# ★ role 구분 필수: 같은 #1A1A1A 가 fill(검정 면)로 쓰일 때와 text(검정 글자)로
#   쓰일 때 다르게 매핑돼야 split 의 검정면+흰글씨 ↔ 흰면+검정글씨 자연 반전 성립.
# ★ DARK_MAP 에 없는 색은 그대로 통과 (임의 변환 금지) — 매핑 표가 진실의 단일 출처.

def _norm_hex(h):
    """hex 정규화 — 대문자 + 6자리 풀어쓰기 (#FFF → #FFFFFF). DARK_MAP 키 매칭용."""
    if not h:
        return h
    h = str(h).upper().lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return "#" + h


DARK_MAP = {
    # 글자색 (role='text') — 검정 배경 위에서 보이도록 반전 / 회색 단계는 토큰 매핑.
    ("#1A1A1A", "text"): "#FFFFFF",   # 본문/헤드라인 검정 → 흰
    ("#444444", "text"): "#DDDDDD",   # 부차 진회색 → FG_SUB
    ("#666666", "text"): "#DDDDDD",   # 부차 중회색 → FG_SUB
    ("#999999", "text"): "#888888",   # 메타/캡션 옅은 회색 → FG_META
    ("#BBBBBB", "text"): "#888888",   # eyebrow 옅은 → FG_META
    ("#333333", "text"): "#DDDDDD",   # 진회색 글자 → FG_SUB
    ("#555555", "text"): "#DDDDDD",   # 진회색 글자 → FG_SUB
    ("#FFFFFF", "text"): "#1A1A1A",   # 검정 면 위 흰 글자(split 좌측) → 흰 면 위 검정 글자
    ("#DDDDDD", "text"): "#444444",   # 검정 면 위 부차 글자(split 좌측 pts) → 흰 면 위 부차 글자
    # 채움 (role='fill') — 면 강조 반전 / 흰 면은 다크 배경에 묻히도록 BG 와 동일.
    ("#1A1A1A", "fill"): "#FFFFFF",   # 강조 검정 면(split 좌측 / asymmetric 우측) → 흰 면
    ("#000000", "fill"): "#FFFFFF",   # 같은 의미의 #000 → 흰 면
    ("#FFFFFF", "fill"): "#0A0A0A",   # 흰 면(split 우측) → 다크 배경 일치 (자연스럽게 묻힘)
    ("#F5F5F5", "fill"): "#1F1F1F",   # 이미지 placeholder 옅은 회색 면 → 어두운 회색
    ("#FAFAFA", "fill"): "#1F1F1F",   # 같은 의미의 옅은 회색 면 → 어두운 회색
    ("#ECECEC", "fill"): "#1F1F1F",   # 옅은 회색 면 → 어두운 회색
    ("#DDDDDD", "fill"): "#2A2A2A",   # 연한 회색 면 → 다크 LINE
    # 선/테두리 (role='stroke') — 검정 선은 흰 선 / 회색 선은 다크용 회색.
    ("#1A1A1A", "stroke"): "#FFFFFF", # 진한 구분선/말풍선 테두리 → 흰
    ("#DDDDDD", "stroke"): "#2A2A2A", # 연한 구분선 → 다크 LINE
    ("#999999", "stroke"): "#666666", # 중간 회색 선 → 살짝 어둡게
    ("#CCCCCC", "stroke"): "#444444", # placeholder 테두리 → 어두운 회색
}


def _map_color(hex_str, role, theme="light"):
    """role-aware 다크 매핑.

    theme != "dark" → 입력색 그대로 (★ 라이트 비트 단위 무변경 보장).
    hex_str falsy(None/"") → 그대로 (fill=None 같은 "투명/지우기" 시그널 보존).
    DARK_MAP 에 없는 색 → 그대로 (임의 변환 금지 / 명시색 무변환 = 안전).
    """
    if theme != "dark" or not hex_str:
        return hex_str
    return DARK_MAP.get((_norm_hex(hex_str), role), hex_str)


log = logging.getLogger("pptx_gen")


# ─── Paperlogy weight 9 단계 매핑 ──────────────────────────────────────────────
# 도형 JSON 의 weight (100~900) 값을 받아 PowerPoint 의 typeface 명으로 매핑.
# 폰트 이름은 ttf internal Family Name (nameID 1) — Windows PowerPoint /
# LibreOffice 가 표준으로 매칭하는 형식. PostScript name (Paperlogy-1Thin)
# 이 아닌 Family Name (Paperlogy 1 Thin, 스페이스 포함) 으로 정확히 박아야
# 폰트 매칭 실패 시 기본 폰트 (맑은 고딕 등) 로 폴백되지 않는다.
WEIGHT_FONT_MAP: dict[int, str] = {
    100: "Paperlogy 1 Thin",
    200: "Paperlogy 2 ExtraLight",
    300: "Paperlogy 3 Light",
    400: "Paperlogy 4 Regular",
    500: "Paperlogy 5 Medium",
    600: "Paperlogy 6 SemiBold",
    700: "Paperlogy 7 Bold",
    800: "Paperlogy 8 ExtraBold",
    900: "Paperlogy 9 Black",
}
DEFAULT_FONT_FAMILY = WEIGHT_FONT_MAP[400]  # "Paperlogy 4 Regular"


def _normalize_weight(weight) -> int:
    """임의 weight 값 → 가장 가까운 100 단위 (100~900) 로 정규화.

    안전성 ↑:
    - float 입력 안전 (예: 400.5 → 400)
    - 영역 외 값 clamp (예: 1000 → 900, 50 → 100)
    - None / 잘못된 형식 → 400 (Regular fallback)
    """
    if weight is None or weight == "":
        return 400
    try:
        w = int(float(weight))
    except (TypeError, ValueError):
        return 400
    # 100~900 범위 clamp 영역 round
    w = max(100, min(900, w))
    return round(w / 100) * 100


def _resolve_font(font_family, weight) -> str:
    """font_family 명시 시 우선, 미지정 시 weight 기반 자동 매핑.

    fallback 강화 (가설 A fix):
    - font_family 명시 + "Paperlogy" 시작 → WEIGHT_FONT_MAP 표준 형식 강제 정규화
      (AI 가 prompt 예시 따라 "Paperlogy 9Black" (스페이스 없음) 명시해도
       표준 형식 "Paperlogy 9 Black" (스페이스 있음) 으로 정상 매핑 — ttf
       Family Name 매칭 실패로 시스템 fallback 되던 사고 차단)
    - font_family 명시 + 다른 폰트 (사용자 의도) → 그대로 사용 (호환성 유지)
    - 미지정 → weight 기반 자동 매핑
    - weight 매핑 X → DEFAULT (Paperlogy 4 Regular)
    """
    if font_family:
        f = str(font_family).strip()
        # Paperlogy 변형 명시 시 → WEIGHT_FONT_MAP 표준 형식 강제
        if f.lower().startswith("paperlogy"):
            return WEIGHT_FONT_MAP.get(_normalize_weight(weight), DEFAULT_FONT_FAMILY)
        return f  # 다른 폰트 명시 (사용자 의도) → 그대로
    norm_weight = _normalize_weight(weight)
    return WEIGHT_FONT_MAP.get(norm_weight, DEFAULT_FONT_FAMILY)


# ─── AUTO 모드 텍스트 영역 식별 ───────────────────────────────

def _run_size_pt(run) -> float:
    """run 의 폰트 크기 pt 단위 (없으면 기본 12)."""
    try:
        if run.font.size is not None:
            return float(run.font.size.pt)
    except Exception:
        pass
    return 12.0


def extract_text_zones(slide) -> dict:
    """슬라이드의 모든 text run 을 폰트 크기로 분류.

    반환:
      {
        "all_runs":    [{"run", "size", "text", "shape_name"}, ...],
        "largest_size": float,
        "second_size": float,
        "candidates":  {"거버닝": [...], "소제목": [...], "본문": [...]}
      }
    """
    runs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = (run.text or "").strip()
                if not txt:
                    continue
                runs.append({
                    "run": run, "size": _run_size_pt(run),
                    "text": txt, "shape_name": shape.name,
                })
    if not runs:
        return {"all_runs": [], "largest_size": 0, "second_size": 0,
                "candidates": {"거버닝": [], "소제목": [], "본문": []}}

    # 폰트 크기 내림차순 정렬
    sorted_by_size = sorted(runs, key=lambda x: -x["size"])
    largest = sorted_by_size[0]["size"]

    # 두 번째 크기 (가장 큰 폰트와 다른 첫 번째)
    second = next((r["size"] for r in sorted_by_size if r["size"] < largest), largest)

    # 카테고리별 후보
    candidates = {"거버닝": [], "소제목": [], "본문": []}
    for r in runs:
        if r["size"] == largest and len(r["text"]) >= 3:
            candidates["거버닝"].append(r)
        elif r["size"] == second and len(r["text"]) >= 3:
            candidates["소제목"].append(r)
        else:
            candidates["본문"].append(r)

    return {
        "all_runs": runs,
        "largest_size": largest,
        "second_size": second,
        "candidates": candidates,
    }


# ─── 텍스트 안전 치환 (폰트 스타일 유지) ─────────────────────

def _replace_run_text(run, new_text: str):
    """run 의 텍스트만 교체. 폰트 스타일(size/bold/color) 유지."""
    try:
        run.text = new_text
    except Exception as e:
        log.warning("run.text 치환 실패: %s", e)


def _replace_text_frame_simple(text_frame, new_text: str):
    """text_frame 의 모든 텍스트를 단일 텍스트로 교체.
    첫 paragraph 의 첫 run 폰트만 살리고 나머지 다 제거.
    """
    if not text_frame.paragraphs:
        return
    first_p = text_frame.paragraphs[0]
    if not first_p.runs:
        # run 없으면 paragraph 의 text 만 교체
        first_p.text = new_text
        return
    # 추가 paragraph 제거
    p_elements = text_frame._txBody.findall(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}p")
    for p in p_elements[1:]:
        try:
            p.getparent().remove(p)
        except Exception:
            pass
    # 첫 paragraph 의 추가 run 제거
    runs_to_remove = list(first_p.runs)[1:]
    for r in runs_to_remove:
        try:
            r._r.getparent().remove(r._r)
        except Exception:
            pass
    # 첫 run 의 텍스트만 변경
    first_p.runs[0].text = new_text


def _clear_text_frame(text_frame):
    """text_frame 의 모든 paragraph + run 의 텍스트를 비움 (폰트 스타일은 유지).
    첫 paragraph 의 첫 run 만 남기고 나머지 모두 제거. 첫 run 의 text 도 ""."""
    if not text_frame or not text_frame.paragraphs:
        return
    NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    # 추가 paragraph 모두 제거
    p_elements = text_frame._txBody.findall(f".//{NS_A}p")
    for p in p_elements[1:]:
        try:
            p.getparent().remove(p)
        except Exception:
            pass
    # 첫 paragraph 의 추가 run 제거
    first_p = text_frame.paragraphs[0]
    if first_p.runs:
        runs_to_remove = list(first_p.runs)[1:]
        for r in runs_to_remove:
            try:
                r._r.getparent().remove(r._r)
            except Exception:
                pass
        # 첫 run 의 텍스트 비움
        try:
            first_p.runs[0].text = ""
        except Exception:
            pass
    else:
        # run 없는 paragraph — 그냥 text 비움
        try:
            first_p.text = ""
        except Exception:
            pass


def _frame_max_font_size(text_frame) -> float:
    """text_frame 안의 최대 폰트 사이즈 (pt). 비어있으면 0."""
    max_sz = 0.0
    for p in text_frame.paragraphs:
        for r in p.runs:
            if (r.text or "").strip():
                max_sz = max(max_sz, _run_size_pt(r))
    return max_sz


def fill_slide_clearing_master(slide, content: dict) -> dict:
    """[옵션 A 핵심] 마스터의 모든 텍스트를 비우고 AI 콘텐츠로 채움.

    매핑 전략 (폰트 사이즈 내림차순):
      - 가장 큰 사이즈 frame    → 거버닝 (governing)
      - 두 번째 큰 frame        → 소제목 (subtitle)
      - 그 외 frame들           → 본문 (body) 항목 순서대로 분배
      - AI 콘텐츠 부족 시       → 빈 문자열 (디자인만 남김. 짬뽕 방지)

    기존 replace_text_in_slide 와의 차이:
      - 기존: 큰 글자만 치환, 본문 박스는 마스터 원본 그대로 (= 짬뽕)
      - 신규: 모든 frame 비우고 AI 콘텐츠로 정확히 매핑 (= 깨끗)

    content 형식:
      {"거버닝": "...", "소제목": "...", "본문": ["...", "..."], "summary": "..."}
    """
    # 1. 슬라이드의 모든 text_frame 수집 + 폰트 사이즈
    frames = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # 비어있는 frame 도 일단 포함 (디자인 자리 표시일 수 있음)
        frames.append({
            "shape": shape,
            "tf": tf,
            "size": _frame_max_font_size(tf),
            "had_text": bool(tf.text and tf.text.strip()),
        })

    if not frames:
        return {"replaced": 0, "cleared": 0, "errors": ["no text frames"]}

    # 2. 사이즈 내림차순 정렬 — 큰 글자가 거버닝/소제목 후보
    frames.sort(key=lambda f: -f["size"])

    # 3. 모든 frame 의 텍스트 비우기 (마스터 원본 잔재 0)
    cleared = 0
    for f in frames:
        if f["had_text"]:
            try:
                _clear_text_frame(f["tf"])
                cleared += 1
            except Exception as e:
                log.warning("frame clear 실패 (shape=%s): %s", f["shape"].name, e)

    # 4. AI 콘텐츠 정리
    governing = (content.get("거버닝") or "").strip() if content else ""
    subtitle = (content.get("소제목") or "").strip() if content else ""
    body_raw = content.get("본문") if content else None
    if isinstance(body_raw, str):
        body = [body_raw.strip()] if body_raw.strip() else []
    elif isinstance(body_raw, list):
        body = [str(b).strip() for b in body_raw if str(b).strip()]
    else:
        body = []
    summary = (content.get("summary") or "").strip() if content else ""

    # 5. 매핑 — 사이즈 내림차순 frame 들에 콘텐츠 채움
    # had_text 가 True 인 frame 만 채움 (디자인 전용 빈 frame 은 그대로)
    fillable = [f for f in frames if f["had_text"]]

    replaced = 0
    fill_idx = 0

    # 5-a. 거버닝 → 가장 큰 frame
    if fill_idx < len(fillable) and governing:
        try:
            _replace_text_frame_simple(fillable[fill_idx]["tf"], governing)
            replaced += 1
        except Exception as e:
            log.warning("거버닝 채움 실패: %s", e)
        fill_idx += 1

    # 5-b. 소제목 → 다음 frame
    if fill_idx < len(fillable) and subtitle:
        try:
            _replace_text_frame_simple(fillable[fill_idx]["tf"], subtitle)
            replaced += 1
        except Exception as e:
            log.warning("소제목 채움 실패: %s", e)
        fill_idx += 1

    # 5-c. 본문 → 남은 frame 들에 순서대로
    body_idx = 0
    while fill_idx < len(fillable) and body_idx < len(body):
        try:
            _replace_text_frame_simple(fillable[fill_idx]["tf"], body[body_idx])
            replaced += 1
        except Exception as e:
            log.warning("본문 [%d] 채움 실패: %s", body_idx, e)
        body_idx += 1
        fill_idx += 1

    # 5-d. summary 있으면 다음 frame 에
    if fill_idx < len(fillable) and summary:
        try:
            _replace_text_frame_simple(fillable[fill_idx]["tf"], "💡 " + summary)
            replaced += 1
        except Exception as e:
            log.warning("summary 채움 실패: %s", e)
        fill_idx += 1

    # 5-e. 본문 항목이 frame 보다 많이 남았으면 — 마지막 frame 에 합쳐서
    if body_idx < len(body) and fill_idx > 0:
        leftover = body[body_idx:]
        # 마지막으로 채운 frame 의 텍스트 뒤에 줄바꿈으로 추가
        try:
            last_tf = fillable[fill_idx - 1]["tf"]
            current = last_tf.text or ""
            combined = current + "\n" + "\n".join(leftover)
            _replace_text_frame_simple(last_tf, combined)
        except Exception as e:
            log.warning("leftover 합침 실패: %s", e)

    # fill_idx 부터 끝까지의 frame 은 비어있는 채로 둠 (마스터 디자인만 남김)
    return {
        "replaced": replaced,
        "cleared": cleared,
        "frames_total": len(frames),
        "frames_fillable": len(fillable),
        "errors": [],
    }


################################################################################
# Placeholder 모드 — 마커 기반 정확 매핑 (NightOff 정석 모드)
#
# 마스터 PPTX 안에 디자이너가 박은 마커:
#   {{거버닝}}              — 단순. max 미명시
#   {{거버닝|max:25}}       — 글자수 명시
#   {{본문_1|max:50}}       — 인덱스 + max
#   {{이미지_1|hint:콜센터}} — 이미지 자리 (코드는 비움 + hint 메타로 보관)
#   {{회사명}} {{발주처}}    — 동적 필드 (런타임에 채움)
#
# 동작 원칙:
#   - 마커 있는 자리만 치환 (마커 없는 박스는 절대 안 건드림)
#   - 마스터 디자인 100% 보존
#   - 같은 마커가 여러 곳이면 모두 같은 값으로 치환
################################################################################

import re

PLACEHOLDER_RE = re.compile(
    r"\{\{\s*(?P<key>[\w가-힣]+(?:_\d+)?)"
    r"(?:\s*\|\s*max\s*:\s*(?P<max>\d+))?"
    r"(?:\s*\|\s*hint\s*:\s*(?P<hint>[^}|]+))?"
    r"\s*\}\}"
)


def parse_placeholder(match: "re.Match") -> dict:
    """{{...}} 매치를 dict 로 변환."""
    return {
        "raw": match.group(0),
        "key": match.group("key"),
        "max": int(match.group("max")) if match.group("max") else None,
        "hint": match.group("hint").strip() if match.group("hint") else None,
    }


def find_placeholders_in_text(text: str) -> list[dict]:
    """문자열 안 모든 {{...}} 마커 추출."""
    return [parse_placeholder(m) for m in PLACEHOLDER_RE.finditer(text or "")]


def has_any_placeholder(prs) -> bool:
    """프레젠테이션 전체에 {{...}} 마커가 하나라도 있는지."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if PLACEHOLDER_RE.search(shape.text_frame.text or ""):
                return True
    return False


def collect_placeholders_in_slide(slide) -> list[dict]:
    """슬라이드의 모든 마커 수집 — 분석/검증/디버깅용.

    반환: [{"shape_name", "key", "max", "hint", "raw"}, ...]
    """
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text or ""
        for m in PLACEHOLDER_RE.finditer(text):
            ph = parse_placeholder(m)
            ph["shape_name"] = shape.name
            out.append(ph)
    return out


def _truncate(text: str, max_len: int | None) -> str:
    """글자수 제한 — 넘으면 '…' 으로 자름."""
    if max_len is None or len(text) <= max_len:
        return text
    return text[: max_len - 1].rstrip() + "…"


def _is_image_marker(key: str) -> bool:
    """이미지 placeholder 키인지 ('이미지', '이미지_1', 'image_2' 등)."""
    if not key:
        return False
    k = key.lower()
    return key.startswith("이미지") or k.startswith("image")


def _resolve_placeholder_value(ph: dict, content: dict) -> str:
    """placeholder 매치를 실제 채울 텍스트로 변환.

    우선순위:
      1. content[key] 가 있으면 그걸로 (AI 가 동적 결정)
      2. 이미지 마커면 hint 기반 안내 ("🖼 콜센터_여성")
         → 사용자가 PowerPoint 에서 박스 더블클릭 → 이미지 삽입
      3. 그 외엔 빈 문자열 (마스터 디자인만 남김)
    """
    key = ph["key"]
    if key in content and content[key] is not None:
        return _truncate(str(content[key]), ph["max"])
    if _is_image_marker(key):
        hint = ph.get("hint") or "이미지 추가"
        return f"🖼 {hint}"
    return ""


def _replace_placeholders_in_text_frame(text_frame, content: dict) -> dict:
    """text_frame 안의 모든 {{...}} 마커를 content 값으로 치환.

    동작:
      - 각 paragraph 의 각 run 의 텍스트에서 {{...}} 패턴 찾음
      - run.text 에 마커가 *완전히* 들어있으면 그 안에서 치환 (가장 흔한 케이스)
      - run 경계로 마커가 잘려있으면 paragraph.text 통째로 처리 (폴백)

    content:
      {"거버닝": "...", "본문_1": "...", "회사명": "...", ...}

    반환:
      {"replaced": int, "missing_keys": [...], "preserved_keys": [...]}
    """
    result = {"replaced": 0, "missing_keys": [], "preserved_keys": []}

    if not text_frame.paragraphs:
        return result

    # 각 paragraph 별로 처리
    for para in text_frame.paragraphs:
        # 1단계: run 단위 치환 시도 (run 안에 마커가 완전히 있는 경우)
        for run in para.runs:
            new_text = run.text or ""
            matches = list(PLACEHOLDER_RE.finditer(new_text))
            if not matches:
                continue
            # 뒤에서부터 치환 (offset 안 꼬이게)
            for m in reversed(matches):
                ph = parse_placeholder(m)
                key = ph["key"]
                val = _resolve_placeholder_value(ph, content)
                new_text = new_text[: m.start()] + val + new_text[m.end():]
                # 통계
                if key in content and content[key] is not None:
                    result["replaced"] += 1
                elif _is_image_marker(key):
                    result["replaced"] += 1  # 이미지도 치환된 것으로 카운트
                elif key not in result["missing_keys"]:
                    result["missing_keys"].append(key)
            try:
                run.text = new_text
            except Exception as e:
                log.warning("run.text 치환 실패: %s", e)

        # 2단계: paragraph 전체 텍스트에 여전히 마커 남아있으면 (run 경계 잘림)
        # paragraph.text 통째로 대체. 단 첫 run 의 스타일만 살아남음 — 트레이드오프
        full = para.text or ""
        if PLACEHOLDER_RE.search(full):
            new_full = full
            for m in reversed(list(PLACEHOLDER_RE.finditer(full))):
                ph = parse_placeholder(m)
                key = ph["key"]
                val = _resolve_placeholder_value(ph, content)
                new_full = new_full[: m.start()] + val + new_full[m.end():]
                if key in content and content[key] is not None:
                    result["replaced"] += 1
                elif _is_image_marker(key):
                    result["replaced"] += 1
                elif key not in result["missing_keys"]:
                    result["missing_keys"].append(key)
            # paragraph 의 run 들을 비우고 첫 run 에 새 텍스트
            if para.runs:
                first_run = para.runs[0]
                # 추가 run 제거
                for r in list(para.runs)[1:]:
                    try:
                        r._r.getparent().remove(r._r)
                    except Exception:
                        pass
                try:
                    first_run.text = new_full
                except Exception:
                    pass
            else:
                try:
                    para.text = new_full
                except Exception:
                    pass

    return result


def auto_inject_markers(
    master_path: "str | Path",
    output_path: "str | Path",
    *,
    dry_run: bool = False,
) -> dict:
    """[옵션 A] 빈 텍스트 박스를 가진 마스터 PPTX 에 자동으로 마커 텍스트 박기.

    크리스가 만든 페이퍼템플릿_1 같은 *빈 박스 + 디자인* 형태의 마스터를
    *자동으로* placeholder 모드 마스터로 변환.

    분류 알고리즘 (휴리스틱):
      1. 푸터 분리: y >= sh*0.85 + h <= 0.6  → 무시 (페이지 번호)
      2. 큰 가로 박스 (w >= sw*0.6 + 슬라이드 상단 50%) → 거버닝 후보
      3. 면적 기준 정렬:
         - 가장 큰 박스 → 거버닝 (이미 거버닝 후보 있으면 그걸로)
         - 두 번째 큰 박스 → 소제목 (높이가 작고 폭이 큰 경우만)
      4. 나머지 박스 → 본문_N (위→아래, 좌→우 정렬)
      5. 모든 박스 면적이 비슷하면 (예: cards layout) → 거버닝 X, 모두 본문_N

    Args:
      master_path: 입력 PPTX (빈 박스 형태)
      output_path: 출력 PPTX (마커 박힌 버전)
      dry_run: True 면 분석만 하고 저장 안 함

    Returns:
      {"slides": [{"idx", "markers_added": [{"shape_name", "marker", "x", "y", "w", "h"}]}],
       "total_markers": int}
    """
    from pptx import Presentation as _P
    p = Path(master_path)
    out = Path(output_path)
    prs = _P(str(p))
    sw_in = prs.slide_width / 914400
    sh_in = prs.slide_height / 914400

    report = {"slides": [], "total_markers": 0}

    for slide_idx, slide in enumerate(prs.slides):
        # 빈 텍스트 박스 수집 + z-order 인덱스 (앞쪽 = 뒤. 배경 식별용)
        empty_boxes = []
        all_shapes_list = list(slide.shapes)
        for z_idx, shape in enumerate(all_shapes_list):
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if (tf.text or "").strip():
                continue  # 이미 텍스트 있으면 안 건드림
            x = (shape.left or 0) / 914400
            y = (shape.top or 0) / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            empty_boxes.append({
                "shape": shape, "tf": tf, "name": shape.name, "z": z_idx,
                "x": x, "y": y, "w": w, "h": h,
                "area": w * h,
                "marker": None,
            })

        if not empty_boxes:
            report["slides"].append({"idx": slide_idx, "markers_added": []})
            continue

        # 1. 디자인 요소 식별 (마커 박지 않음)
        # 단, 빈 텍스트 박스가 적은 슬라이드 (≤3개) 는 디자인 분류 안 함
        # — 표지처럼 큰 박스 1~2개가 *진짜 텍스트 자리* 인 경우 보호
        protect_all = (len(empty_boxes) <= 3)
        for b in empty_boxes:
            if protect_all:
                # 적은 박스 슬라이드 — 분류 X, 모두 텍스트 자리로
                # 단 너무 얇은 strip 은 그래도 디자인 (헤어라인 등)
                if b["h"] < 0.15:
                    b["marker"] = "_design"
                continue
            # 일반 슬라이드 (박스 4개 이상)
            # 1-a. 배경 박스 — 슬라이드 전체 크기 + z-order 0~2 (맨 뒤)
            is_full_size = (b["w"] >= sw_in * 0.93 and b["h"] >= sh_in * 0.93)
            is_back_layer = (b["z"] <= 2)
            if is_full_size and is_back_layer:
                b["marker"] = "_background"
                continue
            # 1-b. 너무 좁고 긴 strip (예: 좌측 세로 띠 1.4x8.3)
            is_vertical_strip = (b["w"] < 1.5 and b["h"] >= sh_in * 0.7)
            is_horizontal_strip = (b["h"] < 0.25 and b["w"] >= sw_in * 0.5)
            if is_vertical_strip or is_horizontal_strip:
                b["marker"] = "_design"
                continue
            # 1-c. 너무 작음 (디자인 dot 또는 빈 셀)
            if b["area"] < 0.25:
                b["marker"] = "_tiny"
                continue

        # 2. 푸터 분리 (하단 + 작은 높이) — h 임계값 0.6→1.0 완화
        footer_y = sh_in * 0.85
        for b in empty_boxes:
            if b["marker"] is None and b["y"] >= footer_y and b["h"] <= 1.0:
                b["marker"] = "_footer"

        candidates = [b for b in empty_boxes if b["marker"] is None]
        if not candidates:
            report["slides"].append({"idx": slide_idx, "markers_added": []})
            continue

        # 1-d. 거대 박스 (배경 디자인) 추가 식별:
        # 가장 큰 박스의 면적 > 다른 박스 면적 합 → 배경 가능성 (거버닝 후보 X)
        # protect_all (박스 ≤3개 슬라이드) 은 면제
        if not protect_all and len(candidates) >= 4:
            cands_sorted_tmp = sorted(candidates, key=lambda b: -b["area"])
            biggest = cands_sorted_tmp[0]
            other_total = sum(b["area"] for b in cands_sorted_tmp[1:])
            if biggest["area"] > other_total * 1.2:
                biggest["marker"] = "_design"
                candidates = [b for b in candidates if b is not biggest]

        # 2. 거버닝 후보 — 상단 50% + 가로 60% 이상 + 면적 큼
        cands_sorted = sorted(candidates, key=lambda b: -b["area"])
        top_half_y = sh_in * 0.5

        # 면적 분포 분석 — 모든 박스가 비슷한 크기면 cards layout (거버닝 X)
        max_area = cands_sorted[0]["area"]
        similar_count = sum(1 for b in candidates if b["area"] >= max_area * 0.7)
        is_cards_layout = (similar_count >= 4 and len(candidates) >= 4)

        governing_box = None
        subtitle_box = None
        if not is_cards_layout:
            # 거버닝: 가장 큰 박스 + 상단 + 가로 큰 것
            for b in cands_sorted:
                if b["w"] >= sw_in * 0.4 and b["y"] <= top_half_y + 1.0:
                    governing_box = b
                    b["marker"] = "거버닝"
                    break
            # 그 외 가장 큰 박스 (2번째) — 거버닝 면적의 70% 이하 + 상단 가까움 → 소제목
            if governing_box:
                for b in cands_sorted:
                    if b is governing_box or b["marker"]:
                        continue
                    # 소제목 = 거버닝보다 작고 + 폭은 적당히 + 위쪽
                    if (b["area"] < governing_box["area"] * 0.7
                        and b["w"] >= sw_in * 0.3
                        and b["y"] <= top_half_y + 1.5):
                        subtitle_box = b
                        b["marker"] = "소제목"
                        break

        # 3. 본문 (위→아래, 좌→우)
        body_candidates = [b for b in candidates if b["marker"] is None]
        body_candidates.sort(key=lambda b: (round(b["y"], 1), round(b["x"], 1)))
        for i, b in enumerate(body_candidates, 1):
            b["marker"] = f"본문_{i}"

        # 4. 마커 박기 (dry_run 이면 skip)
        # underscore prefix (_background, _design, _tiny, _footer) 는 모두 skip
        slide_report = {"idx": slide_idx, "markers_added": []}
        for b in empty_boxes:
            if not b["marker"] or b["marker"].startswith("_"):
                continue
            marker_text = "{{" + b["marker"] + "}}"
            if not dry_run:
                tf = b["tf"]
                if tf.paragraphs:
                    para = tf.paragraphs[0]
                    if para.runs:
                        # 첫 run 의 텍스트만 변경
                        para.runs[0].text = marker_text
                    else:
                        # run 없는 빈 paragraph — 새 run 만들기
                        try:
                            para.text = marker_text
                        except Exception:
                            try:
                                tf.text = marker_text
                            except Exception:
                                pass
                else:
                    try:
                        tf.text = marker_text
                    except Exception:
                        pass
            slide_report["markers_added"].append({
                "shape_name": b["name"],
                "marker": b["marker"],
                "x": round(b["x"], 1), "y": round(b["y"], 1),
                "w": round(b["w"], 1), "h": round(b["h"], 1),
            })
            report["total_markers"] += 1
        report["slides"].append(slide_report)

    if not dry_run:
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))

    return report


def extract_master_slot_guide(master_path: "str | Path") -> list[dict]:
    """마스터 PPTX 를 스캔해서 슬라이드별 마커 목록 추출.

    AI 호출 시 시스템 프롬프트에 주입 → AI 가 마스터 슬롯 개수에 맞춰 콘텐츠 짬.
    (콘텐츠 N개 vs 슬롯 M개 mismatch 방지)

    반환:
      [
        {
          "idx": 0,
          "markers": [
            {"key": "거버닝", "max": 25, "hint": None},
            {"key": "회사명", "max": None, "hint": None},
          ],
          "body_count": 0,        # 본문_N 개수 (UI 표시용)
          "image_count": 0,       # 이미지_N 개수
          "section_hint": "표지",  # 섹션 추정 (마스터 노트 또는 첫 텍스트 기반)
        },
        ...
      ]
    """
    from pptx import Presentation as _P
    p = Path(master_path)
    if not p.exists():
        return []
    try:
        prs = _P(str(p))
    except Exception as e:
        log.warning("마스터 슬롯 추출 실패: %s", e)
        return []

    out: list[dict] = []
    for idx, slide in enumerate(prs.slides):
        markers: list[dict] = []
        seen_keys: set[str] = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            for m in PLACEHOLDER_RE.finditer(text):
                ph = parse_placeholder(m)
                # 같은 key 중복은 1번만 (같은 마커가 여러 자리에 있어도)
                if ph["key"] in seen_keys:
                    continue
                seen_keys.add(ph["key"])
                markers.append({
                    "key": ph["key"],
                    "max": ph["max"],
                    "hint": ph["hint"],
                })

        # 본문/이미지 개수 카운트
        body_count = sum(1 for m in markers if m["key"].startswith("본문_") or m["key"] == "본문")
        image_count = sum(1 for m in markers if m["key"].startswith("이미지_") or m["key"] == "이미지")

        # 섹션 추정 — 슬라이드 노트 (presenter notes) 우선, 없으면 ""
        section_hint = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text or ""
                # "section: ..." 형태 또는 첫 줄
                m_sec = re.search(r"section\s*:\s*(\S[^\n]+)", notes)
                if m_sec:
                    section_hint = m_sec.group(1).strip()[:30]
                elif notes.strip():
                    section_hint = notes.strip().split("\n")[0][:30]
        except Exception:
            pass

        out.append({
            "idx": idx,
            "markers": markers,
            "body_count": body_count,
            "image_count": image_count,
            "section_hint": section_hint,
        })
    return out


def format_slot_guide_for_prompt(slots: list[dict]) -> str:
    """추출된 슬롯 가이드를 AI 프롬프트에 주입할 텍스트로 변환.

    예시 출력:
      [마스터 슬라이드 슬롯 가이드]
      슬라이드 0 (표지): 거버닝, 소제목, 회사명
      슬라이드 1 (사업이해): 거버닝, 본문_1, 본문_2, 본문_3 (본문 3개)
      슬라이드 2: 거버닝, 본문_1, 본문_2, 본문_3, 본문_4, 이미지_1 (본문 4개 + 이미지)
      ...
    """
    if not slots:
        return ""
    lines = ["[마스터 슬라이드 슬롯 가이드 — 자동 추출]"]
    for s in slots:
        if not s["markers"]:
            continue  # 마커 없는 슬라이드는 가이드 안 만듦
        marker_keys = [m["key"] for m in s["markers"]]
        # 글자수 제한도 표시 (있으면)
        marker_strs = []
        for m in s["markers"]:
            if m["max"]:
                marker_strs.append(f"{m['key']}(max:{m['max']})")
            else:
                marker_strs.append(m["key"])
        section_part = f" ({s['section_hint']})" if s["section_hint"] else ""
        suffix_parts = []
        if s["body_count"]:
            suffix_parts.append(f"본문 {s['body_count']}개")
        if s["image_count"]:
            suffix_parts.append(f"이미지 {s['image_count']}개")
        suffix = f" — {' + '.join(suffix_parts)}" if suffix_parts else ""
        lines.append(f"  슬라이드 {s['idx']}{section_part}: {', '.join(marker_strs)}{suffix}")

    lines.append("")
    lines.append("⚠ 규칙: 위 슬롯 개수 정확히 맞춰서 콘텐츠 작성. 슬롯 부족하면 다른 슬라이드로 분산. 본문 개수 어기지 말 것.")
    return "\n".join(lines)


def fill_slide_with_placeholders(slide, content: dict) -> dict:
    """[Placeholder 모드 — 슬라이드 단위] 마커 있는 자리만 치환.

    마스터 디자인 100% 보존. 마커 없는 박스는 손대지 않음.

    content:
      {"거버닝": "...", "본문_1": "...", "회사명": "...", ...}

    반환:
      {"replaced": int, "missing_keys": [...], "frames_processed": int}
    """
    total = {"replaced": 0, "missing_keys": [], "frames_processed": 0}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # 마커가 하나라도 있는 frame 만 처리 (성능 + 안전)
        if not PLACEHOLDER_RE.search(shape.text_frame.text or ""):
            continue
        r = _replace_placeholders_in_text_frame(shape.text_frame, content)
        total["replaced"] += r["replaced"]
        total["frames_processed"] += 1
        for k in r["missing_keys"]:
            if k not in total["missing_keys"]:
                total["missing_keys"].append(k)
    return total


def replace_text_in_slide(slide, content: dict) -> dict:
    """[LEGACY] 기존 AUTO 모드 텍스트 치환 — 큰 글자만 치환.

    ⚠ 이 함수는 마스터의 작은 텍스트(본문, 서브헤더)를 그대로 둠 → 짬뽕 발생.
    신규 코드는 fill_slide_clearing_master() 사용.

    content 형식:
      {"거버닝": "...", "소제목": "...", "본문": ["...", "..."]}

    반환:
      {"replaced": int, "skipped": int, "errors": [...]}
    """
    zones = extract_text_zones(slide)
    if not zones["all_runs"]:
        return {"replaced": 0, "skipped": 0, "errors": ["no text runs"]}

    result = {"replaced": 0, "skipped": 0, "errors": []}
    candidates = zones["candidates"]

    # 거버닝 — 가장 큰 폰트 첫 번째 run 의 텍스트 박스 통째 교체
    if "거버닝" in content and candidates["거버닝"]:
        first = candidates["거버닝"][0]
        try:
            tf = first["run"]._r.getparent().getparent()  # txBody 의 parent (sp)
            # 해당 shape 찾기
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r is first["run"]:
                            _replace_text_frame_simple(shape.text_frame, content["거버닝"])
                            result["replaced"] += 1
                            break
        except Exception as e:
            result["errors"].append(f"거버닝: {e}")

    # 소제목 — 두 번째 크기 첫 번째 run
    if "소제목" in content and candidates["소제목"]:
        first = candidates["소제목"][0]
        try:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r is first["run"]:
                            _replace_text_frame_simple(shape.text_frame, content["소제목"])
                            result["replaced"] += 1
                            break
        except Exception as e:
            result["errors"].append(f"소제목: {e}")

    # 본문 — 본문 후보 run 들의 텍스트박스에 순서대로 교체
    if "본문" in content and candidates["본문"]:
        body_texts = content["본문"] if isinstance(content["본문"], list) else [content["본문"]]
        # 본문 텍스트박스 모음 (shape 단위, 중복 제거)
        seen_shape_ids = set()
        body_shapes = []
        for r in candidates["본문"]:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # run 이 이 shape 에 속하는지 체크
                if r["run"]._r in [run._r for p in shape.text_frame.paragraphs for run in p.runs]:
                    sid = id(shape)
                    if sid not in seen_shape_ids:
                        seen_shape_ids.add(sid)
                        body_shapes.append(shape)
                    break
        for shape, txt in zip(body_shapes, body_texts):
            try:
                _replace_text_frame_simple(shape.text_frame, txt)
                result["replaced"] += 1
            except Exception as e:
                result["errors"].append(f"본문: {e}")

    return result


# ─── 슬라이드 삭제 (인덱스 기반) ─────────────────────────────

def remove_slides_keep(prs: Presentation, keep_indices: list[int]):
    """keep_indices 에 없는 슬라이드 모두 삭제."""
    keep_set = set(keep_indices)
    xml_slides = prs.slides._sldIdLst
    slide_id_elements = list(xml_slides)
    # 뒤에서부터 삭제 (인덱스 안 꼬이게)
    for idx in range(len(slide_id_elements) - 1, -1, -1):
        if idx not in keep_set:
            try:
                xml_slides.remove(slide_id_elements[idx])
            except Exception as e:
                log.warning("슬라이드 %d 삭제 실패: %s", idx, e)


# ─── 미디어 garbage collection (PPTX 사이즈 축소) ────────────

import zipfile
import shutil as _shutil
import xml.etree.ElementTree as ET


def garbage_collect_media(pptx_path: str | Path) -> dict:
    """PPTX 안의 사용 안 하는 미디어 (이미지/동영상) 제거.

    PPTX = ZIP. 슬라이드 삭제 후엔 ppt/media/ 안 일부 파일이 더 이상
    어떤 슬라이드에도 참조 안 됨. 이걸 ZIP 에서 빼서 사이즈 축소.

    참조 추적:
      - ppt/slides/_rels/slideN.xml.rels 안의 Target → ppt/media/imageX.*
      - 살아있는 슬라이드들이 참조하는 미디어만 keep
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)

    tmp_path = pptx_path.with_suffix(".tmp.pptx")
    used_media: set[str] = set()
    live_slide_xml: set[str] = set()      # ppt/slides/slideN.xml — 살아있는 것만
    live_slide_rel: set[str] = set()      # ppt/slides/_rels/slideN.xml.rels — 살아있는 것만

    # 0. presentation.xml 의 sldIdLst → 살아있는 r:id 들 → presentation.xml.rels 의 Target 매핑
    NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    NS_PKG = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    with zipfile.ZipFile(pptx_path, "r") as zin:
        all_names = set(zin.namelist())
        live_rids: set[str] = set()
        # 1단계 — presentation.xml 의 sldIdLst 안 sldId 의 r:id
        try:
            with zin.open("ppt/presentation.xml") as f:
                pres_root = ET.fromstring(f.read())
            for sld_id in pres_root.iter(f"{NS_P}sldId"):
                rid = sld_id.attrib.get(f"{NS_R}id")
                if rid:
                    live_rids.add(rid)
        except Exception as e:
            log.warning("presentation.xml 파싱 실패: %s", e)

        # 2단계 — presentation.xml.rels 에서 r:id → Target 매핑 (살아있는 r:id 만)
        try:
            with zin.open("ppt/_rels/presentation.xml.rels") as f:
                rels_root = ET.fromstring(f.read())
            for el in rels_root.iter(f"{NS_PKG}Relationship"):
                rel_id = el.attrib.get("Id", "")
                target = el.attrib.get("Target", "")
                if rel_id in live_rids and target.startswith("slides/slide"):
                    full = "ppt/" + target  # "ppt/slides/slide1.xml"
                    live_slide_xml.add(full)
                    rel_name = "ppt/slides/_rels/" + target.split("/")[-1] + ".rels"
                    live_slide_rel.add(rel_name)
        except Exception as e:
            log.warning("presentation.xml.rels 파싱 실패: %s", e)

        # 폴백 — 살아있는 슬라이드 못 찾으면 모든 슬라이드 keep
        if not live_slide_xml:
            log.warning("살아있는 슬라이드 추적 실패 — 전체 keep")
            live_slide_xml = {n for n in all_names
                              if n.startswith("ppt/slides/slide") and n.endswith(".xml")}
            live_slide_rel = {n for n in all_names
                              if n.startswith("ppt/slides/_rels/") and n.endswith(".xml.rels")}

    log.info("살아있는 슬라이드: %d (xml) / %d (rel)", len(live_slide_xml), len(live_slide_rel))

    # 1차 스캔 — 살아있는 슬라이드의 rel 들만 참조하는 media 파일 수집
    with zipfile.ZipFile(pptx_path, "r") as zin:
        for rel_name in sorted(live_slide_rel):
            if rel_name not in all_names:
                continue
            try:
                with zin.open(rel_name) as f:
                    content = f.read()
                # XML 파싱 — Target 속성이 미디어 가리키는 것 수집
                try:
                    root = ET.fromstring(content)
                    for el in root.iter():
                        target = el.attrib.get("Target", "")
                        if "media/" in target:
                            # 상대경로 → 절대 zip path
                            # 예: "../media/image1.jpg" → "ppt/media/image1.jpg"
                            media_name = target.split("media/")[-1]
                            used_media.add(f"ppt/media/{media_name}")
                except ET.ParseError:
                    pass
            except Exception as e:
                log.warning("rel 스캔 실패 %s: %s", rel_name, e)

        # slideLayout / slideMaster 의 rels 도 추가 (테마 이미지)
        for rel_name in [n for n in all_names
                         if (n.startswith("ppt/slideLayouts/_rels/")
                             or n.startswith("ppt/slideMasters/_rels/")
                             or n.startswith("ppt/theme/_rels/"))
                         and n.endswith(".xml.rels")]:
            try:
                with zin.open(rel_name) as f:
                    content = f.read()
                try:
                    root = ET.fromstring(content)
                    for el in root.iter():
                        target = el.attrib.get("Target", "")
                        if "media/" in target:
                            media_name = target.split("media/")[-1]
                            used_media.add(f"ppt/media/{media_name}")
                except ET.ParseError:
                    pass
            except Exception:
                pass

    # 2차 — keep 할 미디어 + 살아있는 슬라이드 XML/rel 만 남기고 ZIP 다시 쓰기
    removed_media = []
    removed_slides = []
    kept_media = []
    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                name = item.filename
                # 죽은 슬라이드 XML 제거
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    if live_slide_xml and name not in live_slide_xml:
                        removed_slides.append(name)
                        continue
                # 죽은 슬라이드 .rels 제거
                if name.startswith("ppt/slides/_rels/") and name.endswith(".xml.rels"):
                    if live_slide_rel and name not in live_slide_rel:
                        removed_slides.append(name)
                        continue
                # 미디어 GC
                if name.startswith("ppt/media/"):
                    if name in used_media:
                        zout.writestr(item, zin.read(name))
                        kept_media.append(name)
                    else:
                        removed_media.append(name)
                        continue
                else:
                    zout.writestr(item, zin.read(name))

    # 원본을 GC된 사본으로 교체
    orig_size = pptx_path.stat().st_size
    _shutil.move(str(tmp_path), str(pptx_path))
    new_size = pptx_path.stat().st_size

    log.info("GC · 미디어 제거 %d, 슬라이드 제거 %d · %.1fMB → %.1fMB (%.0f%% ↓)",
             len(removed_media), len(removed_slides),
             orig_size / 1024 / 1024, new_size / 1024 / 1024,
             100 * (1 - new_size / orig_size) if orig_size else 0)

    return {
        "media_removed": len(removed_media),
        "media_kept": len(kept_media),
        "slides_removed": len(removed_slides),
        "size_before_mb": round(orig_size / 1024 / 1024, 1),
        "size_after_mb": round(new_size / 1024 / 1024, 1),
        "size_reduction_pct": round(100 * (1 - new_size / orig_size), 1) if orig_size else 0,
    }


# ─── 메인 함수 ──────────────────────────────────────────────

def generate_from_master(
    master_path: str | Path,
    content_per_slide: dict[int, dict],
    output_path: str | Path,
    keep_indices: Optional[list[int]] = None,
) -> dict:
    """마스터 PPTX 기반 새 제안서 생성.

    Args:
      master_path: 마스터 PPTX 파일 경로
      content_per_slide: {slide_idx: {"거버닝": ..., "소제목": ..., "본문": [...]}, ...}
      output_path: 결과 PPTX 저장 경로
      keep_indices: 유지할 슬라이드 인덱스 (None 이면 content_per_slide 의 키만 유지)

    Returns:
      {"slide_count": int, "replaced_total": int, "errors": [...]}
    """
    master_path = Path(master_path)
    output_path = Path(output_path)
    if not master_path.exists():
        raise FileNotFoundError(f"마스터 PPTX 없음: {master_path}")

    # 1. 마스터 통째 파일 복사 (이미지/도형/차트/표/시각화 100% 보존)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(master_path, output_path)

    # 2. 사본 열기
    prs = Presentation(str(output_path))
    slides = list(prs.slides)
    log.info("마스터 로드 · %d 슬라이드", len(slides))

    # 3. 텍스트 치환 — 마스터에 placeholder 마커 있으면 placeholder 모드, 없으면 AUTO 모드
    #    placeholder 모드: {{거버닝|max:25}} 같은 마커만 정확히 치환 (디자인 100% 보존)
    #    AUTO 모드:        모든 텍스트 비우고 사이즈 매핑으로 채움 (legacy, 짬뽕 가능)
    use_placeholder_mode = has_any_placeholder(prs)
    log.info("렌더 모드: %s", "PLACEHOLDER" if use_placeholder_mode else "AUTO")

    replaced_total = 0
    cleared_total = 0
    missing_keys_total: set[str] = set()
    errors_total = []
    for slide_idx, content in content_per_slide.items():
        if slide_idx >= len(slides):
            log.warning("슬라이드 %d 인덱스 초과 (전체 %d)", slide_idx, len(slides))
            continue
        slide = slides[slide_idx]
        try:
            if use_placeholder_mode:
                r = fill_slide_with_placeholders(slide, content)
                replaced_total += r["replaced"]
                missing_keys_total.update(r.get("missing_keys", []))
            else:
                r = fill_slide_clearing_master(slide, content)
                replaced_total += r["replaced"]
                cleared_total += r.get("cleared", 0)
                errors_total.extend([f"slide{slide_idx}: {e}" for e in r.get("errors", [])])
        except Exception as e:
            log.exception("슬라이드 %d 치환 실패: %s", slide_idx, e)
            errors_total.append(f"slide{slide_idx}: {e}")

    # 3-b. content 가 지정되지 않은 슬라이드 (keep_indices 에는 있지만 content_per_slide 에 없는)
    #      AUTO 모드: 모든 텍스트 비움 (마스터 원본 잔재 방지)
    #      placeholder 모드: 손대지 않음 (마커 없으면 자연히 안 건드림)
    if keep_indices is not None and not use_placeholder_mode:
        for idx in keep_indices:
            if idx in content_per_slide:
                continue
            if idx >= len(slides):
                continue
            slide = slides[idx]
            try:
                r = fill_slide_clearing_master(slide, {})
                cleared_total += r.get("cleared", 0)
            except Exception as e:
                log.warning("슬라이드 %d 빈값 처리 실패: %s", idx, e)

    # 4. keep_indices 지정 시 그 외 삭제
    if keep_indices is None:
        keep_indices = sorted(content_per_slide.keys())
    if keep_indices:
        remove_slides_keep(prs, keep_indices)

    # 5. 저장
    prs.save(str(output_path))

    final_count = len(prs.slides)
    if use_placeholder_mode:
        log.info("저장 · %d 슬라이드 / 마커 치환 %d · 누락키 %d · 에러 %d",
                 final_count, replaced_total, len(missing_keys_total), len(errors_total))
        if missing_keys_total:
            log.info("누락된 마커 키 (마스터엔 있는데 content 에 없음): %s",
                     sorted(missing_keys_total)[:20])
    else:
        log.info("저장 · %d 슬라이드 / 치환 %d · 비움 %d · 에러 %d",
                 final_count, replaced_total, cleared_total, len(errors_total))

    # 6. 미디어 garbage collection — 사용 안 하는 이미지/동영상 제거 (사이즈 축소)
    gc_result = {}
    try:
        gc_result = garbage_collect_media(output_path)
    except Exception as e:
        log.warning("미디어 GC 실패 (무시): %s", e)

    return {
        "slide_count": final_count,
        "replaced_total": replaced_total,
        "errors": errors_total[:10],
        "output_path": str(output_path),
        "media_gc": gc_result,
    }


# ─── PPTX → PNG 미리보기 변환 ────────────────────────────────

LIBREOFFICE_CANDIDATES = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/libreoffice",
    "/usr/bin/soffice",
    "soffice",
]


def _find_soffice() -> Optional[str]:
    """LibreOffice 실행파일 자동 탐색."""
    import subprocess as sp
    for c in LIBREOFFICE_CANDIDATES:
        p = Path(c)
        if p.exists():
            return str(p)
        try:
            r = sp.run([c, "--version"], capture_output=True, timeout=5)
            if r.returncode == 0:
                return c
        except Exception:
            continue
    return None


def pptx_to_png_previews(
    pptx_path: str | Path,
    out_dir: str | Path,
    *,
    width: int = 1280,
    timeout_sec: int = 90,
) -> list[Path]:
    """PPTX → PDF (LibreOffice) → 페이지별 PNG (pypdfium2).

    Args:
      pptx_path: 입력 PPTX
      out_dir: PNG 저장 디렉토리 — slide_01.png, slide_02.png, ...
      width: PNG 가로 픽셀 (높이는 종횡비 유지)
      timeout_sec: LibreOffice 변환 타임아웃

    Returns:
      생성된 PNG 경로 리스트 (slide_idx 순)
    """
    import subprocess as sp
    import tempfile

    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    soffice = _find_soffice()
    if not soffice:
        log.warning("LibreOffice 못 찾음 — PNG 미리보기 생성 불가")
        return []

    # 1. PPTX → PDF (LibreOffice headless)
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        # 임시 user-profile (다른 soffice 인스턴스와 충돌 방지)
        profile_dir = tmp_dir / "profile"
        profile_dir.mkdir()
        cmd = [
            soffice, "--headless", "--norestore", "--nologo", "--nofirststartwizard",
            f"-env:UserInstallation=file:///{str(profile_dir).replace(chr(92), '/')}",
            "--convert-to", "pdf",
            "--outdir", str(tmp_dir),
            str(pptx_path),
        ]
        try:
            r = sp.run(cmd, capture_output=True, timeout=timeout_sec)
            if r.returncode != 0:
                log.warning("LibreOffice PDF 변환 실패: %s",
                            r.stderr.decode("utf-8", errors="replace")[:200])
                return []
        except sp.TimeoutExpired:
            log.warning("LibreOffice PDF 변환 타임아웃 (%ds)", timeout_sec)
            try:
                sp.run(["taskkill", "/F", "/IM", "soffice.bin"],
                       capture_output=True, timeout=10)
            except Exception:
                pass
            return []

        pdf_files = list(tmp_dir.glob("*.pdf"))
        if not pdf_files:
            log.warning("PDF 출력 못 찾음")
            return []
        pdf_path = pdf_files[0]

        # 2. PDF → PNG 페이지별 (pypdfium2)
        try:
            import pypdfium2 as pdfium
        except ImportError:
            log.warning("pypdfium2 미설치 — pip install pypdfium2")
            return []

        png_paths: list[Path] = []
        try:
            pdf = pdfium.PdfDocument(str(pdf_path))
            n_pages = len(pdf)
            for i in range(n_pages):
                page = pdf[i]
                # 1280px 너비 기준 scale 계산 (PDF 1pt = 1/72인치)
                pdf_w_pt = page.get_width()
                scale = width / pdf_w_pt if pdf_w_pt else 2.0
                bitmap = page.render(scale=scale)
                pil = bitmap.to_pil()
                out_path = out_dir / f"slide_{i+1:02d}.png"
                # JPEG 가 PNG 보다 작지만 미리보기는 PNG 가 무손실
                pil.save(out_path, "PNG", optimize=True)
                png_paths.append(out_path)
            pdf.close()
        except Exception as e:
            log.exception("PDF → PNG 변환 실패: %s", e)
            return []

    log.info("PNG 미리보기 생성 · %d 슬라이드 → %s", len(png_paths), out_dir)
    return png_paths


def find_master_template(domain: Optional[str] = None) -> Optional[Path]:
    """분야에 맞는 마스터 PPTX 파일 찾기.

    조회 순서 (placeholder 모드 우선):
      1. master_templates/paperlogy_default.pptx (auto_inject_markers 적용된 placeholder 마스터)
      2. master_templates/dmz_default.pptx (legacy, AUTO 모드 fallback)
      3. master_templates/ 안의 첫 *.pptx
      4. R2_LOCAL_CACHE_DIR 환경변수가 가리키는 디렉토리

    차후 domain 별 매핑 확장 예정.
    """
    import os
    candidates: list[Path] = []
    base = Path(__file__).parent / "master_templates"
    if base.is_dir():
        # 1. placeholder 마스터 우선 (paperlogy_default 또는 *_placeholder)
        candidates.append(base / "paperlogy_default.pptx")
        # 그 외 *_placeholder.pptx 패턴
        candidates.extend(sorted(base.glob("*_placeholder.pptx")))
        # 2. legacy 마스터 (AUTO 모드)
        candidates.append(base / "dmz_default.pptx")
        # 3. 그 외 모든 pptx
        candidates.extend(sorted(base.glob("*.pptx")))
    cache_env = os.environ.get("R2_LOCAL_CACHE_DIR")
    if cache_env:
        cache = Path(cache_env)
        if cache.is_dir():
            candidates.append(cache / "paperlogy_default.pptx")
            candidates.append(cache / "dmz_default.pptx")
            candidates.extend(sorted(cache.glob("*.pptx")))
    seen: set[str] = set()
    for c in candidates:
        if str(c) in seen:
            continue
        seen.add(str(c))
        if c.exists() and c.stat().st_size > 0:
            return c
    return None


################################################################################
# 🎨 도형 JSON 모드 — Claude 가 layout 자유 결정 → 원시 도형 그리기
#
# 마스터 PPTX 와 무관. AI 가 슬라이드별로 도형 + 위치 + 텍스트 자유롭게 정함.
# 우리 코드는 *원시 도형 그리기 함수* 만 제공 (rect/text/line/circle/arrow/image)
# 입력 형식: 도형 JSON 스펙
################################################################################

from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn


def _hex_to_rgb(hex_color):
    """#RRGGBB → RGBColor (잘못된 입력은 검정으로 폴백)."""
    if not hex_color:
        return RGBColor(0, 0, 0)
    h = str(hex_color).lstrip("#").strip()
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return RGBColor(0, 0, 0)


def _set_no_fill(shape) -> None:
    """투명 채움."""
    try:
        shape.fill.background()
    except Exception:
        pass


def _set_no_line(shape) -> None:
    """테두리 없음."""
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _add_rect(slide, x, y, w, h, *, fill="#FFFFFF", stroke=None, stroke_width=None, radius=None,
              theme="light"):
    """사각형 (옵션: rounded, 테두리, 채움).

    Spec D-Build-ThemeColorMap (1-c, 옵션 A): theme='dark' 면 fill/stroke 를 role 별 다크 매핑.
    light 면 입력색 그대로 (_map_color 가 theme 가드).
    """
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(_map_color(fill, "fill", theme))
    else:
        _set_no_fill(shape)
    if stroke:
        shape.line.color.rgb = _hex_to_rgb(_map_color(stroke, "stroke", theme))
        if stroke_width:
            shape.line.width = Pt(float(stroke_width))
    else:
        _set_no_line(shape)
    return shape


def _add_text(slide, x, y, w, h, text, *,
              size=14, weight=400, color="#1A1A1A",
              align="left", valign="top",
              font_family=None, italic=False,
              theme="light",
              text_runs=None):
    """텍스트 박스. 줄바꿈 \\n 으로 멀티라인 지원.

    Spec D-Build-ThemeColorMap (1-c, 옵션 A): theme='dark' 면 color 를 text role 다크 매핑.
    light 면 입력색 그대로.

    Spec D-Build-TextRunsRender (1-d-①): text_runs 가 있으면 부분 강조 경로.
      형식: [{"t":"...", "accent":bool?}, ...]
      run 단위로 색 분기 — accent=True && theme=='dark' 만 accent 색(#A78BFA), 그 외 일반 color.
      ★ text_runs 없음/빈 list → 기존 경로 가드 (비트 단위 무변경 / 한 글자도 안 건드림).
      ★ text 필드는 fallback 으로 유지 — text_runs 미지원/디버깅용. 새 경로는 text 무시.
      ★ segment 의 "t" 안에 \\n 이 있으면: split 후 각 \\n 경계마다 paragraph 분리,
         accent 속성은 분리된 양쪽 모두 유지(빈 segment 는 paragraph 경계 역할만).
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE_FIT  # D-Fix-AutoSize: 텍스트가 박스 넘치면 폰트 자동 축소 (넘침 방지)
    except Exception:
        pass
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)     # Phase 4 — 본문 시각 여백 ↑ (옵션 1)
    tf.margin_bottom = Inches(0.04)  # Phase 4 — 본문 시각 여백 ↑ (옵션 1)
    valign_map = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    try:
        tf.vertical_anchor = valign_map.get(str(valign).lower(), MSO_ANCHOR.TOP)
    except Exception:
        pass

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    # weight 기반 자동 폰트 매핑 (font_family 명시 시 그것이 우선)
    target_font = _resolve_font(font_family, weight)
    # 안전망: target_font 가 빈 문자열 / None 으로 새어 나올 경우 DEFAULT 강제
    if not target_font or not str(target_font).strip():
        target_font = DEFAULT_FONT_FAMILY
    weight_norm = _normalize_weight(weight)

    # Spec D-Build-TextRunsRender (1-d-①) — 신규 경로 분기.
    # text_runs 가 None/빈값/list 아님 → 기존 경로(아래)로 떨어짐 → 비트 단위 무변경.
    if text_runs and isinstance(text_runs, list):
        # 1) 평탄화 — 각 segment 의 t 를 \n 으로 split, accent 속성 보존.
        #    빈 part 는 run 생성 X (paragraph 경계 역할만), 정상 part 만 run 으로 추가.
        paragraphs: list[list[dict]] = [[]]
        for seg in text_runs:
            if not isinstance(seg, dict):
                continue
            t = str(seg.get("t", ""))
            accent = bool(seg.get("accent", False))
            parts = t.split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    paragraphs.append([])
                if part:
                    paragraphs[-1].append({"t": part, "accent": accent})

        # 2) 모든 paragraph 가 비면 빈 박스만 반환 (안전망).
        if not any(paragraphs):
            return box

        # 3) paragraph 단위로 그리기 — 같은 paragraph 내 run 색만 분기, 나머지 폰트는 동일.
        for i, segs in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.15  # 기존 경로와 동일 — Phase 4 일관
            try:
                p.alignment = align_map.get(str(align).lower(), PP_ALIGN.LEFT)
            except Exception:
                pass
            # 빈 paragraph(연속 \n) → run 0개 = 빈 줄
            for seg in segs:
                run = p.add_run()
                run.text = seg["t"]
                try:
                    run.font.size = Pt(float(size))
                except Exception:
                    run.font.size = Pt(14)
                run.font.bold = weight_norm >= 600
                if italic:
                    run.font.italic = True
                # ★ 색 결정 — accent && dark 만 accent 색(#A78BFA). 라이트는 accent 무시 → 일반 color.
                if seg["accent"] and theme == "dark":
                    seg_color = "#A78BFA"
                else:
                    seg_color = color
                run.font.color.rgb = _hex_to_rgb(_map_color(seg_color, "text", theme))
                try:
                    run.font.name = target_font
                except Exception:
                    pass
        return box

    # 기존 경로 (text_runs 없음) — 한 글자도 안 바뀜 / 비트 단위 무변경 가드.
    lines = (text or "").split("\n")

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15  # Phase 4 — 본문 시각 여백 ↑ (옵션 1, 모든 paragraph)
        try:
            p.alignment = align_map.get(str(align).lower(), PP_ALIGN.LEFT)
        except Exception:
            pass
        run = p.add_run()
        run.text = line
        try:
            run.font.size = Pt(float(size))
        except Exception:
            run.font.size = Pt(14)
        # PowerPoint 호환 bool bold 유지 (weight 600 이상이면 True)
        # — 매핑된 폰트가 이미 굵기 표현하지만 PowerPoint 기본 렌더링 호환용
        run.font.bold = weight_norm >= 600
        if italic:
            run.font.italic = True
        run.font.color.rgb = _hex_to_rgb(_map_color(color, "text", theme))
        # 폰트 매핑: font_family 명시 우선 → weight 자동 매핑 (Paperlogy 9 단계)
        try:
            run.font.name = target_font
        except Exception:
            pass
    return box


def _add_line(slide, x1, y1, x2, y2, *, color="#1A1A1A", width=1.0, theme="light"):
    """직선 (커넥터).

    Spec D-Build-ThemeColorMap (1-c): theme='dark' 면 color 를 stroke role 다크 매핑.
    """
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = _hex_to_rgb(_map_color(color, "stroke", theme))
    try:
        line.line.width = Pt(float(width))
    except Exception:
        line.line.width = Pt(1)
    return line


def _add_arrow(slide, x1, y1, x2, y2, *, color="#1A1A1A", width=1.5, theme="light"):
    """화살표 — 직선 + tail 끝에 삼각형.

    Spec D-Build-ThemeColorMap (1-c): theme='dark' 면 color 를 stroke role 다크 매핑.
    """
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = _hex_to_rgb(_map_color(color, "stroke", theme))
    try:
        line.line.width = Pt(float(width))
    except Exception:
        line.line.width = Pt(1.5)
    # XML 직접 조작 — tail 에 화살촉 추가
    try:
        ln = line.line._get_or_add_ln()
        # 기존 헤드/테일 제거
        for tag in ("a:headEnd", "a:tailEnd"):
            existing = ln.find(qn(tag))
            if existing is not None:
                ln.remove(existing)
        from lxml import etree
        head_end = etree.SubElement(ln, qn("a:headEnd"))
        head_end.set("type", "none")
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
        tail_end.set("type", "triangle")
        tail_end.set("w", "med")
        tail_end.set("len", "med")
    except Exception as e:
        log.warning("화살촉 추가 실패 (선만 표시): %s", e)
    return line


def _add_circle(slide, x, y, w, h, *, fill="#000000", stroke=None, stroke_width=None,
                theme="light"):
    """원/타원.

    Spec D-Build-ThemeColorMap (1-c): theme='dark' 면 fill/stroke 를 role 별 다크 매핑.

    Spec D-Build-PresetCircles — fill 정규화: "none"/"transparent"/"" → 채움 없음(투명).
      render_shape_to_slide 의 circle 분기가 str() 로 강제 변환하므로 None 으로는 안 도달.
      preset 'circles' 가 fill="none" 으로 테두리만 있는 원을 요청할 때 필수 가드.
      이 가드 없으면 _hex_to_rgb("none") 가 6자리 hex 아님 → 검정 폴백 → 원이 검정으로 꽉 참.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill_norm = fill
    if isinstance(fill_norm, str) and fill_norm.strip().lower() in ("none", "transparent", ""):
        fill_norm = None
    if fill_norm:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(_map_color(fill_norm, "fill", theme))
    else:
        _set_no_fill(shape)
    if stroke:
        shape.line.color.rgb = _hex_to_rgb(_map_color(stroke, "stroke", theme))
        if stroke_width:
            shape.line.width = Pt(float(stroke_width))
    else:
        _set_no_line(shape)
    return shape


def _add_image_placeholder(slide, x, y, w, h, hint="이미지 추가", *, theme="light"):
    """이미지 자리 — 회색 박스 + 안내. 사용자가 PowerPoint 에서 더블클릭으로 이미지 삽입.

    Spec D-Build-ThemeColorMap (1-c, 옵션 A — 예외 처리):
      이 함수는 RGBColor 직접 호출(하드코딩)이라 _map_color 로 감쌀 수 없음.
      theme='dark' 면 본문 내 직접 분기 — 다크 placeholder 색 사용.
      light(현재): fill=#ECECEC, stroke=#CCCCCC, text=#888888
      dark       : fill=#1F1F1F, stroke=#444444, text=#888888 (text 는 동일 — 안내문구 가독 유지)
    """
    if theme == "dark":
        _fill_rgb = RGBColor(0x1F, 0x1F, 0x1F)   # 어두운 회색 면 (다크 배경보다 살짝 들뜸)
        _stroke_rgb = RGBColor(0x44, 0x44, 0x44) # 어두운 테두리
        _text_rgb = RGBColor(0x88, 0x88, 0x88)   # 동일한 중회색 안내문구
    else:
        _fill_rgb = RGBColor(0xEC, 0xEC, 0xEC)   # 현재값 그대로 (라이트 무변경)
        _stroke_rgb = RGBColor(0xCC, 0xCC, 0xCC)
        _text_rgb = RGBColor(0x88, 0x88, 0x88)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _fill_rgb
    box.line.color.rgb = _stroke_rgb
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "🖼  " + str(hint)
    run.font.size = Pt(11)
    run.font.color.rgb = _text_rgb
    run.font.italic = True
    return box


# ============================================================
# Spec D-Fix-11a Stage A (5/19) — 신규 도형 헬퍼 5종
# 시각화 다양성 영역 확장: chevron / pentagon / callout / block_arrow / star
# ⚠ Stage A 만 적용 시 효과 0 — LLM 영역 신규 type 출력 X (시스템 프롬프트 무변경).
# ⚠ Stage B 영역 (시스템 프롬프트 안내) 진입 전 안전 영역 확보 단계.
# 패턴 정합: _add_rect (fill/stroke) + _add_image_placeholder (text_frame 영역).
# ============================================================

def _add_shape_with_text(slide, mso_shape_type, x, y, w, h, *,
                          fill="#FFFFFF", stroke=None, stroke_width=None,
                          text=None, text_color="#1A1A1A", text_size=12,
                          text_weight=400, text_align="center",
                          theme="light"):
    """내부 헬퍼 — 신규 도형 5종 영역 공통 본문 (DRY 영역).

    fill / stroke 영역 = _add_rect 패턴 정확 정합.
    텍스트 영역 = _add_image_placeholder 패턴 + _add_text 정렬 정합.
    Spec D-Build-ThemeColorMap (1-c): theme='dark' 면 fill/stroke/text_color 모두 role 별 다크 매핑.
    """
    shape = slide.shapes.add_shape(
        mso_shape_type, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    # fill / stroke (_add_rect 패턴)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(_map_color(fill, "fill", theme))
    else:
        _set_no_fill(shape)
    if stroke:
        shape.line.color.rgb = _hex_to_rgb(_map_color(stroke, "stroke", theme))
        if stroke_width:
            shape.line.width = Pt(float(stroke_width))
    else:
        _set_no_line(shape)
    # 텍스트 영역 (옵션, _add_image_placeholder 패턴 + _add_text 정렬)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        try:
            p.alignment = align_map.get(str(text_align).lower(), PP_ALIGN.CENTER)
        except Exception:
            pass
        target_font = _resolve_font(None, text_weight) or DEFAULT_FONT_FAMILY
        weight_norm = _normalize_weight(text_weight)
        run = p.add_run()
        run.text = str(text)
        try:
            run.font.size = Pt(float(text_size))
        except Exception:
            run.font.size = Pt(12)
        run.font.bold = weight_norm >= 600
        run.font.color.rgb = _hex_to_rgb(_map_color(text_color, "text", theme))
        try:
            run.font.name = target_font
        except Exception:
            pass
    return shape


def _add_chevron(slide, x, y, w, h, *, fill="#FFFFFF", stroke=None, stroke_width=None,
                 text=None, text_color="#1A1A1A", text_size=12, text_weight=400, text_align="center",
                 theme="light"):
    """CHEVRON 도형 (프로세스 단계 영역). 텍스트는 도형 안 가운데 정렬.

    Spec D-Build-ThemeColorMap (1-c): theme 를 _add_shape_with_text 로 통과.
    """
    return _add_shape_with_text(
        slide, MSO_SHAPE.CHEVRON, x, y, w, h,
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        text=text, text_color=text_color, text_size=text_size,
        text_weight=text_weight, text_align=text_align, theme=theme,
    )


def _add_pentagon(slide, x, y, w, h, *, fill="#FFFFFF", stroke=None, stroke_width=None,
                  text=None, text_color="#1A1A1A", text_size=12, text_weight=400, text_align="center",
                  theme="light"):
    """PENTAGON 도형 (Phase / 단계 영역). 텍스트는 도형 안 가운데 정렬.

    Spec D-Build-ThemeColorMap (1-c): theme 를 _add_shape_with_text 로 통과.
    """
    return _add_shape_with_text(
        slide, MSO_SHAPE.PENTAGON, x, y, w, h,
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        text=text, text_color=text_color, text_size=text_size,
        text_weight=text_weight, text_align=text_align, theme=theme,
    )


def _add_callout(slide, x, y, w, h, *, fill="#FFFFFF", stroke="#1A1A1A", stroke_width=1.0,
                 text=None, text_color="#1A1A1A", text_size=12, text_weight=400, text_align="center",
                 theme="light"):
    """RECTANGULAR_CALLOUT 도형 (말풍선 / 핵심 메시지 강조).

    기본 stroke = #1A1A1A — 말풍선 영역 본질 (테두리 있음).
    Spec D-Build-ThemeColorMap (1-c): theme 를 _add_shape_with_text 로 통과.
    """
    return _add_shape_with_text(
        slide, MSO_SHAPE.RECTANGULAR_CALLOUT, x, y, w, h,
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        text=text, text_color=text_color, text_size=text_size,
        text_weight=text_weight, text_align=text_align, theme=theme,
    )


def _add_block_arrow(slide, x, y, w, h, *, fill="#1A1A1A", stroke=None, stroke_width=None,
                     text=None, text_color="#FFFFFF", text_size=12, text_weight=400, text_align="center",
                     theme="light"):
    """RIGHT_ARROW 도형 (굵은 화살표 — line/arrow 영역 보다 시각 ↑).

    기본 fill = #1A1A1A / text_color = #FFFFFF — 검은 화살표 + 흰 글자 본질.
    Spec D-Build-ThemeColorMap (1-c): theme 를 _add_shape_with_text 로 통과.
    """
    return _add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h,
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        text=text, text_color=text_color, text_size=text_size,
        text_weight=text_weight, text_align=text_align, theme=theme,
    )


def _add_star(slide, x, y, w, h, *, fill="#1A1A1A", stroke=None, stroke_width=None,
              text=None, text_color="#FFFFFF", text_size=14, text_weight=700, text_align="center",
              theme="light"):
    """STAR_5_POINT 도형 (차별화 / 강조 포인트).

    기본 text_weight=700 / text_size=14 — 별 영역 강조 본질.
    Spec D-Build-ThemeColorMap (1-c): theme 를 _add_shape_with_text 로 통과.
    """
    return _add_shape_with_text(
        slide, MSO_SHAPE.STAR_5_POINT, x, y, w, h,
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        text=text, text_color=text_color, text_size=text_size,
        text_weight=text_weight, text_align=text_align, theme=theme,
    )


def render_shape_to_slide(slide, shape_def, *, default_text_color="#1A1A1A", theme="light"):
    """단일 도형 스펙(JSON) → 슬라이드에 그림.

    지원 type:
      [기존] rect, text, line, arrow, circle/ellipse/oval, image/image_placeholder
      [Stage A (5/19) 신규] chevron, pentagon, callout, block_arrow, star
    실패 시 None 반환 (다른 도형 렌더링은 계속됨).

    [Spec D-Build-ThemeConnect 1-b]
    default_text_color — text 도형의 color 가 누락됐을 때의 기본 글자색.
      라이트(기본): "#1A1A1A" — 현재 운영과 동일.
      다크: "#FFFFFF" — generate_from_shape_json 에서 theme="dark" 일 때 주입.
    ★ shape_def 에 color 가 명시돼 있으면 그 값을 그대로 사용 — 기본값만 영향.

    [Spec D-Build-ThemeColorMap 1-c (옵션 A)]
    theme — 'light'(기본) / 'dark'. 모든 helper 호출에 그대로 전달 →
      helper 내부에서 _map_color 가 role 별로 다크 매핑(라이트면 입력색 그대로).
    """
    if not isinstance(shape_def, dict):
        return None
    t = str(shape_def.get("type", "")).lower().strip()
    try:
        if t in ("rect", "rectangle"):
            return _add_rect(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 1)), float(shape_def.get("h", 1)),
                fill=shape_def.get("fill"),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                radius=shape_def.get("radius"),
                theme=theme,
            )
        if t == "text":
            # Spec D-Fix-GovColorRemove-1 — 거버닝 파랑(#1E40AF) 강제 휴리스틱 제거.
            # AI 가 SLIDE_SYSTEM_PROMPT 의 흑백 6색 규칙으로 출력하니 그 색 그대로 사용.
            # Spec D-Build-ThemeConnect 1-b — color 누락 시 기본값을 theme 토큰으로(라이트=#1A1A1A 유지).
            # Spec D-Build-ThemeColorMap 1-c — 명시 color 도 다크에서 role 별 매핑(_add_text 내부 _map_color).
            # Spec D-Build-TextRunsRender 1-d-① — text_runs(선택) 그대로 전달, _add_text 가 분기 처리.
            #   text_runs 없는 도형(기존 100%) → _add_text 내부에서 기존 경로로 들어가 비트 단위 무변경.
            # Spec Governing-Purple — role=="governing" 시 색을 브랜드 accent 로 강제.
            #   각 preset builder 가 거버닝(메인 제목) 도형에만 "role":"governing" 마킹.
            #   theme.py 의 ACCENT 토큰 참조 (라이트 #6B46E5 / 다크 #A78BFA).
            #   미마킹 도형(divider 로마·quantitative value·아이템 head 등)은 무영향.
            color_val = str(shape_def.get("color", default_text_color))
            if shape_def.get("role") == "governing":
                _accent = _get_theme(theme).get("ACCENT")
                if _accent:
                    color_val = _accent
            return _add_text(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 5)), float(shape_def.get("h", 1)),
                str(shape_def.get("text", "")),
                size=float(shape_def.get("size", 14)),
                weight=int(shape_def.get("weight", 400)),
                color=color_val,
                align=str(shape_def.get("align", "left")),
                valign=str(shape_def.get("valign", "top")),
                font_family=shape_def.get("font_family"),
                italic=bool(shape_def.get("italic", False)),
                theme=theme,
                text_runs=shape_def.get("text_runs"),
            )
        if t == "line":
            return _add_line(
                slide,
                float(shape_def.get("x1", 0)), float(shape_def.get("y1", 0)),
                float(shape_def.get("x2", 1)), float(shape_def.get("y2", 0)),
                color=str(shape_def.get("color", "#1A1A1A")),
                width=float(shape_def.get("width", 1.0)),
                theme=theme,
            )
        if t == "arrow":
            return _add_arrow(
                slide,
                float(shape_def.get("x1", 0)), float(shape_def.get("y1", 0)),
                float(shape_def.get("x2", 1)), float(shape_def.get("y2", 0)),
                color=str(shape_def.get("color", "#1A1A1A")),
                width=float(shape_def.get("width", 1.5)),
                theme=theme,
            )
        if t in ("circle", "ellipse", "oval"):
            return _add_circle(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 1)), float(shape_def.get("h", 1)),
                fill=str(shape_def.get("fill", "#000000")),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                theme=theme,
            )
        if t in ("image", "image_placeholder"):
            return _add_image_placeholder(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 4)), float(shape_def.get("h", 3)),
                hint=str(shape_def.get("hint", "이미지 추가")),
                theme=theme,
            )
        # Spec D-Fix-11a Stage A (5/19) — 신규 도형 5종 분기
        # ⚠ 효과 0 (의도된 안전 — LLM 영역 본 type 출력 X, Stage B 영역 안내 추가 후 활성).
        if t in ("chevron", "process_step"):
            return _add_chevron(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 1)), float(shape_def.get("h", 0.5)),
                fill=shape_def.get("fill", "#FFFFFF"),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                text=shape_def.get("text"),
                text_color=str(shape_def.get("text_color", "#1A1A1A")),
                text_size=float(shape_def.get("text_size", 12)),
                text_weight=int(shape_def.get("text_weight", 400)),
                text_align=str(shape_def.get("text_align", "center")),
                theme=theme,
            )
        if t in ("pentagon", "step"):
            return _add_pentagon(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 1)), float(shape_def.get("h", 0.8)),
                fill=shape_def.get("fill", "#FFFFFF"),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                text=shape_def.get("text"),
                text_color=str(shape_def.get("text_color", "#1A1A1A")),
                text_size=float(shape_def.get("text_size", 12)),
                text_weight=int(shape_def.get("text_weight", 400)),
                text_align=str(shape_def.get("text_align", "center")),
                theme=theme,
            )
        if t in ("callout", "rect_callout"):
            return _add_callout(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 3)), float(shape_def.get("h", 1)),
                fill=shape_def.get("fill", "#FFFFFF"),
                stroke=shape_def.get("stroke", "#1A1A1A"),
                stroke_width=float(shape_def.get("stroke_width", 1.0)),
                text=shape_def.get("text"),
                text_color=str(shape_def.get("text_color", "#1A1A1A")),
                text_size=float(shape_def.get("text_size", 12)),
                text_weight=int(shape_def.get("text_weight", 400)),
                text_align=str(shape_def.get("text_align", "center")),
                theme=theme,
            )
        if t in ("block_arrow", "right_arrow"):
            return _add_block_arrow(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 2)), float(shape_def.get("h", 0.5)),
                fill=shape_def.get("fill", "#1A1A1A"),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                text=shape_def.get("text"),
                text_color=str(shape_def.get("text_color", "#FFFFFF")),
                text_size=float(shape_def.get("text_size", 12)),
                text_weight=int(shape_def.get("text_weight", 400)),
                text_align=str(shape_def.get("text_align", "center")),
                theme=theme,
            )
        if t in ("star", "star_5_point"):
            return _add_star(
                slide,
                float(shape_def.get("x", 0)), float(shape_def.get("y", 0)),
                float(shape_def.get("w", 1)), float(shape_def.get("h", 1)),
                fill=shape_def.get("fill", "#1A1A1A"),
                stroke=shape_def.get("stroke"),
                stroke_width=shape_def.get("stroke_width"),
                text=shape_def.get("text"),
                text_color=str(shape_def.get("text_color", "#FFFFFF")),
                text_size=float(shape_def.get("text_size", 14)),
                text_weight=int(shape_def.get("text_weight", 700)),
                text_align=str(shape_def.get("text_align", "center")),
                theme=theme,
            )
    except Exception as e:
        log.warning("도형 렌더링 실패 (type=%s): %s", t, e)
    return None


def _build_preset_quantitative_emphasis(slide_data: dict) -> list:
    """Spec D-Fix-Preset1 — 정량 강조(큰 숫자) 레이아웃 프리셋.

    slide_data["metrics"] = [{"value": "50억", "label": "총 예산"}, ...] (1~3개).
    반환 = render_shape_to_slide 가 처리하는 JSON 도형 리스트 (type/x/y/w/h/text/...).
    형식 오류 / 데이터 없음 시 빈 리스트 (호출부 try/except 와 별개 안전망).

    좌표 설계 (A4 가로 11.69 × 8.27 인치):
      · 상단 분리선 y=2.5 (강조)
      · 큰 숫자 1~3개: y=2.8~5.0, size 90~120pt, 중앙·균등 배치
      · 라벨: 숫자 바로 아래 y=5.1, 18pt, 회색
    """
    metrics = slide_data.get("metrics") or []
    if not isinstance(metrics, list) or not metrics:
        return []
    valid: list = []
    for m in metrics:
        if not isinstance(m, dict):
            continue
        v = str(m.get("value", "")).strip()
        lab = str(m.get("label", "")).strip()
        if v:
            valid.append((v, lab))
        if len(valid) >= 3:
            break
    if not valid:
        return []

    n = len(valid)
    if n == 1:
        positions = [(4.0, 3.7, 120)]
    elif n == 2:
        positions = [(1.3, 4.4, 100), (6.0, 4.4, 100)]
    else:  # 3
        positions = [(0.9, 3.3, 90), (4.2, 3.3, 90), (7.5, 3.3, 90)]

    # ★ Spec Governing-Title-Three-Presets — 조건부 거버닝 (하위 호환: title 없으면 미출력).
    #   상단 empty 2.5" 여유 (line y=2.5) → y=1.0 h=1.0 자리. 기존 도형과 겹침 없음.
    shapes: list = []
    title = str(slide_data.get("title", "")).strip()
    if title:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 1.0, "w": 9.89, "h": 1.0,
            "text": title,
            "size": 28, "weight": 800, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
            "role": "governing",
        })
    shapes.append({
        "type": "line",
        "x1": 0.9, "y1": 2.5, "x2": 10.8, "y2": 2.5,
        "color": "#1A1A1A", "width": 1.5,
    })
    for (x, w, size), (val, lab) in zip(positions, valid):
        shapes.append({
            "type": "text",
            "x": x, "y": 2.8, "w": w, "h": 2.2,
            "text": val,
            "size": size, "weight": 900, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
        })
        if lab:
            shapes.append({
                "type": "text",
                "x": x, "y": 5.1, "w": w, "h": 0.6,
                "text": lab,
                "size": 18, "weight": 500, "color": "#666",
                "align": "center", "valign": "top",
            })
    return shapes


def _build_preset_horizontal_process(slide_data: dict) -> list:
    """Spec D-Fix-Preset2 — 가로 프로세스(추진 절차/운영 흐름) 레이아웃 프리셋.

    slide_data["steps"] = [{"label": "분석", "desc": "현황 진단" (선택)}, ...] (3~7개).
    반환 = render_shape_to_slide 가 처리하는 JSON 도형 리스트 (chevron + 선택 desc text).
    형식 오류 / 데이터 없음 / 단계 수 3~7 밖 → 빈 리스트 (호출부 try/except 와 별개 안전망).

    좌표 설계 (A4 가로 11.69 × 8.27 / 좌우 여백 0.9 / 본문 폭 9.89):
      · chevron 가로 균등 배치 (단계 수에 따라 box_w 자동 계산)
      · gap = 0.1 (n<=5) / 0.05 (n=6~7)
      · box_y = 3.5 / box_h = 1.2
      · 통일 흰 fill + 검정 stroke 1.5pt + label 14pt 700weight
      · desc 있으면 chevron 아래 11pt 회색 (y=4.9, h=1.3)

    ⚠ chevron 은 auto_size 미적용 (텍스트 박스 전용) — label 짧게 강제 권장
       (3단계 5~10자 / 5단계 4~6자 / 7단계 2~4자).
    """
    steps = slide_data.get("steps") or []
    if not isinstance(steps, list):
        return []
    valid: list = []
    for s in steps:
        if not isinstance(s, dict):
            continue
        label = str(s.get("label", "")).strip()
        desc = str(s.get("desc", "")).strip()
        if label:
            valid.append((label, desc))
    n = len(valid)
    if n < 3 or n > 7:
        return []

    margin = 0.9
    gap = 0.1 if n <= 5 else 0.05
    inner_w = 11.69 - 2 * margin
    box_w = (inner_w - gap * (n - 1)) / n
    box_h = 1.2
    box_y = 3.5

    shapes: list = []
    # ★ Spec Governing-Title-Three-Presets — 조건부 거버닝 (하위 호환: title 없으면 미출력).
    #   상단 empty 3.5" 여유 (chevron y=3.5) → y=1.0 h=1.0 자리. 기존 도형과 겹침 없음.
    title = str(slide_data.get("title", "")).strip()
    if title:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 1.0, "w": 9.89, "h": 1.0,
            "text": title,
            "size": 28, "weight": 800, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
            "role": "governing",
        })
    for i, (label, desc) in enumerate(valid):
        x = margin + i * (box_w + gap)
        shapes.append({
            "type": "chevron",
            "x": x, "y": box_y, "w": box_w, "h": box_h,
            "fill": "#FFFFFF",
            "stroke": "#1A1A1A",
            "stroke_width": 1.5,
            "text": label,
            "text_color": "#1A1A1A",
            "text_size": 14,
            "text_weight": 700,
            "text_align": "center",
        })
        if desc:
            shapes.append({
                "type": "text",
                "x": x, "y": 4.9, "w": box_w, "h": 1.3,
                "text": desc,
                "size": 11, "weight": 400, "color": "#666",
                "align": "center", "valign": "top",
            })
    return shapes


def _build_preset_two_column(slide_data: dict) -> list:
    """Spec D-Fix-Preset3 — 2분할(AS-IS/TO-BE / 현재→개선 / 문제→해결) 레이아웃 프리셋.

    slide_data["columns"] = [
      {"title": "AS-IS", "items": ["현재 상태 1", "현재 상태 2", ...]},
      {"title": "TO-BE", "items": ["개선 결과 1", "개선 결과 2", ...]}
    ] (정확히 2개 / 좌=AS-IS / 우=TO-BE).
    반환 = render_shape_to_slide 가 처리하는 JSON 도형 리스트
           (좌우 rect + 제목 text + 항목 text + 중앙 block_arrow).
    형식 오류 / columns 2개 아님 / 빈 데이터 → 빈 리스트 (호출부 try/except 와 별개 안전망).

    좌표 설계 (A4 가로 11.69 × 8.27 / 좌우 여백 0.9 / 본문 폭 9.89):
      · 중앙 화살표 자리 1.0 / panel_w = (9.89 - 1.0) / 2 = 4.445
      · 좌 패널 x=0.9 / 우 패널 x=6.345 / panel_y=2.8 / panel_h=4.0
      · 좌 = 흰 fill + 회색 stroke (#999) / 우 = 흰 fill + 검정 stroke (#1A1A1A 강조)
      · 제목 18pt 700w / 항목 13pt 400w (auto_size 적용 — 넘침 방지)
      · 중앙 block_arrow (x=5.35 y=4.4 w=1.0 h=0.8 / 검정 fill + 흰 글자 "→")
    """
    columns = slide_data.get("columns") or []
    if not isinstance(columns, list) or len(columns) != 2:
        return []
    parsed: list = []
    for c in columns:
        if not isinstance(c, dict):
            return []
        title = str(c.get("title", "")).strip()
        items_raw = c.get("items") or []
        if not isinstance(items_raw, list):
            items_raw = []
        items = [str(i).strip() for i in items_raw if str(i).strip()]
        parsed.append((title, items))
    if not any(t or its for t, its in parsed):
        return []

    margin = 0.9
    center_gap = 1.0
    panel_w = (11.69 - 2 * margin - center_gap) / 2
    panel_h = 4.0
    panel_y = 2.8
    left_x = margin
    right_x = margin + panel_w + center_gap

    shapes: list = []
    # ★ Spec Governing-Title-Three-Presets — 조건부 거버닝 (하위 호환: title 없으면 미출력).
    #   상단 empty 2.8" 여유 (panel y=2.8) → y=1.0 h=1.0 자리. 기존 도형과 겹침 없음.
    title = str(slide_data.get("title", "")).strip()
    if title:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 1.0, "w": 9.89, "h": 1.0,
            "text": title,
            "size": 28, "weight": 800, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
            "role": "governing",
        })
    panel_styles = [
        {"x": left_x,  "stroke": "#999999", "stroke_width": 1.0},
        {"x": right_x, "stroke": "#1A1A1A", "stroke_width": 1.5},
    ]
    for (title, items), style in zip(parsed, panel_styles):
        x = style["x"]
        # 패널 사각형
        shapes.append({
            "type": "rect",
            "x": x, "y": panel_y, "w": panel_w, "h": panel_h,
            "fill": "#FFFFFF",
            "stroke": style["stroke"],
            "stroke_width": style["stroke_width"],
        })
        # 제목
        if title:
            shapes.append({
                "type": "text",
                "x": x + 0.2, "y": panel_y + 0.15, "w": panel_w - 0.4, "h": 0.5,
                "text": title,
                "size": 18, "weight": 700, "color": "#1A1A1A",
                "align": "left", "valign": "top",
            })
        # 항목 리스트 (한 텍스트 박스에 줄바꿈 — auto_size 자동 적용)
        if items:
            body = "\n".join("· " + it for it in items)
            shapes.append({
                "type": "text",
                "x": x + 0.2, "y": panel_y + 0.85, "w": panel_w - 0.4, "h": panel_h - 1.05,
                "text": body,
                "size": 13, "weight": 400, "color": "#1A1A1A",
                "align": "left", "valign": "top",
            })

    # 중앙 block_arrow (좌 → 우 전환)
    shapes.append({
        "type": "block_arrow",
        "x": left_x + panel_w + 0.05, "y": panel_y + (panel_h - 0.8) / 2,
        "w": center_gap - 0.1, "h": 0.8,
        "fill": "#1A1A1A",
        "text": "→",
        "text_color": "#FFFFFF",
        "text_size": 20,
        "text_weight": 700,
    })
    return shapes


def _build_preset_narrative(slide_data: dict) -> list:
    """Spec D-Fix-Preset4 / D-Fix-NarrativeV3 — 내러티브 흐름형 레이아웃 프리셋.

    도식 없이 텍스트 위계로 흐름을 보여주는 단순 슬롯 구조.
    slide_data["quote"]      = "큰 인용 문구"               (필수)
    slide_data["eyebrow"]    = "좌상단 메타 라벨"           (선택 — D-Fix-NarrativeV3)
    slide_data["flow"]       = ["흐름 설명1", "흐름 설명2", ...] (선택 / 1~3개)
    slide_data["conclusion"] = "결론 강조 문구"             (선택)
    반환 = render_shape_to_slide 가 처리하는 JSON 도형 리스트.
    quote 없음 → 빈 리스트 (호출부 try/except 와 별개 안전망).

    좌표 설계 v3 (Spec D-Fix-NarrativeV3 — 거버닝 없는 전제 + 수직 중앙정렬):
      · eyebrow (선택)  : x=0.9 y=0.5 w=9.89 h=0.4 / 11pt 400w #BBBBBB left top
      · 인용 (필수)    : x=1.2 y=2.0 w=9.29 h=1.8 / 40pt 700w 검정 center middle
      · 흐름 (1~3개)  : x=1.2 w=9.29 h=0.55 / 15pt 400w #666 center / y=4.1 +0.62 간격
      · 결론 (선택)   : 가는 구분선(rect h=0.02) + text 18pt 700w center middle

    Spec D-Fix-Preset5 — style 분기:
      · "quote" (기본) → 본 로직 (D-Fix-NarrativeV3)
      · "declaration" → 큰 선언 + 근거 2~3개 (D-Fix-NarrativeV3 v3 좌표)
      · "qa"          → 질문 + 답변 1~3 (기존 유지)
      · "emphasis"    → 소제목 + 본문 + 핵심 강조 (기존 유지)
      · "contrast"    → "A가 아니라 B" 대비 (기존 유지)
    """
    style = str(slide_data.get("style", "quote")).strip().lower()
    if style == "declaration":
        return _narrative_declaration(slide_data)
    if style == "qa":
        return _narrative_qa(slide_data)
    if style == "emphasis":
        return _narrative_emphasis(slide_data)
    if style == "contrast":
        return _narrative_contrast(slide_data)
    # style == "quote" 또는 미지정 — Spec D-Fix-NarrativeV3:
    #   거버닝 없는 전제(컨셉 슬로건은 별도 hero 페이지가 담당) + 수직 중앙정렬.
    #   eyebrow 옵션 — 있으면 좌상단 회색 작은 메타, 없으면 생략.
    quote = str(slide_data.get("quote", "")).strip()
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    flow_raw = slide_data.get("flow") or []
    if not isinstance(flow_raw, list):
        flow_raw = []
    flow = [str(f).strip() for f in flow_raw if str(f).strip()][:3]
    conclusion = str(slide_data.get("conclusion", "")).strip()

    if not quote:
        return []

    shapes: list = []
    # eyebrow (선택) — 좌상단 회색 메타 1줄
    if eyebrow:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 0.5, "w": 9.89, "h": 0.4,
            "text": eyebrow,
            "size": 11, "weight": 400, "color": "#BBBBBB",
            "align": "left", "valign": "top",
        })
    # 큰 인용 (필수) — 40pt 수직 중앙
    shapes.append({
        "type": "text",
        "x": 1.2, "y": 2.0, "w": 9.29, "h": 1.8,
        "text": quote,
        "size": 40, "weight": 700, "color": "#1A1A1A",
        "align": "center", "valign": "middle",
        "role": "governing",
    })
    # 흐름 설명 (선택 / 최대 3개) — 0.62 간격
    y = 4.1
    for i, line in enumerate(flow):
        shapes.append({
            "type": "text",
            "x": 1.2, "y": y + i * 0.62, "w": 9.29, "h": 0.55,
            "text": line,
            "size": 15, "weight": 400, "color": "#666666",
            "align": "center", "valign": "middle",
        })
    # 결론 (선택) — 가는 구분선 + 18pt 텍스트
    yend = y + len(flow) * 0.62 + 0.25
    if conclusion:
        shapes.append({
            "type": "rect",
            "x": 4.34, "y": yend, "w": 3.0, "h": 0.02,
            "fill": "#1A1A1A",
        })
        shapes.append({
            "type": "text",
            "x": 1.2, "y": yend + 0.15, "w": 9.29, "h": 0.6,
            "text": conclusion,
            "size": 18, "weight": 700, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
        })
    return shapes


def _narrative_declaration(slide_data: dict) -> list:
    """D-Fix-Preset5 / D-Fix-NarrativeV3 narrative style=declaration —
    거버닝 없는 전제 + 수직 중앙정렬 + eyebrow 옵션.
    큰 선언(44pt) + 근거 2~3개("— " 접두).
    """
    declaration = str(slide_data.get("declaration", "")).strip()
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    grounds_raw = slide_data.get("grounds") or []
    if not isinstance(grounds_raw, list):
        grounds_raw = []
    grounds = [str(g).strip() for g in grounds_raw if str(g).strip()][:3]
    if not declaration:
        return []
    shapes: list = []
    # eyebrow (선택)
    if eyebrow:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 0.5, "w": 9.89, "h": 0.4,
            "text": eyebrow,
            "size": 11, "weight": 400, "color": "#BBBBBB",
            "align": "left", "valign": "top",
        })
    # 큰 선언 (필수) — 44pt
    shapes.append({
        "type": "text",
        "x": 1.2, "y": 2.2, "w": 9.29, "h": 2.0,
        "text": declaration,
        "size": 44, "weight": 700, "color": "#1A1A1A",
        "align": "center", "valign": "middle",
        "role": "governing",
    })
    # 근거 (선택 / 최대 3개) — Spec Preset-Fix-Declaration A안:
    #   옅은 테두리 박스 + 왼쪽 정렬 + em-dash 접두 제거 (박스가 구분 역할 대체).
    #   box_h 0.75 (16pt 2줄 + 여유 0.24") / gap 0.9 (박스 사이 0.15" 여백) / y 4.6 시작.
    #   3 박스 끝 = 4.6 + 2*0.9 + 0.75 = 7.15 → 하단(8.27) 여유 1.12".
    y = 4.6
    box_h = 0.75
    gap = 0.9
    for i, g in enumerate(grounds):
        gy = y + i * gap
        # 옅은 테두리 박스 — 근거 구분
        shapes.append({
            "type": "rect", "x": 1.2, "y": gy, "w": 9.29, "h": box_h,
            "fill": "none", "stroke": "#DDDDDD", "stroke_width": 1,
        })
        # 근거 텍스트 — 왼쪽 정렬 + 박스 안 패딩(좌 0.25")
        shapes.append({
            "type": "text", "x": 1.45, "y": gy, "w": 8.79, "h": box_h,
            "text": g,
            "size": 16, "weight": 400, "color": "#555555",
            "align": "left", "valign": "middle",
        })
    return shapes


def _narrative_qa(slide_data: dict) -> list:
    """D-Fix-Preset5 narrative style=qa — 질문 + 답변 1~3."""
    question = str(slide_data.get("question", "")).strip()
    answers_raw = slide_data.get("answers") or []
    if not isinstance(answers_raw, list):
        answers_raw = []
    answers = [str(a).strip() for a in answers_raw if str(a).strip()][:3]
    if not question:
        return []
    shapes: list = [{
        "type": "text",
        "x": 0.9, "y": 2.0, "w": 9.89, "h": 1.4,
        "text": "Q. " + question,
        "size": 28, "weight": 700, "color": "#1A1A1A",
        "align": "center", "valign": "middle",
        "role": "governing",
    }]
    for i, a in enumerate(answers):
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 4.0 + i * 0.9, "w": 9.89, "h": 0.8,
            "text": a,
            "size": 16, "weight": 400, "color": "#444444",
            "align": "center", "valign": "middle",
        })
    return shapes


def _narrative_emphasis(slide_data: dict) -> list:
    """D-Fix-Preset5 narrative style=emphasis — 소제목 + 본문 + 핵심 강조."""
    subtitle = str(slide_data.get("subtitle", "")).strip()
    body = str(slide_data.get("body", "")).strip()
    highlight = str(slide_data.get("highlight", "")).strip()
    if not (subtitle or body or highlight):
        return []
    shapes: list = []
    if subtitle:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 1.6, "w": 9.89, "h": 0.8,
            "text": subtitle,
            "size": 18, "weight": 700, "color": "#666666",
            "align": "center", "valign": "middle",
        })
    if body:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 2.6, "w": 9.89, "h": 2.5,
            "text": body,
            "size": 15, "weight": 400, "color": "#333333",
            "align": "center", "valign": "top",
        })
    if highlight:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 5.5, "w": 9.89, "h": 1.0,
            "text": highlight,
            "size": 24, "weight": 700, "color": "#1A1A1A",
            "align": "center", "valign": "middle",
            "role": "governing",
        })
    return shapes


def _narrative_contrast(slide_data: dict) -> list:
    """D-Fix-Preset5 narrative style=contrast — 'A가 아니라 B' 대비."""
    not_this = str(slide_data.get("not_this", "")).strip()
    but_this = str(slide_data.get("but_this", "")).strip()
    if not but_this:
        return []
    shapes: list = []
    if not_this:
        shapes.append({
            "type": "text",
            "x": 0.9, "y": 2.8, "w": 9.89, "h": 1.0,
            "text": not_this,
            "size": 22, "weight": 400, "color": "#999999",
            "align": "center", "valign": "middle",
        })
    shapes.append({
        "type": "text",
        "x": 0.9, "y": 4.2, "w": 9.89, "h": 1.6,
        "text": but_this,
        "size": 34, "weight": 700, "color": "#1A1A1A",
        "align": "center", "valign": "middle",
        "role": "governing",
    })
    return shapes


# ─── Spec D-Build-PresetSplit — split (색면 2분할) 레이아웃 프리셋 ─────────────
# 좌 검정/우 흰 색면 대비 + 좌우 각각 (label / head / points 1~4) 텍스트 위계.
# 입력 스키마:
#   slide_data["left"]  = {"label": "...", "head": "...(필수)", "points": [..., ...]}
#   slide_data["right"] = {"label": "...", "head": "...(필수)", "points": [..., ...]}
# 안전망: left/right 가 dict 아님 / head 누락 → 빈 리스트 반환 (호출부 try/except 와
#         별개 안전 fallback). 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec.
def _build_preset_split(slide_data):
    left  = slide_data.get("left")  or {}
    right = slide_data.get("right") or {}
    if not isinstance(left, dict) or not isinstance(right, dict):
        return []
    l_label = str(left.get("label", "")).strip()
    l_head  = str(left.get("head", "")).strip()
    l_pts = [str(p).strip() for p in (left.get("points") or []) if str(p).strip()][:4]
    r_label = str(right.get("label", "")).strip()
    r_head  = str(right.get("head", "")).strip()
    r_pts = [str(p).strip() for p in (right.get("points") or []) if str(p).strip()][:4]
    if not l_head or not r_head:
        return []
    W, H = 11.69, 8.27
    half = W / 2
    shapes = [
        {"type":"rect","x":0,"y":0,"w":half,"h":H,"fill":"#1A1A1A"},
        {"type":"rect","x":half,"y":0,"w":half,"h":H,"fill":"#FFFFFF"},
    ]
    def _start_y(label, pts):
        block = 0.0
        if label: block += 0.5 + 0.2
        block += 1.2
        if pts: block += 0.5 + len(pts) * 0.62
        return max(0.8, (H - block) / 2)
    def _build_side(x, w, label, head, pts, label_color, head_color, pt_color, head_weight):
        out = []
        y = _start_y(label, pts)
        if label:
            out.append({"type":"text","x":x,"y":y,"w":w,"h":0.5,"text":label,"size":13,"weight":700,"color":label_color,"align":"center","valign":"middle"})
            y += 0.7
        out.append({"type":"text","x":x,"y":y,"w":w,"h":1.2,"text":head,"size":26,"weight":head_weight,"color":head_color,"align":"center","valign":"middle","role":"governing"})
        y += 1.7
        for p in pts:
            out.append({"type":"text","x":x,"y":y,"w":w,"h":0.55,"text":p,"size":13,"weight":400,"color":pt_color,"align":"center","valign":"middle"})
            y += 0.62
        return out
    lx, lw = 0.6, half - 1.2
    rx, rw = half + 0.6, half - 1.2
    shapes += _build_side(lx, lw, l_label, l_head, l_pts, "#FFFFFF", "#FFFFFF", "#DDDDDD", 800)
    shapes += _build_side(rx, rw, r_label, r_head, r_pts, "#1A1A1A", "#1A1A1A", "#444444", 800)
    return shapes


# ─── Spec Preset-Remove-Matrix — _build_governing_block 헬퍼 제거 (옵션 B 완전삭제) ──
# 본 위치에 있던 _build_governing_block 헬퍼는 matrix 유일 사용자가 함께 삭제되며
# dead code 화되어 제거. 향후 donut/venn 등 신규 preset 이 거버닝 패턴을 재활용하려면
# git history (이전 commit) 에서 함수 복원하거나 timeline (자체 거버닝 처리) 패턴 참고.


# ─── Spec D-Build-PresetTimeline — timeline (세로 점 단계형) 레이아웃 프리셋 ───
# 좌측 세로 라인 + 점(circle) 단계 마커 + 각 단계 (label / head / desc) 텍스트 위계,
# 우측은 이미지 영역 placeholder (회색 box). 단계 1~6.
# 입력 스키마:
#   slide_data["eyebrow"] = "좌상단 메타 라벨"        (선택)
#   slide_data["title"]   = "상단 페이지 제목"        (선택)
#   slide_data["steps"]   = [{"label","head"(필수),"desc"}, ...]  (1~6)
# 안전망: steps list 아님 / 모든 step.head 누락 → 빈 리스트 반환.
# 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec.
def _build_preset_timeline(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    title   = str(slide_data.get("title", "")).strip()
    steps_raw = slide_data.get("steps") or []
    if not isinstance(steps_raw, list):
        return []
    steps = []
    for s in steps_raw:
        if not isinstance(s, dict):
            continue
        head = str(s.get("head", "")).strip()
        if not head:
            continue
        steps.append({
            "label": str(s.get("label", "")).strip(),
            "head": head,
            "desc": str(s.get("desc", "")).strip(),
        })
    steps = steps[:6]
    if not steps:
        return []
    W, H = 11.69, 8.27
    shapes = []
    if eyebrow:
        # Spec D-Fix-TimelineGoverningPosition — eyebrow x 0.9 → 0.5 로 일반 본문 페이지 정합.
        #   일반 본문 거버닝 표준 좌표 (proposal_multi_pass.py:1737-1738 SLIDE 예시) x=0.5 와
        #   상단 좌측 정렬 통일. w=9.89 / y=0.5 / h=0.4 / size 11 / color #BBBBBB 무변경.
        #   ★ 본 변경은 timeline preset 만 — 다른 preset (asymmetric/circles/zigzag/hsplit/quad 등)의
        #     동일한 eyebrow 도형은 무변경 (사용자 명령).
        shapes.append({"type":"text","x":0.5,"y":0.5,"w":9.89,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    top = 1.1
    if title:
        # Spec D-Fix-TimelineGoverningPartialAccent — 직전 spec(D-Fix-TimelineGoverningFollowup)의
        # "전체 강조(text_runs 단일 run accent=True)" 를 "핵심 명사구 한 곳만 형광" 으로 교체.
        #   다른 거버닝 표준: D-Build-TextRunsInject (proposal_multi_pass.py:3057-3083) —
        #   "거버닝당 형광 1곳만 / 핵심 가치를 압축한 명사구 한 곳에만 accent".
        #   timeline 도 같은 방식: LLM 이 SLIDE 출력에 slide_data["title_runs"] 키로 segment 직접 결정.
        #     [{"t":"앞부분 "},{"t":"핵심 명사구","accent":true},{"t":" 뒷부분"}]
        #   title_runs 유효(list + 비어있지 않음) → title 도형의 text_runs 키로 그대로 박음
        #     → dispatch(L2040) → _add_text(L1638) → run 단위 색 분기 → accent && dark 만 #A78BFA.
        #   title_runs 없음/빈 list/비-list → text_runs 키 자체 누락 → _add_text L1690 기존 경로 →
        #     일반 거버닝 fallback (형광 없음, 거버닝 자체는 28pt 로 표시 — 안전).
        #   ★ size 28 / 좌표 (y=0.8, h=1.2, top=2.4, area_top=2.7) 무변경 — 단계 분포 0 영향.
        # Spec D-Fix-TimelineGoverningPosition — title 거버닝 x 0.9 → 0.5, w 10.0 → 10.5 로
        #   일반 본문 페이지 거버닝 표준 좌표 (proposal_multi_pass.py:1737-1738 SLIDE 예시
        #   x=0.5, y=0.8, w=10.5, h=1.2) 와 정렬 통일. eyebrow x=0.5 와도 좌측 정렬 일관.
        #   ★ y=0.8 / h=1.2 / size=28 / 형광 title_runs / weight=800 / color 무변경 — 위치 (x/w) 만.
        #   ★ 거버닝 끝 = 0.8 + 1.2 = 2.0 그대로 → top=2.4, area_top=2.7, area_bot=7.4, ys 분포,
        #     박스 마커, 세로 라인, 우측 placeholder 모두 무변경 (좌표 0 영향).
        #   ★ 단계 박스 line_x=1.6 무변경 — 거버닝/eyebrow 만 좌측 정렬 (단계는 그대로).
        title_runs_raw = slide_data.get("title_runs")
        title_shape = {
            "type": "text", "x": 0.5, "y": 0.8, "w": 10.5, "h": 1.2,
            "text": title,
            "size": 28, "weight": 800, "color": "#1A1A1A",
            "align": "left", "valign": "top",
            "role": "governing",
        }
        if isinstance(title_runs_raw, list) and title_runs_raw:
            title_shape["text_runs"] = title_runs_raw
        shapes.append(title_shape)
        top = 2.4
    line_x = 1.2
    dot_r = 0.04
    # Spec D-Fix-TimelineDotsLine — 마커를 원(circle) → 작은 네모(rect, 한 변 0.1") 로 교체.
    #   box_half 는 박스 한 변의 절반 (= 0.05"). 박스 그릴 때 좌상단 = (line_x - box_half, cy - box_half).
    #   dot_r 자체는 무변경 — placeholder 영역(ph_top/ph_bot, L2717-2718) 계산이 dot_r 에 의존하므로
    #   ph_x/ph_w/ph_h 좌표 정합을 위해 dot_r=0.04 그대로 유지.
    box_half = 0.05
    area_top = top + 0.3
    # Spec D-Fix-TimelineGoverning — 슬라이드 바닥(8.27") 침범 안전망.
    #   진단 발견: 옛 area_bot=7.5 + label+head+desc 모두 있는 마지막 단계 →
    #   ys[-1] + 0.82 (desc 끝) = 8.32 → 바닥 0.05" 침범 (기존 버그).
    #   변경: 7.5 → 7.4 로 0.1" 올림. n=4 (cap 1.5 적용) 영향 0, n=5 desc 끝 = 8.22 (안전).
    #   n=3·4 의 ys[-1] 는 cap 1.5 로 결정되므로 area_bot 0.1" 변경 영향 거의 없음 (안전 여유 ↑).
    area_bot = 7.4
    n = len(steps)
    if n == 1:
        ys = [(area_top + area_bot) / 2]
    else:
        gap = (area_bot - area_top) / (n - 1)
        gap = min(gap, 1.2)
        total = gap * (n - 1)
        start = area_top + ((area_bot - area_top) - total) / 2
        ys = [start + i * gap for i in range(n)]
    text_x = line_x + 0.6
    text_w = 6.7 - text_x - 0.3
    # Spec D-Fix-TimelineDotsLine — 세로 라인 추가 (첫 박스 중심 ~ 마지막 박스 중심).
    #   박스보다 먼저 append → z-order 아래 → 박스가 라인 위에 도드라짐 (PPTX 도형 추가 순 = 하단).
    #   n==1 일 때는 라인 안 그림 (시작=끝, h=0 → 의미 없음).
    if n >= 2:
        shapes.append({"type":"rect","x":line_x - 0.005,"y":ys[0],"w":0.01,"h":ys[-1] - ys[0],"fill":"#DDDDDD"})
    for i, st in enumerate(steps):
        cy = ys[i]
        # Spec D-Fix-TimelineDotsLine — 마커: 원(circle) → 작은 네모(rect, 한 변 0.1").
        #   기존 점 (지름 0.08" ≈ 2mm) 이 너무 작아 시각 신호 약함 → 박스(한 변 0.1") + 라인 조합.
        #   x/y 좌상단 = (line_x - box_half, cy - box_half) → 중심이 기존 점 중심과 동일.
        shapes.append({"type":"rect","x":line_x - box_half,"y":cy - box_half,"w":box_half*2,"h":box_half*2,"fill":"#1A1A1A"})
        ty = cy - 0.45
        if st["label"]:
            shapes.append({"type":"text","x":text_x,"y":ty,"w":text_w,"h":0.3,"text":st["label"],"size":11,"weight":700,"color":"#999999","align":"left","valign":"top"})
            ty += 0.32
        shapes.append({"type":"text","x":text_x,"y":ty,"w":text_w,"h":0.45,"text":st["head"],"size":17,"weight":700,"color":"#1A1A1A","align":"left","valign":"top"})
        ty += 0.45
        if st["desc"]:
            shapes.append({"type":"text","x":text_x,"y":ty,"w":text_w,"h":0.5,"text":st["desc"],"size":12,"weight":400,"color":"#666666","align":"left","valign":"top"})
    ph_x = 6.7
    ph_w = W - ph_x - 0.8
    ph_top = ys[0] - dot_r
    ph_bot = ys[-1] + dot_r
    ph_h = ph_bot - ph_top
    shapes.append({"type":"rect","x":ph_x,"y":ph_top,"w":ph_w,"h":ph_h,"fill":"#F5F5F5","stroke":"#DDDDDD","stroke_width":1})
    shapes.append({"type":"text","x":ph_x,"y":ph_top + ph_h/2 - 0.2,"w":ph_w,"h":0.4,"text":"이미지 영역","size":12,"weight":400,"color":"#BBBBBB","align":"center","valign":"middle"})
    return shapes


# ─── Spec D-Build-PresetAsymmetric — asymmetric (비대칭 2분할) 레이아웃 프리셋 ─
# 좌측 흰 영역(약 65%) + 우측 검정 영역(약 35%) — 좌측은 큰 숫자/메인 헤드라인/포인트
# 텍스트 위계, 우측 검정면은 흰 글씨 보조 항목 리스트(label/desc).
# 입력 스키마:
#   slide_data["eyebrow"] = "좌상단 메타 라벨"      (선택)
#   slide_data["number"]  = "큰 숫자/번호"          (선택, 예 "01")
#   slide_data["head"]    = "큰 헤드라인"           (필수)
#   slide_data["points"]  = ["좌측 본문 포인트", ...] (선택, 최대 3)
#   slide_data["items"]   = [{"label","desc"}, ...]   (선택, 최대 4 — 우측 검정면)
# 안전망: head 누락 → 빈 리스트 반환. 1단계는 코드만 등록 — viz_pattern 연결은 별도 spec.
def _build_preset_asymmetric(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    number  = str(slide_data.get("number", "")).strip()
    head    = str(slide_data.get("head", "")).strip()
    pts = [str(p).strip() for p in (slide_data.get("points") or []) if str(p).strip()][:3]
    items_raw = slide_data.get("items") or []
    items = []
    if isinstance(items_raw, list):
        for it in items_raw:
            if not isinstance(it, dict):
                continue
            lab = str(it.get("label", "")).strip()
            dsc = str(it.get("desc", "")).strip()
            if lab or dsc:
                items.append({"label": lab, "desc": dsc})
    items = items[:4]
    if not head:
        return []
    W, H = 11.69, 8.27
    divider_x = 7.6
    shapes = []
    shapes.append({"type":"rect","x":divider_x,"y":0,"w":W-divider_x,"h":H,"fill":"#1A1A1A"})
    if eyebrow:
        shapes.append({"type":"text","x":0.9,"y":0.5,"w":6.0,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    lx = 0.9
    lw = divider_x - lx - 0.5
    y = 1.6
    if number:
        shapes.append({"type":"text","x":lx,"y":y,"w":lw,"h":1.0,"text":number,"size":54,"weight":800,"color":"#1A1A1A","align":"left","valign":"top"})
        y += 1.3
    shapes.append({"type":"text","x":lx,"y":y,"w":lw,"h":1.4,"text":head,"size":28,"weight":800,"color":"#1A1A1A","align":"left","valign":"top","role":"governing"})
    y += 1.5
    for p in pts:
        shapes.append({"type":"text","x":lx,"y":y,"w":lw,"h":0.7,"text":p,"size":14,"weight":400,"color":"#444444","align":"left","valign":"top"})
        y += 0.8
    rx = divider_x + 0.5
    rw = W - rx - 0.5
    if items:
        ry_top = 1.7
        ry_bot = H - 1.0
        n = len(items)
        gap = (ry_bot - ry_top) / n
        for i, it in enumerate(items):
            iy = ry_top + i * gap
            if it["label"]:
                shapes.append({"type":"text","x":rx,"y":iy,"w":rw,"h":0.4,"text":it["label"],"size":14,"weight":700,"color":"#FFFFFF","align":"left","valign":"top"})
                iy += 0.42
            if it["desc"]:
                shapes.append({"type":"text","x":rx,"y":iy,"w":rw,"h":0.8,"text":it["desc"],"size":12,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    return shapes


# ─── Spec D-Build-PresetZigzag — zigzag (가운데 축 좌우 번갈아) 레이아웃 프리셋 ─
# 슬라이드 중앙 세로 라인 + 점 단계 마커. 각 단계 텍스트(label/head/desc)를 좌·우 번갈아
# 배치 — 짝수 번째 좌측(right-align), 홀수 번째 우측(left-align). 아이템 1~6개.
# 입력 스키마:
#   slide_data["eyebrow"] = "좌상단 메타 라벨"      (선택)
#   slide_data["title"]   = "상단 페이지 제목"      (선택)
#   slide_data["items"]   = [{"label","head"(필수),"desc"}, ...]  (1~6)
# 안전망: items list 아님 또는 모든 head 누락 → 빈 리스트 반환.
# 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec.
def _build_preset_zigzag(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    title   = str(slide_data.get("title", "")).strip()
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        items.append({
            "label": str(it.get("label", "")).strip(),
            "head": head,
            "desc": str(it.get("desc", "")).strip(),
        })
    items = items[:6]
    if not items:
        return []
    W, H = 11.69, 8.27
    cx = W / 2
    shapes = []
    if eyebrow:
        shapes.append({"type":"text","x":0.9,"y":0.5,"w":9.89,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    top = 1.1
    if title:
        shapes.append({"type":"text","x":0.9,"y":1.0,"w":10.0,"h":0.9,"text":title,"size":26,"weight":800,"color":"#1A1A1A","align":"left","valign":"top","role":"governing"})
        top = 2.3
    n = len(items)
    area_top = top + 0.3
    area_bot = 7.5
    gap = (area_bot - area_top) / n
    dot_r = 0.05
    shapes.append({"type":"line","x1":cx,"y1":area_top,"x2":cx,"y2":area_bot,"color":"#DDDDDD","width":1.5})
    for i, it in enumerate(items):
        cy = area_top + gap * i + gap / 2
        is_left = (i % 2 == 0)
        shapes.append({"type":"circle","x":cx - dot_r,"y":cy - dot_r,"w":dot_r*2,"h":dot_r*2,"fill":"#1A1A1A"})
        if is_left:
            tx = 0.9
            tw = cx - tx - 0.7
            al = "right"
        else:
            tx = cx + 0.7
            tw = W - tx - 0.9
            al = "left"
        ty = cy - 0.5
        if it["label"]:
            shapes.append({"type":"text","x":tx,"y":ty,"w":tw,"h":0.3,"text":it["label"],"size":11,"weight":700,"color":"#999999","align":al,"valign":"top"})
            ty += 0.32
        shapes.append({"type":"text","x":tx,"y":ty,"w":tw,"h":0.45,"text":it["head"],"size":16,"weight":700,"color":"#1A1A1A","align":al,"valign":"top"})
        ty += 0.45
        if it["desc"]:
            shapes.append({"type":"text","x":tx,"y":ty,"w":tw,"h":0.6,"text":it["desc"],"size":12,"weight":400,"color":"#666666","align":al,"valign":"top"})
    return shapes


# ─── Spec D-Build-PresetHsplit — hsplit (가로 분할 — 위 이미지/아래 텍스트) ─────
# 상단 이미지 영역 placeholder(회색 박스) + 하단 텍스트(헤드라인 + 설명).
# 한 페이지가 "주제 + 시각" 단순 구성일 때 사용 — 이미지 1장 + 본문 한 묶음.
# 입력 스키마:
#   slide_data["eyebrow"] = "좌상단 메타 라벨"        (선택)
#   slide_data["head"]    = "큰 헤드라인"             (필수)
#   slide_data["desc"]    = "한 줄 설명"              (선택)
#   slide_data["caption"] = "이미지 placeholder 안 안내" (선택, 기본 "이미지 영역")
# 안전망: head 누락 시 빈 리스트 반환.
# 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec.
def _build_preset_hsplit(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    head    = str(slide_data.get("head", "")).strip()
    desc    = str(slide_data.get("desc", "")).strip()
    caption = str(slide_data.get("caption", "")).strip() or "이미지 영역"
    if not head:
        return []
    W, H = 11.69, 8.27
    shapes = []
    if eyebrow:
        shapes.append({"type":"text","x":0.9,"y":0.5,"w":9.89,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    img_x = 0.9
    img_w = W - img_x * 2
    img_top = 1.1
    img_h = 3.8
    shapes.append({"type":"rect","x":img_x,"y":img_top,"w":img_w,"h":img_h,"fill":"#F5F5F5","stroke":"#DDDDDD","stroke_width":1})
    shapes.append({"type":"text","x":img_x,"y":img_top + img_h/2 - 0.2,"w":img_w,"h":0.4,"text":caption,"size":12,"weight":400,"color":"#BBBBBB","align":"center","valign":"middle"})
    txt_top = img_top + img_h + 0.5
    shapes.append({"type":"text","x":img_x,"y":txt_top,"w":img_w,"h":0.8,"text":head,"size":26,"weight":800,"color":"#1A1A1A","align":"left","valign":"top","role":"governing"})
    if desc:
        shapes.append({"type":"text","x":img_x,"y":txt_top + 0.9,"w":img_w,"h":0.9,"text":desc,"size":14,"weight":400,"color":"#444444","align":"left","valign":"top"})
    return shapes


# ─── Spec D-Build-PresetCircles — circles (원형 가로 정렬) 레이아웃 프리셋 ──────
# 가로로 동일 크기 원 1~4 개를 균등 배치 + 각 원 안에 큰 수치/키워드 + 아래 라벨/설명.
# 원은 "테두리만 있는 원"(fill="none") — _add_circle 의 fill 정규화 가드가 처리.
# 입력 스키마:
#   slide_data["eyebrow"] = "좌상단 메타 라벨"        (선택)
#   slide_data["title"]   = "상단 페이지 제목"        (선택)
#   slide_data["items"]   = [{"value"(필수),"label","desc"}, ...]  (1~4)
# 안전망: items list 아님 / value 누락 → 빈 리스트 반환.
# 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec (2단계).
def _build_preset_circles(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    title   = str(slide_data.get("title", "")).strip()
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        val = str(it.get("value", it.get("keyword", ""))).strip()
        if not val:
            continue
        items.append({
            "value": val,
            "label": str(it.get("label", it.get("head", ""))).strip(),
            "desc": str(it.get("desc", "")).strip(),
        })
    items = items[:4]
    if not items:
        return []
    W, H = 11.69, 8.27
    shapes = []
    # ① 배경 2장 (z-order 최하단) — 상단 검정 밴드 2.6" + 하단 흰 영역 5.67"
    #    흰 함정 회피: #1B1B1B / #FEFEFE 모두 DARK_MAP 미매핑 → light/dark 그대로 통과
    shapes.append({"type":"rect","x":0,"y":0,"w":11.69,"h":2.6,"fill":"#1B1B1B"})
    shapes.append({"type":"rect","x":0,"y":2.6,"w":11.69,"h":5.67,"fill":"#FEFEFE"})
    # ② eyebrow (선택) — 검정 밴드 안 중앙 정렬
    if eyebrow:
        shapes.append({"type":"text","x":0,"y":0.5,"w":11.69,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#999999","align":"center","valign":"top"})
    # ③ title (선택) — 검정 밴드 안 중앙 정렬, title_runs 형광 호환
    #    timeline _build_governing_block L2691-2701 패턴 정합: title_runs list 면 text_runs 키로 박음.
    #    accent 세그먼트는 _add_text(L1679)가 dark 에서 #A78BFA 자동 처리 → 본 함수 추가 코드 X.
    if title:
        title_shape = {"type":"text","x":0,"y":1.05,"w":11.69,"h":1.2,"text":title,"size":28,"weight":800,"color":"#FEFEFE","align":"center","valign":"top","role":"governing"}
        title_runs_raw = slide_data.get("title_runs")
        if isinstance(title_runs_raw, list) and title_runs_raw:
            title_shape["text_runs"] = title_runs_raw
        shapes.append(title_shape)
    # ④ 원 영역 — 흰 영역(y=2.6~8.27, h=5.67) 안에서 세로 중앙 배치
    #    묶음 높이 = circle_d + (label gap+h) + (desc gap+h). label/desc 유무로 가변.
    n = len(items)
    circle_d = 2.0
    margin = 0.9
    usable = W - margin * 2
    gap = (usable - circle_d * n) / (n - 1) if n > 1 else 0
    has_label = any(it["label"] for it in items)
    has_desc  = any(it["desc"]  for it in items)
    block_h = circle_d  # 2.0
    if has_label:
        block_h += 0.4 + 0.5  # circle→label gap + label h
    if has_desc:
        block_h += (0.05 if has_label else 0.4) + 1.0  # label→desc gap(또는 circle→desc) + desc h
    circle_y = (5.67 - block_h) / 2 + 2.6
    cy_center = circle_y + circle_d / 2
    for i, it in enumerate(items):
        cx = margin + i * (circle_d + gap)
        shapes.append({"type":"circle","x":cx,"y":circle_y,"w":circle_d,"h":circle_d,"fill":"none","stroke":"#1B1B1B","stroke_width":1.5})
        vlen = len(it["value"])
        vsize = 30 if vlen <= 4 else (22 if vlen <= 7 else 16)
        shapes.append({"type":"text","x":cx,"y":cy_center - 0.45,"w":circle_d,"h":0.9,"text":it["value"],"size":vsize,"weight":800,"color":"#1B1B1B","align":"center","valign":"middle"})
        ty = circle_y + circle_d + 0.4
        if it["label"]:
            shapes.append({"type":"text","x":cx - 0.3,"y":ty,"w":circle_d + 0.6,"h":0.5,"text":it["label"],"size":14,"weight":700,"color":"#1B1B1B","align":"center","valign":"top"})
            ty += 0.55
        if it["desc"]:
            shapes.append({"type":"text","x":cx - 0.3,"y":ty,"w":circle_d + 0.6,"h":1.0,"text":it["desc"],"size":11,"weight":400,"color":"#5A5A5A","align":"center","valign":"top"})
    return shapes


# ─── Spec Preset-New-Triad — triad (좌 거버닝 + 우 원 3 + 실선 + 하단 라벨/설명) ──
# 비대칭 입체 골격. 좌측 거버닝 (governing_block 패턴 — eyebrow + title + subtitle, w 축소)
# + 우측 원(이미지 placeholder) 3개 + 각 원 → 라벨 실선 연결 + 하단 라벨/설명 2단.
# 입력 스키마:
#   slide_data["eyebrow"]    = "좌상단 메타 라벨"        (선택)
#   slide_data["title"]      = "좌측 거버닝 메시지"      (필수)
#   slide_data["title_runs"] = 형광 segment list        (선택, accent 1곳)
#   slide_data["subtitle"]   = "서브 거버닝"            (선택)
#   slide_data["items"]      = [{"label","desc"}, ...]  (3 cap — 원 안 placeholder 이므로 value 키 없음)
# 안전망 (엄격): items list 아님 / 모든 item.label·desc 누락 → 빈 리스트 반환 (preset 미성립 → LLM 백업).
# 색 정합 (운영 흑백 6색 — DARK_MAP 자연 매핑):
#   거버닝 #1A1A1A / #444444 / #BBBBBB / 원 placeholder #ECECEC fill + #CCCCCC stroke
#   (_add_image_placeholder L1808 정합) / 실선 #DDDDDD / 라벨 #1A1A1A / desc #555555.
# theme 인자 안 받음 — circles 정합 (_add_* helper 내부 _map_color 자동 변환에 위임).
# ★ 1 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트 미연결 (다음 spec 으로 켤 것).
def _build_preset_triad(slide_data):
    eyebrow  = str(slide_data.get("eyebrow", "")).strip()
    title    = str(slide_data.get("title", "")).strip()
    subtitle = str(slide_data.get("subtitle", "")).strip()
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        items.append({
            "label": str(it.get("label", it.get("head", ""))).strip(),
            "desc":  str(it.get("desc", "")).strip(),
        })
    items = items[:3]
    if not items or not any(it["label"] or it["desc"] for it in items):
        return []
    W, H = 11.69, 8.27
    shapes = []

    # ── ① 좌측 거버닝 영역 (x 0.5 ~ 4.8) — governing_block 패턴, w 축소
    gov_w = 4.3
    if eyebrow:
        shapes.append({"type":"text","x":0.5,"y":0.5,"w":gov_w,"h":0.4,
                       "text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB","align":"left","valign":"top"})
    if title:
        t = {"type":"text","x":0.5,"y":2.2,"w":gov_w,"h":3.0,
             "text":title,"size":28,"weight":800,"color":"#1A1A1A","align":"left","valign":"top",
             "role":"governing"}
        tr = slide_data.get("title_runs")
        if isinstance(tr, list) and tr:
            t["text_runs"] = tr
        shapes.append(t)
    if subtitle:
        shapes.append({"type":"text","x":0.5,"y":5.0,"w":gov_w,"h":1.0,
                       "text":subtitle,"size":14,"weight":500,"color":"#444444","align":"left","valign":"top"})

    # ── ② 우측 3요소 영역 (x 5.3 ~ 11.2) — 셀 균등 분배, 셀당 원+선+라벨+desc 수직 스택
    area_x0, area_x1 = 5.3, 11.2
    n = len(items)
    cell_w = (area_x1 - area_x0) / n
    circle_d = 1.5
    cy = 3.1                       # 원 top
    cy_center = cy + circle_d / 2
    label_y = 5.4
    for i, it in enumerate(items):
        cell_x = area_x0 + cell_w * i
        cx = cell_x + cell_w / 2   # 셀 가로 중앙
        # 원 (이미지 placeholder — _add_image_placeholder 색 정합: 회색 fill + stroke + 안내문)
        shapes.append({"type":"circle","x":cx - circle_d / 2,"y":cy,"w":circle_d,"h":circle_d,
                       "fill":"#ECECEC","stroke":"#CCCCCC","stroke_width":0.75})
        shapes.append({"type":"text","x":cx - circle_d / 2,"y":cy_center - 0.2,"w":circle_d,"h":0.4,
                       "text":"이미지","size":11,"weight":400,"color":"#888888","align":"center","valign":"middle"})
        # 실선 연결 (원 하단 → 라벨 위) — 점선 인프라 없음, 1단계는 실선 (zigzag L2961 / hero_cards L3197 정합)
        shapes.append({"type":"line","x1":cx,"y1":cy + circle_d,"x2":cx,"y2":label_y - 0.1,
                       "color":"#DDDDDD","width":1})
        # 하단 라벨 (size 14 / 700 / #1A1A1A) + desc (size 12 / 400 / #555555) — 2단 (circles 정합)
        if it["label"]:
            shapes.append({"type":"text","x":cell_x,"y":label_y,"w":cell_w,"h":0.5,
                           "text":it["label"],"size":14,"weight":700,"color":"#1A1A1A","align":"center","valign":"top"})
        if it["desc"]:
            shapes.append({"type":"text","x":cell_x,"y":label_y + 0.55,"w":cell_w,"h":1.5,
                           "text":it["desc"],"size":12,"weight":400,"color":"#555555","align":"center","valign":"top"})
    return shapes


# ─── Spec Preset-New-StrategyMap — strategy_map (대전략 3 + 하위실행 5, 4단 분할) ──
# 한 슬라이드에 (상) 거버닝 + (중상) 박스 3 + "+" 결합 + (중하) chevron 5 흐름 + (하) 마무리 거버닝.
# 대전략 + 하위실행을 한 화면 조망 — 비교적 복잡한 골격 (도형 약 23~28개).
# 입력 스키마:
#   slide_data["eyebrow"]      = "상단 메타 라벨"          (선택)
#   slide_data["title"]        = "상단 메인 거버닝"        (필수)
#   slide_data["title_runs"]   = 형광 segment list       (선택, accent 1곳)
#   slide_data["pillars"]      = [{"head"(필수),"desc","caption"}, ...]  (정확히 3개)
#   slide_data["steps"]        = [{"label"(필수)}, ...]                 (3~6개)
#   slide_data["footer"]       = "하단 마무리 거버닝"      (선택)
# 안전망 (엄격): title 없음 / pillars 정확히 3개 아님 / steps 3개 미만 → 빈 list (LLM 백업).
# 색 정합 (운영 흑백 6색 — DARK_MAP 자연 매핑):
#   거버닝 #1A1A1A / 박스 fill #FFFFFF + stroke #DDDDDD / chevron stroke #1A1A1A /
#   "+" #1A1A1A / footer #444444 / caption #999999 / desc #555555.
# theme 인자 안 받음 — circles/triad 정합 (_add_* helper 내부 _map_color 자동 변환에 위임).
# 흐름은 chevron 자체가 → 모양이라 화살표(arrow) 별도 없음 — process L2220 패턴 정합.
# ★ 1 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트 미연결 (다음 spec 으로 켤 것).
def _build_preset_strategy_map(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    title   = str(slide_data.get("title", "")).strip()
    footer  = str(slide_data.get("footer", "")).strip()
    pillars_raw = slide_data.get("pillars") or []
    steps_raw   = slide_data.get("steps") or []
    if not title:
        return []
    if not isinstance(pillars_raw, list):
        return []
    # ★ Spec Strategy-Map-Expand-And-Loosen — 정확 3 강제 폐지. 2~3 허용.
    pillars = [p for p in pillars_raw if isinstance(p, dict)][:3]
    if len(pillars) < 2:
        return []
    steps = []
    for s in (steps_raw if isinstance(steps_raw, list) else []):
        if isinstance(s, dict):
            lab = str(s.get("label", "")).strip()
            if lab:
                steps.append(lab)
    steps = steps[:6]
    # ★ Spec Strategy-Map-Expand-And-Loosen — 최소 3 강제 폐지. 2~6 허용.
    if len(steps) < 2:
        return []

    W, H = 11.69, 8.27
    margin = 0.9
    usable = W - 2 * margin
    shapes = []

    # ── ① 상단 거버닝 (eyebrow + title)
    if eyebrow:
        shapes.append({"type":"text","x":margin,"y":0.5,"w":usable,"h":0.4,
                       "text":eyebrow,"size":11,"weight":400,"color":"#BBBBBB",
                       "align":"center","valign":"top"})
    t = {"type":"text","x":margin,"y":0.9,"w":usable,"h":0.9,
         "text":title,"size":28,"weight":800,"color":"#1A1A1A",
         "align":"center","valign":"top",
         "role":"governing"}
    tr = slide_data.get("title_runs")
    if isinstance(tr, list) and tr:
        t["text_runs"] = tr
    shapes.append(t)

    # ── ② 박스 N개(2 또는 3) + "+" 결합 (Spec Strategy-Map-Expand-And-Loosen)
    box_y, box_h = 1.9, 1.9
    n_p = len(pillars)
    box_w = 2.9
    # n_p=3 → gap=0.595 (기존), n_p=2 → gap=0.6 + 좌우 중앙 정렬 (묶음 usable 미만 → 중앙).
    if n_p == 2:
        gap_p = 0.6
        total_p = box_w * n_p + gap_p * (n_p - 1)   # 6.4
        start_p_x = margin + (usable - total_p) / 2
    else:  # n_p == 3
        gap_p = (usable - box_w * n_p) / (n_p - 1)  # 0.595
        start_p_x = margin
    for i, p in enumerate(pillars):
        bx = start_p_x + i * (box_w + gap_p)
        # 박스 본체
        shapes.append({"type":"rect","x":bx,"y":box_y,"w":box_w,"h":box_h,
                       "fill":"#FFFFFF","stroke":"#DDDDDD","stroke_width":1})
        # head (16/700/#1A1A1A)
        shapes.append({"type":"text","x":bx + 0.15,"y":box_y + 0.2,"w":box_w - 0.3,"h":0.5,
                       "text":str(p.get("head","")).strip(),
                       "size":16,"weight":700,"color":"#1A1A1A",
                       "align":"center","valign":"top"})
        # desc (12/400/#555)
        shapes.append({"type":"text","x":bx + 0.15,"y":box_y + 0.75,"w":box_w - 0.3,"h":box_h - 0.9,
                       "text":str(p.get("desc","")).strip(),
                       "size":12,"weight":400,"color":"#555555",
                       "align":"center","valign":"top"})
        # caption (박스 밖 아래, 11/400/#999)
        cap = str(p.get("caption","")).strip()
        if cap:
            shapes.append({"type":"text","x":bx,"y":box_y + box_h + 0.1,"w":box_w,"h":0.5,
                           "text":cap,"size":11,"weight":400,"color":"#999999",
                           "align":"center","valign":"top"})
        # "+" 기호 (박스 사이 세로중앙, 마지막 박스 제외 → 2개)
        if i < n_p - 1:
            plus_x = bx + box_w
            shapes.append({"type":"text","x":plus_x,"y":box_y + box_h / 2 - 0.3,"w":gap_p,"h":0.6,
                           "text":"+","size":32,"weight":800,"color":"#1A1A1A",
                           "align":"center","valign":"middle"})

    # ── ③ 흐름 — 원(이미지 placeholder) + 화살표 + 라벨 (Spec Strategy-Map-Expand-And-Loosen: 2~6)
    flow_y = 5.2  # 원 top
    n_s = len(steps)
    # 원 직경 동적 — n 많아질수록 작게 (가용 영역 안 화살표 공간 확보)
    if n_s <= 3:
        circle_d = 1.3
    elif n_s == 4:
        circle_d = 1.2
    elif n_s == 5:
        circle_d = 1.1
    else:  # 6
        circle_d = 0.95
    # 균등 분배 (원 + 화살표 사이 gap). n_s=2 는 usable 전폭에 원 2 개면 gap 이 너무 벌어져
    # 화살표가 길게 늘어지므로 gap 고정 + 좌우 중앙 정렬 (pillars n_p=2 와 동일 패턴).
    if n_s == 2:
        gap_s = 3.0
        total_s = circle_d * n_s + gap_s * (n_s - 1)   # 5.6
        start_s_x = margin + (usable - total_s) / 2
    else:
        gap_s = (usable - circle_d * n_s) / (n_s - 1)
        start_s_x = margin
    cy_center = flow_y + circle_d / 2
    label_y = flow_y + circle_d + 0.15
    for i, label in enumerate(steps):
        cell_x = start_s_x + i * (circle_d + gap_s)  # 원 좌측
        # 원 (이미지 placeholder — triad 색 정합: #ECECEC fill + #CCCCCC stroke)
        shapes.append({"type":"circle","x":cell_x,"y":flow_y,"w":circle_d,"h":circle_d,
                       "fill":"#ECECEC","stroke":"#CCCCCC","stroke_width":0.75})
        # 원 안 안내문 (triad 정합 — 11/400/#888 center)
        shapes.append({"type":"text","x":cell_x,"y":cy_center - 0.2,"w":circle_d,"h":0.4,
                       "text":"이미지","size":11,"weight":400,"color":"#888888",
                       "align":"center","valign":"middle"})
        # 원 아래 라벨 (14/700/#1A1A1A center)
        shapes.append({"type":"text","x":cell_x,"y":label_y,"w":circle_d,"h":0.4,
                       "text":label,"size":14,"weight":700,"color":"#1A1A1A",
                       "align":"center","valign":"top"})
        # 화살표 (원 사이 가로, 마지막 제외 = n-1개) — y = 원 세로중앙
        if i < n_s - 1:
            x1 = cell_x + circle_d + 0.1                    # 원 우측 + 여유
            x2 = cell_x + circle_d + gap_s - 0.1            # 다음 원 좌측 - 여유
            shapes.append({"type":"arrow","x1":x1,"y1":cy_center,"x2":x2,"y2":cy_center,
                           "color":"#1A1A1A","width":1.5})

    # ── ④ 하단 거버닝 (선택)
    if footer:
        shapes.append({"type":"text","x":margin,"y":7.2,"w":usable,"h":0.7,
                       "text":footer,"size":16,"weight":500,"color":"#444444",
                       "align":"center","valign":"top"})
    return shapes


# ─── Spec D-Build-PresetHeroCards — hero_cards (상단 거버닝 히어로 + 카드 N개) ────
# 상단 검정 밴드(거버닝 중앙정렬, title_runs 노랑 호환) + 거버닝 아래 ↓ 화살표 +
# 하단 흰 영역에 카드 N개(2~4, 5+ cap) 가로 배치.
# 각 카드 = 위 배지 라벨(선택) + 카드 본체(살짝 둥근 rect) 안에 제목 + 짧은 구분선 + 설명.
# circles 패턴 정합: 배경 2분할 + 거버닝 중앙정렬 + items 가로 분배 + 묶음 세로 중앙 +
# title_runs 직접 처리(_build_governing_block 헬퍼 X — 검정 밴드 안 흰 거버닝 자체 처리).
# 입력 스키마:
#   slide_data["eyebrow"]    = "상단 메타 라벨"               (선택)
#   slide_data["title"]      = "상단 거버닝 메시지"           (필수)
#   slide_data["title_runs"] = 형광 segment list             (선택, accent:true 1곳)
#   slide_data["items"]      = [{"label"(선택 배지), "head"(필수), "desc"(선택)}, ...]
#                              (2~4 권장, 5+ cap, 0개 → 빈 리스트 fallback)
# 안전망 (엄격): items list 아님 / 모든 item.head 누락 → 빈 리스트 반환 (preset 미성립 → LLM 백업).
# 색 정합 (circles 흰 함정 회피 — DARK_MAP 미매핑/자연매핑):
#   검정 배경 #1B1B1B / 흰 배경 #FEFEFE / 카드 fill #FBFBFB (미매핑) /
#   카드 stroke #DDDDDD (stroke role 자연 매핑 → dark #2A2A2A).
#   글자 #FEFEFE / #1B1B1B / #999999(eyebrow) / #5A5A5A(desc) 미매핑 회색.
# ★ 1 단계: 코드만 등록 — viz_pattern 화이트리스트 미연결 (다음 spec 으로 켤 것).
def _build_preset_hero_cards(slide_data):
    eyebrow = str(slide_data.get("eyebrow", "")).strip()
    title   = str(slide_data.get("title", "")).strip()
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        items.append({
            "label": str(it.get("label", "")).strip(),
            "head": head,
            "desc": str(it.get("desc", "")).strip(),
        })
    items = items[:4]  # 5+ cap (circles 패턴 정합)
    if not items:
        return []
    W, H = 11.69, 8.27
    shapes = []
    # ① 배경 2장 (z-order 최하단, circles 정합) — 검정 밴드 2.6" + 흰 영역 5.67". 미매핑 색.
    shapes.append({"type":"rect","x":0,"y":0,"w":11.69,"h":2.6,"fill":"#1B1B1B"})
    shapes.append({"type":"rect","x":0,"y":2.6,"w":11.69,"h":5.67,"fill":"#FEFEFE"})
    # ② eyebrow (선택) — 검정 밴드 안 중앙정렬
    if eyebrow:
        shapes.append({"type":"text","x":0,"y":0.5,"w":11.69,"h":0.4,"text":eyebrow,"size":11,"weight":400,"color":"#999999","align":"center","valign":"top"})
    # ③ title (선택) — 검정 밴드 안 중앙정렬, title_runs 형광 호환 (circles 패턴 정합).
    #   timeline _build_governing_block L2691-2701 패턴: title_runs list → text_runs 키.
    #   accent 세그먼트는 _add_text L1679 분기에서 dark 시 #A78BFA 자동 처리 → 추가 코드 X.
    if title:
        title_shape = {"type":"text","x":0,"y":1.0,"w":11.69,"h":1.0,"text":title,"size":28,"weight":800,"color":"#FEFEFE","align":"center","valign":"middle","role":"governing"}
        title_runs_raw = slide_data.get("title_runs")
        if isinstance(title_runs_raw, list) and title_runs_raw:
            title_shape["text_runs"] = title_runs_raw
        shapes.append(title_shape)
    # ④ ↓ 화살표 (검정 밴드 안 하단 중앙) — arrow type, 가로중앙 W/2, y=2.05~2.45 (검정밴드 하단부)
    shapes.append({"type":"arrow","x1":W/2,"y1":2.05,"x2":W/2,"y2":2.45,"color":"#FEFEFE","width":1.5})
    # ⑤ 카드 가로 분배 — circles 정합 (margin + gap) + card_w cap 3.2 (n=2 너무 안 넓어지게)
    margin = 0.9
    usable = W - margin * 2  # 9.89
    n = len(items)
    gap_min = 0.4
    if n == 1:
        card_w = min(usable, 3.2)
    else:
        card_w_raw = (usable - gap_min * (n - 1)) / n
        card_w = min(card_w_raw, 3.2)
    # cap 발동 시 묶음 너비 < usable → 좌우 중앙 정렬로 시작 x 재계산
    total_w = card_w * n + gap_min * (n - 1)
    start_x = margin + (usable - total_w) / 2
    # ⑥ 카드 묶음 세로 중앙 (흰 영역 5.67 안) — circles block_h 패턴 정합. desc 유무로 카드 h 가변.
    has_desc = any(it["desc"] for it in items)
    has_label = any(it["label"] for it in items)
    card_h = 2.6 if has_desc else 1.8
    badge_h = 0.4
    badge_pad = 0.1
    block_h = (badge_h + badge_pad if has_label else 0) + card_h
    block_top = (5.67 - block_h) / 2 + 2.6
    card_y = block_top + (badge_h + badge_pad if has_label else 0)
    badge_y = block_top
    # ⑦ 각 카드 그리기 (배지(선택) → 본체 → 제목 → 구분선 → 설명(선택))
    for i, it in enumerate(items):
        cx = start_x + i * (card_w + gap_min)
        cx_center = cx + card_w / 2
        # ⑦-a 배지 라벨 (label 있을 때) — 카드 가로중앙, 살짝 둥근 검정 rect + 흰 글자.
        if it["label"]:
            badge_w = 1.4
            badge_x = cx_center - badge_w / 2
            shapes.append({"type":"rect","x":badge_x,"y":badge_y,"w":badge_w,"h":badge_h,"fill":"#1B1B1B","radius":0.06})
            shapes.append({"type":"text","x":badge_x,"y":badge_y,"w":badge_w,"h":badge_h,"text":it["label"],"size":11,"weight":700,"color":"#FEFEFE","align":"center","valign":"middle"})
        # ⑦-b 카드 본체 — 살짝 둥근(radius 0.08), fill #FBFBFB(미매핑 옅은 회색) + stroke #DDDDDD(자연 매핑).
        shapes.append({"type":"rect","x":cx,"y":card_y,"w":card_w,"h":card_h,"fill":"#FBFBFB","stroke":"#DDDDDD","stroke_width":1,"radius":0.08})
        # ⑦-c 제목 (head) — 카드 상단, #1B1B1B, size 16, weight 800, center
        head_y = card_y + 0.3
        shapes.append({"type":"text","x":cx,"y":head_y,"w":card_w,"h":0.5,"text":it["head"],"size":16,"weight":800,"color":"#1B1B1B","align":"center","valign":"top"})
        # ⑦-d 짧은 구분선 (line) — 제목 아래 가로중앙 짧게, #1B1B1B
        sep_w = 0.6
        sep_y = head_y + 0.6
        shapes.append({"type":"line","x1":cx_center - sep_w/2,"y1":sep_y,"x2":cx_center + sep_w/2,"y2":sep_y,"color":"#1B1B1B","width":1})
        # ⑦-e 설명 (desc, 선택) — 구분선 아래, #5A5A5A(미매핑 회색), size 11, weight 400, center
        if it["desc"]:
            desc_y = sep_y + 0.15
            desc_h = card_y + card_h - desc_y - 0.2  # 카드 하단까지 여유 0.2"
            shapes.append({"type":"text","x":cx + 0.15,"y":desc_y,"w":card_w - 0.3,"h":desc_h,"text":it["desc"],"size":11,"weight":400,"color":"#5A5A5A","align":"center","valign":"top"})
    return shapes


# ─── Spec Preset-Remove-Matrix — _build_preset_matrix 제거 자리 (옵션 B 완전삭제) ──
# 본 위치에 있던 matrix preset 함수는 2×2 사분면 범용성 낮음 + 첫 생성 빈 페이지 버그로
# 제거. quad/cards3 등으로 대체 가능. 화이트리스트/카탈로그/dispatch/SLIDE 분기/거버닝
# 헬퍼 (_build_governing_block, 유일 사용자) 모두 동시 삭제.
# ─── Spec D-Build-PresetHsplitTop — hsplit_top (위 검정 거버닝 + 아래 흰 좌우 2항목) ─
# split (세로 좌우 색면 2분할) 의 가로 비대칭 버전 — 위 1/3 검정 색면 + 아래 2/3 흰 바탕.
# 위 검정 영역: 중앙 정렬 거버닝 (eyebrow + 메인 1줄 + 형광 2줄 + subtitle).
# 아래 흰 영역: 좌/우 2항목 (각 head + body), 선·라벨·구분선 없음 (위치로만 구분).
#
# ★ 본 preset 은 거버닝 헬퍼 _build_governing_block 사용 X —
#   헬퍼는 흰 배경·검정 글자·좌측 정렬 가정. 본 preset 은 검정 배경·흰 글자·중앙 정렬 →
#   거버닝을 함수 안에서 직접 그림.
#
# ★ 거버닝 강조: title (1줄) + title_accent (2줄) 둘 다 role:"governing" 마커.
#   Spec Preset-Purple-Accent-Cleanup — 이전엔 title_accent 만 하드코딩 #A78BFA 로 색차 유도,
#   그러나 라이트에서 title #6B46E5 와 title_accent #A78BFA 가 두 톤으로 튐. 마커 통일로 해결.
#   dispatch (L2027-2036) 가 role:"governing" 도형을 _get_theme(theme)["ACCENT"] 로 override →
#   라이트 #6B46E5 / 다크 #A78BFA. 2줄 시각 임팩트는 위치·크기 그대로.
#
# ★ z-order (shapes 리스트 순서 = 렌더 순서, 먼저 = 아래):
#   ① 검정 rect → ② eyebrow → ③ 메인 1줄 → ④ 메인 2줄(형광) → ⑤ subtitle
#   → ⑥ 좌 head → ⑦ 좌 body → ⑧ 우 head → ⑨ 우 body
#
# ★ dark 자동 반전 (split 패턴 정합):
#   #1A1A1A(fill) ↔ #FFFFFF(text) 자동 반전 → light "검정 면+흰 글자" ↔ dark "흰 면+검정 글자".
#   #999999(text), #1A1A1A(text), #444444(text) 도 DARK_MAP 매핑으로 안전.
#   subtitle #999999 채택 (#CCCCCC 는 DARK_MAP 매핑 없어 dark 흰 면 위 가독성 약함).
#
# 입력 스키마:
#   slide_data["eyebrow"]      = "위 메타 라벨" (선택, color #999999)
#   slide_data["title"]        = "메인 거버닝 1줄 (흰 글자)" (필수)
#   slide_data["title_accent"] = "거버닝 2줄 강조 키워드 (role:'governing' → ACCENT)" (선택)
#   slide_data["subtitle"]     = "서브 거버닝" (선택, color #999999)
#   slide_data["left"]         = {"head"(필수), "body"} (필수)
#   slide_data["right"]        = {"head"(필수), "body"} (필수)
# 안전망:
#   - title 누락 → 빈 list (거버닝 본질)
#   - left/right 둘 중 하나라도 dict 아니거나 head 누락 → 빈 list (좌우 비교 본질)
def _build_preset_hsplit_top(slide_data):
    title = str(slide_data.get("title", "")).strip()
    if not title:
        return []   # 거버닝 본질 — title 없으면 미성립

    left  = slide_data.get("left")  or {}
    right = slide_data.get("right") or {}
    if not isinstance(left, dict) or not isinstance(right, dict):
        return []
    l_head = str(left.get("head", "")).strip()
    r_head = str(right.get("head", "")).strip()
    if not l_head or not r_head:
        return []   # 좌우 비교 본질 — 둘 다 head 필수

    eyebrow      = str(slide_data.get("eyebrow", "")).strip()
    title_accent = str(slide_data.get("title_accent", "")).strip()
    subtitle     = str(slide_data.get("subtitle", "")).strip()
    l_body = str(left.get("body", "")).strip()
    r_body = str(right.get("body", "")).strip()

    W, H = 11.69, 8.27
    BLACK_H = H / 3  # 2.7567 — 위 검정 영역 정확히 1/3

    shapes = []

    # ① 검정 색면 rect (가장 먼저 = z-order 맨 아래)
    shapes.append({
        "type": "rect",
        "x": 0, "y": 0, "w": W, "h": BLACK_H,
        "fill": "#1A1A1A",
    })

    # ② eyebrow (선택, 검정 영역 안 중앙 정렬, 옅은 회색)
    if eyebrow:
        shapes.append({
            "type": "text",
            "x": 0, "y": 0.5, "w": W, "h": 0.4,
            "text": eyebrow,
            "size": 11, "weight": 400, "color": "#999999",
            "align": "center", "valign": "top",
        })

    # ③ 메인 거버닝 1줄 (필수, 흰 글자, 중앙)
    shapes.append({
        "type": "text",
        "x": 0, "y": 1.0, "w": W, "h": 0.7,
        "text": title,
        "size": 28, "weight": 800, "color": "#FFFFFF",
        "align": "center", "valign": "middle",
        "role": "governing",
    })

    # ④ 메인 거버닝 2줄 (선택, role:"governing" → 라이트 #6B46E5 / 다크 #A78BFA 자동)
    #   Spec Preset-Purple-Accent-Cleanup — 하드코딩 #A78BFA 제거, title 과 동일하게 마커 통일.
    #   color:"#FFFFFF" 는 마커 override 되기 전 fallback 힌트 (title 과 정합).
    if title_accent:
        shapes.append({
            "type": "text",
            "x": 0, "y": 1.7, "w": W, "h": 0.7,
            "text": title_accent,
            "size": 28, "weight": 800, "color": "#FFFFFF",
            "align": "center", "valign": "middle",
            "role": "governing",
        })

    # ⑤ subtitle (선택, 검정 영역 끝쪽, #999999 — dark 가독성 위해 #CCC 회피)
    if subtitle:
        shapes.append({
            "type": "text",
            "x": 0, "y": 2.40, "w": W, "h": 0.3,
            "text": subtitle,
            "size": 13, "weight": 500, "color": "#999999",
            "align": "center", "valign": "top",
        })

    # ⑥~⑨ 아래 흰 영역 좌/우 2항목 (선·구분선 없음, 위치로만 구분)
    #   좌측: x=0.7, w=4.8 / 우측: x=6.2, w=4.8 / 사이 0.7" 시각 호흡
    LX, LW = 0.7, 4.8
    RX, RW = 6.2, 4.8

    # ⑥ 좌 head
    shapes.append({
        "type": "text",
        "x": LX, "y": 3.3, "w": LW, "h": 0.8,
        "text": l_head,
        "size": 22, "weight": 800, "color": "#1A1A1A",
        "align": "left", "valign": "top",
    })
    # ⑦ 좌 body
    if l_body:
        shapes.append({
            "type": "text",
            "x": LX, "y": 4.2, "w": LW, "h": 3.0,
            "text": l_body,
            "size": 14, "weight": 400, "color": "#444444",
            "align": "left", "valign": "top",
        })

    # ⑧ 우 head
    shapes.append({
        "type": "text",
        "x": RX, "y": 3.3, "w": RW, "h": 0.8,
        "text": r_head,
        "size": 22, "weight": 800, "color": "#1A1A1A",
        "align": "left", "valign": "top",
    })
    # ⑨ 우 body
    if r_body:
        shapes.append({
            "type": "text",
            "x": RX, "y": 4.2, "w": RW, "h": 3.0,
            "text": r_body,
            "size": 14, "weight": 400, "color": "#444444",
            "align": "left", "valign": "top",
        })

    return shapes


# ─── Spec D-Build-PresetQuad — quad (색면 4분할 흑/백/흑/백) 레이아웃 프리셋 ────
# split(2분할) 의 확장: 세로 4면을 흑/백/흑/백으로 번갈아 칠하고,
# 각 면 안에 keyword(상단 라벨) + head(메인) + desc(부연)를 가운데정렬.
# 면 색에 따라 글자색 자동 결정 — 검정 면 → 흰 글자, 흰 면 → 검정 글자.
# 다크 테마에서는 _map_color 가 fill/text 각각 role-aware 반전 → split 처럼 자연 반전.
# 입력 스키마:
#   slide_data["items"] = [{"keyword","head"(필수 중 하나),"desc"}, ...]  (2~4, alias: "columns")
# 안전망: items list 아님 / keyword+head 모두 비어있음 / 유효 items < 2 → 빈 리스트 반환.
# 1 단계는 코드만 등록 — viz_pattern 연결은 별도 spec (2단계).
def _build_preset_quad(slide_data):
    items_raw = slide_data.get("items") or slide_data.get("columns") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        kw = str(it.get("keyword", it.get("label", ""))).strip()
        head = str(it.get("head", "")).strip()
        desc = str(it.get("desc", "")).strip()
        if not (kw or head):
            continue
        items.append({"keyword": kw, "head": head, "desc": desc})
    items = items[:4]
    if len(items) < 2:
        return []
    W, H = 11.69, 8.27
    n = len(items)
    panel_w = W / n
    shapes = []
    fills = ["#1A1A1A", "#FFFFFF", "#1A1A1A", "#FFFFFF"]
    for i, it in enumerate(items):
        px = i * panel_w
        fill = fills[i] if n == 4 else fills[i % 2]
        shapes.append({"type":"rect","x":px,"y":0,"w":panel_w,"h":H,"fill":fill})
        is_dark_panel = fill == "#1A1A1A"
        kw_color   = "#FFFFFF" if is_dark_panel else "#1A1A1A"
        head_color = "#FFFFFF" if is_dark_panel else "#1A1A1A"
        desc_color = "#DDDDDD" if is_dark_panel else "#666666"
        pad = 0.45
        tw = panel_w - pad * 2
        if it["keyword"]:
            shapes.append({"type":"text","x":px+pad,"y":1.5,"w":tw,"h":0.5,
                           "text":it["keyword"],"size":13,"weight":500,"color":kw_color,
                           "align":"center","valign":"top"})
        if it["head"]:
            shapes.append({"type":"text","x":px+pad,"y":3.2,"w":tw,"h":2.0,
                           "text":it["head"],"size":20,"weight":800,"color":head_color,
                           "align":"center","valign":"top"})
        if it["desc"]:
            shapes.append({"type":"text","x":px+pad,"y":5.6,"w":tw,"h":2.0,
                           "text":it["desc"],"size":11,"weight":400,"color":desc_color,
                           "align":"center","valign":"top"})
    return shapes


# ─── Spec D-Build-PresetDivider — 챕터 간지(divider) 코드 고정 템플릿 ──────────
# 챕터 간지(Ⅰ/Ⅱ/...) 5개가 LLM 자유 출력으로 매번 제각각 나오는 문제 해결.
# 식별: slide_type="hero" AND (section "(챕터 divider)" OR governing_main "Ⅰ. " 패턴).
# 콘셉트 슬로건("Ⅰ.4 콘셉트")/강조 hero("Ⅰ.3 ...")는 section 형식 차이로 자동 제외.
# 디자인: 왼쪽 정렬, 로마숫자(큰) + 부문명(중간)을 한 줄로, 부제/분할선/페이지번호 X.
# ★ _add_text 의 text_runs 는 run 별 size 다르게 지원 X (1-d-① 색만 분기) →
#   두 개의 text 도형으로 분리 + valign="bottom" baseline 정렬.

# Spec D-Fix-DividerPrefixStrip — section/governing_main 맨 앞 "번호 prefix" 광범위 흡수.
#   기존: "로마숫자(Ⅰ~Ⅴ / I~V) + 점 + 공백" 만 매칭 → 변형 케이스(점 뒤 공백 X / 점 X /
#         Ⅵ 이상 / 아라비아 숫자) 미매칭 → fallback 시 LLM prefix 잔존 → 코드 자동 매김
#         prefix 와 겹침("Ⅰ Ⅲ. 공간구성계획" 사고 원인, 진단 D-Check-DividerRomanCollision).
#   확장: 로마숫자 Ⅰ~Ⅻ(유니코드 단일 문자) / I~XII(ASCII 긴 것 먼저, IX/XI 우선) /
#         아라비아 1~99 + 구분자 [.\s]* (0개 이상 — 점/공백/없음 모두 흡수).
#   호환: group(1) = 번호 prefix / group(2) = 본문 (기존 _extract_divider_title L3068-3070
#         의 m.group(2) 사용 코드와 그대로 호환).
#   주의: ASCII alternation 순서 — XII 가 XI 보다 먼저, IX 가 X 보다 먼저, IV 가 I/II/III 보다
#         먼저 와야 긴 것이 우선 매칭 (greedy alternation 보장).
_DIVIDER_ROMAN_RE = re.compile(
    r"^("
    r"[ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩⅪⅫ]"           # 유니코드 로마숫자 1~12 (단일 문자)
    r"|XII|VIII|VII|XI|IX|IV|III|VI|II|X|V|I"  # ASCII 로마 (긴 것·우선순위 높은 것 먼저)
    r"|\d{1,2}"                          # 아라비아 1~99
    r")"
    r"[.\s]*"                             # 구분자 0개 이상 (점·공백 임의 조합)
    r"(.+)$"                              # 본문 (최소 1자)
)
_SECTION_DOT_DIGIT_RE = re.compile(r"^.{1,5}\.\d")

# Spec D-Fix-DividerActuallyWorks — 1-based 순번 → 유니코드 로마숫자.
# LLM 이 로마숫자를 자유 매김(V 중복 등) 하는 환각 차단용. 코드가 등장 순서대로 부여.
_DIVIDER_ROMAN_NUMERALS = ["Ⅰ", "Ⅱ", "Ⅲ", "Ⅳ", "Ⅴ", "Ⅵ", "Ⅶ", "Ⅷ", "Ⅸ", "Ⅹ"]


def _divider_idx_to_roman(idx):
    """1-based 정수 → 유니코드 로마숫자. 11+ 는 "11", "12" 형태 fallback (안전망)."""
    if not isinstance(idx, int) or idx < 1:
        return "?"
    if idx <= len(_DIVIDER_ROMAN_NUMERALS):
        return _DIVIDER_ROMAN_NUMERALS[idx - 1]
    return str(idx)


# Spec D-Fix-DividerNumberFromLLM — ASCII 로마 → 정수 매핑 (정규화용).
# LLM 이 ASCII 로마(I, II, III, IV, V, ...) 또는 아라비아(1, 2, 3, ...) 어떤 형식으로
# prefix 박아도 유니코드 로마(Ⅰ, Ⅱ, ...) 로 통일해 간지에 표시.
_ASCII_ROMAN_TO_INT = {
    "I": 1, "II": 2, "III": 3, "IV": 4, "V": 5, "VI": 6,
    "VII": 7, "VIII": 8, "IX": 9, "X": 10, "XI": 11, "XII": 12,
}

# 정규화 입력용 — 유니코드 로마 1~12. _DIVIDER_ROMAN_NUMERALS(L3018, Ⅰ~Ⅹ 10개) 는
# _divider_idx_to_roman 의 출력 매핑(fallback) 용이라 무변경. 본 dict 는 LLM 이
# 박은 prefix 가 Ⅺ·Ⅻ 일 때도 정규화 통과시키기 위한 입력 인식용(별도 책임 분리).
_UNICODE_ROMAN_TO_INT = {
    "Ⅰ": 1, "Ⅱ": 2, "Ⅲ": 3, "Ⅳ": 4, "Ⅴ": 5, "Ⅵ": 6,
    "Ⅶ": 7, "Ⅷ": 8, "Ⅸ": 9, "Ⅹ": 10, "Ⅺ": 11, "Ⅻ": 12,
}


def _normalize_divider_prefix_to_roman(prefix):
    """LLM 박은 prefix(유니코드 로마 / ASCII 로마 / 아라비아) → 유니코드 로마숫자 정규화.

    Spec D-Fix-DividerNumberFromLLM:
      간지 번호를 코드 자동 매김(_divider_counter) 이 아니라 LLM 박은 prefix 그대로 사용해
      목차·소항목과 같은 출처 보장 (목차 Ⅲ + 간지 Ⅲ 일치).

    매핑:
      유니코드 로마(Ⅰ~Ⅻ) → 그대로
      ASCII 로마(I~XII)    → 유니코드 로마
      아라비아(1~12)       → 유니코드 로마
      아라비아(13~)         → "13", "14" 형태 fallback (안전망, _divider_idx_to_roman 와 동일)
      매칭 실패              → "" (caller 가 fallback 카운터 사용)
    """
    if not prefix or not isinstance(prefix, str):
        return ""
    p = prefix.strip()
    if not p:
        return ""
    # 유니코드 로마(Ⅰ~Ⅻ) — 1~10 은 _divider_idx_to_roman 출력과 동형이라 그대로,
    # Ⅺ·Ⅻ 도 인식해 그대로 반환 (입력 보존). _UNICODE_ROMAN_TO_INT 가 1~12 다 커버.
    if p in _UNICODE_ROMAN_TO_INT:
        return p
    # ASCII 로마 → 정수 → 유니코드 로마
    n = _ASCII_ROMAN_TO_INT.get(p)
    if n is not None:
        # 11~12 는 _divider_idx_to_roman fallback("11"/"12") 대신 Ⅺ·Ⅻ 직접 매핑
        # (입력 의도 보존 — LLM 이 XI 박았는데 출력에서 "11" 로 강등되면 어긋남).
        if n == 11:
            return "Ⅺ"
        if n == 12:
            return "Ⅻ"
        return _divider_idx_to_roman(n)
    # 아라비아 → 정수 → 유니코드 로마 (1~12). 13+ 는 _divider_idx_to_roman 의 숫자 fallback.
    if p.isdigit():
        try:
            n = int(p)
            if n == 11:
                return "Ⅺ"
            if n == 12:
                return "Ⅻ"
            if n >= 1:
                return _divider_idx_to_roman(n)
        except ValueError:
            pass
    return ""


def _extract_divider_roman(slide_data):
    """LLM 이 section / governing_main 에 박은 prefix → 유니코드 로마숫자.

    Spec D-Fix-DividerNumberFromLLM (간지 번호 ↔ 목차 번호 일치 보장):
      _extract_divider_title 이 이미 _DIVIDER_ROMAN_RE 로 prefix(group 1) + 제목(group 2)
      둘 다 매칭하지만 제목만 사용. 본 함수가 group 1 (LLM 박은 번호) 을 정규화해 반환.

    우선순위:
      1. governing_main 에서 매칭 → group(1) 정규화 (가장 정제된 LLM 출력)
      2. section 에서 "(챕터 divider)" 라벨 떼고 매칭 → group(1) 정규화
      3. 둘 다 실패 (LLM 이 prefix 안 박음 / 잘못 박음) → "" 반환
         → _build_preset_divider 가 _divider_counter 기반 fallback 사용 (어제 의도 보존)
    """
    if not isinstance(slide_data, dict):
        return ""
    # 1순위 governing_main
    gm = str(slide_data.get("governing_main", "")).strip()
    if gm:
        m = _DIVIDER_ROMAN_RE.match(gm)
        if m:
            roman = _normalize_divider_prefix_to_roman(m.group(1))
            if roman:
                return roman
    # 2순위 section (라벨 떼고)
    section = str(slide_data.get("section", "")).strip()
    if section:
        title = section.replace("(챕터 divider)", "").replace("(chapter divider)", "").strip()
        if title:
            m2 = _DIVIDER_ROMAN_RE.match(title)
            if m2:
                roman = _normalize_divider_prefix_to_roman(m2.group(1))
                if roman:
                    return roman
    return ""


def _is_chapter_divider(slide_data):
    """챕터 간지 페이지 식별 — section "(챕터 divider)" 라벨 단독으로 식별.

    Spec D-Fix-DividerActuallyWorks (어제 작업이 한 번도 작동 안 한 회귀 수정):
      어제 spec(0a814ab) 의 1차 가드 `slide_type == "hero"` 가 항상 False 였다.
      이유: SLIDE pass output 스키마는 {section, shapes} 만 강제(L2768) →
            LLM 이 slide_type 키를 자기 JSON 출력에 안 박음 →
            final_slides 변환의 sr.meta(L3339, "shapes/section 제외 키만 보존")에
            slide_type 미포함 → generate_from_shape_json 의 slide_data 에서 부재.
      → slide_type 의존 완전 제거. section 라벨 "(챕터 divider)" 단독 식별.

    이 라벨은 OUTLINE_SYSTEM_PROMPT 출력 예시(L636) 에서 챕터 divider 에만 부여.
    콘셉트 슬로건("Ⅰ.4 콘셉트") / 강조 hero("Ⅰ.3 ...") 는 라벨 안 받음 →
    false positive 0 (어제 진단에서 확인).
    """
    if not isinstance(slide_data, dict):
        return False
    section = str(slide_data.get("section", ""))
    if "(챕터 divider)" in section or "(chapter divider)" in section.lower():
        return True
    return False


def _extract_divider_title(slide_data):
    """챕터 간지 부문명 추출 — governing_main 우선, 없으면 section 에서 추출.

    Spec D-Fix-DividerActuallyWorks:
      governing_main 도 LLM 이 SLIDE output JSON 에 안 박을 수 있음 →
      section 에서 "(챕터 divider)" 라벨을 떼서 부문명으로 사용 (확실한 fallback).

    우선순위:
      1. governing_main 이 "Ⅰ. 제안 개요" 형식이면 prefix 떼고 "제안 개요"
      2. governing_main 이 비-empty 면 그대로
      3. section 에서 "(챕터 divider)" 라벨 떼고 strip
      4. 모두 실패 시 "" (빈 문자열 — _build_preset_divider 가 빈 list 로 처리)
    """
    gm = str(slide_data.get("governing_main", "")).strip()
    if gm:
        m = _DIVIDER_ROMAN_RE.match(gm)
        if m:
            return m.group(2).strip()
        return gm
    # Spec D-Fix-DividerPrefixStrip — section fallback 에도 prefix 떼기 적용.
    #   기존: section 에서 "(챕터 divider)" 라벨만 떼고 통째 반환 → LLM 이 section 에 박은
    #         로마숫자 prefix("Ⅲ. 공간구성계획") 잔존 → 코드 자동 매김 prefix("Ⅰ") 와 합쳐져
    #         "Ⅰ Ⅲ. 공간구성계획" 사고. governing_main 경로와 동일한 정규식 적용해 일관 보장.
    section = str(slide_data.get("section", ""))
    title = section.replace("(챕터 divider)", "").replace("(chapter divider)", "").strip()
    m2 = _DIVIDER_ROMAN_RE.match(title)
    if m2:
        title = m2.group(2).strip()
    return title


def _build_preset_divider(slide_data, divider_idx=1):
    """챕터 간지 페이지 — 큰 로마숫자(코드 자동 매김) + 중간 부문명 한 줄(왼쪽 정렬).

    Spec D-Fix-DividerActuallyWorks (어제 spec 의 회귀 수정):
      어제 작업은 governing_main 의 LLM 매김 로마숫자(_DIVIDER_ROMAN_RE) 를 신뢰했으나,
      행사 도메인에서 LLM 이 로마숫자를 안 박거나 잘못 박음(V 중복 등) → 정규식 실패 →
      빈 list → LLM 백업 fallback → 들쭉날쭉. 이번엔 로마숫자를 코드가 등장 순서로
      자동 매김 (divider_idx) → LLM 자유 매김 환각 차단.

    입력:
      slide_data — chapter divider 식별된 페이지
      divider_idx — 1-based 등장 순번. generate_from_shape_json 의 카운터가 부여.

    부문명 추출 (_extract_divider_title):
      governing_main 우선 (prefix 떼기) → section 에서 라벨 떼기 → "" fallback.
    부문명 없으면 빈 list → 백업 fallback (마지막 안전망).
    ★ 부문명 길이에 따라 폰트 자동 축소 (한 줄 유지).
    ★ valign="bottom" 으로 두 박스의 baseline 정렬 (run 단위 size 다르게 못 하므로).
    """
    # Spec D-Fix-DividerNumberFromLLM — 간지 번호 결정 우선순위:
    #   1순위: LLM 이 section/governing_main 에 박은 prefix (목차·소항목과 같은 출처라 일치 보장)
    #          예: LLM 이 "Ⅲ. 공간구성계획" 박으면 간지에 Ⅲ → 목차 Ⅲ 과 정확히 일치.
    #          아라비아("3. ~") / ASCII 로마("III. ~") 도 유니코드 Ⅲ 으로 정규화.
    #   2순위 fallback: _divider_idx_to_roman(divider_idx) — _divider_counter 등장 순서 기반
    #          (어제 D-Fix-DividerActuallyWorks 의 "LLM 번호 누락/잘못 박음 시 코드 백업" 의도 보존)
    #   회귀 원인 (이번 spec 이 해결): 단순 카운터가 LLM 의도와 어긋남 (예: Ⅰ·Ⅱ 챕터 divider 가
    #          카운트 미진입 시 Ⅲ 챕터부터 카운터 1 → 간지에 Ⅰ 표시, 목차 Ⅲ 와 불일치).
    roman = _extract_divider_roman(slide_data) or _divider_idx_to_roman(divider_idx)
    title = _extract_divider_title(slide_data)
    if not title:
        return []
    # 부문명 길이별 폰트 자동 축소 — 한 줄 안 넘침 보수적 설정.
    # 한글 1자 폭 ≈ 폰트 크기와 비슷(1pt ≈ 1/72인치). 보수적으로 추정.
    n = len(title)
    if n <= 6:
        roman_size, title_size = 150, 54
    elif n <= 9:
        roman_size, title_size = 140, 48
    elif n <= 12:
        roman_size, title_size = 130, 42
    elif n <= 15:
        roman_size, title_size = 120, 36
    else:
        roman_size, title_size = 110, 30
    W, H = 11.69, 8.27
    # 세로 중앙쯤 — 큰 박스 h=2.5, y_box 가 박스 상단 (valign="bottom" 으로 baseline 정렬)
    y_box = 2.9
    box_h = 2.5
    x_left = 0.9
    # 로마숫자 박스 폭 — 한글 1자 폭 ≈ size/72 인치. 보수적으로 size*1.3/72 (안전 여백).
    # Ⅰ(좁음) / Ⅲ(넓음) 변동 흡수.
    roman_w = max(1.8, roman_size * 1.5 / 72.0)
    # 부문명 박스 — 로마숫자 우측 + 0.25 인치 여백
    title_x = x_left + roman_w + 0.25
    title_w = W - title_x - 0.5
    shapes = [
        {"type":"text", "x":x_left, "y":y_box, "w":roman_w, "h":box_h,
         "text":roman, "size":roman_size, "weight":900, "color":"#1A1A1A",
         "align":"left", "valign":"bottom"},
        {"type":"text", "x":title_x, "y":y_box, "w":title_w, "h":box_h,
         "text":title, "size":title_size, "weight":900, "color":"#1A1A1A",
         "align":"left", "valign":"bottom"},
    ]
    return shapes


# ─── Spec Preset-New-ConclusionCards — 3열 카드 + 하단 결론 밴드 ────────────────
# 상단 거버닝 + 중단 3열 카드(라벨바+헤드+옅은회색 desc박스 안 태그·설명) +
# 하단 검정 밴드(아래 화살표 + conclusion). 이미지 placeholder 없음.
# 색 정합 (거버닝 자동 보라 + 미매핑 검정/흰):
#   상단 title / 하단 conclusion → role:"governing" 마킹 → _get_theme(theme)["ACCENT"] 자동.
#   라벨 바·태그·밴드·아이템 흰글씨 → #1B1B1B / #FEFEFE 미매핑 (양 테마 그대로).
#   desc 박스 → #F5F5F5 (다크에선 DARK_MAP 매핑 #1F1F1F).
# 화살표 → chevron/block_arrow 우향 고정이라 사용 X, arrow(직선+삼각형 헤드) 세로 배치.
#
# 입력 스키마:
#   slide_data["items"]           = [{label(필수), head(필수), tag?, desc?}, ...]  ★ 정확히 3개
#   slide_data["conclusion"]      = str (필수)
#   slide_data["eyebrow"]?        = str
#   slide_data["subtitle"]?       = str
#   slide_data["lead"]?           = str
#   slide_data["title"]?          = str  (거버닝 — role:"governing" 자동 보라)
#   slide_data["conclusion_lead"]? = str
# 안전망: items list 아님 / 3개 정확히 아님 / 각 head 누락 / conclusion 없음 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트·SLIDE 프롬프트
#   미연결 (LLM 우연 선택 방지). 갤러리 검증 후 별도 spec 으로 켜기.
def _build_preset_conclusion_cards(slide_data):
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        items.append({
            "label": str(it.get("label", "")).strip(),
            "head": head,
            "tag": str(it.get("tag", "")).strip(),
            "desc": str(it.get("desc", "")).strip(),
        })
    # ★ Spec Preset-Loosen-Count — 정확 3 강제 폐지. 초과 잘라내고 최소 1개면 렌더.
    #   기존 성공 프리셋(timeline/circles/asymmetric)의 관대함 준수 (Law 1).
    if len(items) > 3:
        items = items[:3]
    if not items:
        return []
    conclusion = str(slide_data.get("conclusion", "")).strip()
    if not conclusion:
        return []

    eyebrow         = str(slide_data.get("eyebrow", "")).strip()
    subtitle        = str(slide_data.get("subtitle", "")).strip()
    lead            = str(slide_data.get("lead", "")).strip()
    title           = str(slide_data.get("title", "")).strip()
    conclusion_lead = str(slide_data.get("conclusion_lead", "")).strip()

    W, H = 11.69, 8.27
    shapes = []

    # ── ① 상단 거버닝 블록 (중앙정렬)
    top_x, top_w = 0.9, W - 1.8
    if eyebrow:
        shapes.append({"type":"text","x":top_x,"y":0.3,"w":top_w,"h":0.35,
                       "text":eyebrow,"size":11,"weight":400,"color":"#999999",
                       "align":"center","valign":"middle"})
    if subtitle:
        shapes.append({"type":"text","x":top_x,"y":0.7,"w":top_w,"h":0.35,
                       "text":subtitle,"size":12,"weight":400,"color":"#666666",
                       "align":"center","valign":"middle"})
    if lead:
        shapes.append({"type":"text","x":top_x,"y":1.15,"w":top_w,"h":0.4,
                       "text":lead,"size":14,"weight":500,"color":"#1A1A1A",
                       "align":"center","valign":"middle"})
    if title:
        # ★ role:"governing" → render 시 _get_theme(theme)["ACCENT"] 자동 (라이트 #6B46E5).
        shapes.append({"type":"text","x":top_x,"y":1.55,"w":top_w,"h":0.7,
                       "text":title,"size":28,"weight":800,"color":"#1A1A1A",
                       "align":"center","valign":"middle",
                       "role":"governing"})

    # ── ② 중단 N열 카드 (1~3, 균등 분배 — Spec Preset-Loosen-Count)
    #   n=1: card_w ≈ 9.89 (전폭)  n=2: card_w ≈ 4.795  n=3: card_w ≈ 3.097
    n = len(items)
    margin  = 0.9
    gap     = 0.3
    card_w  = (W - 2 * margin - (n - 1) * gap) / n
    card_top = 2.4
    label_h  = 0.5
    head_top = card_top + label_h + 0.05        # 2.95
    head_h   = 0.65
    desc_top = head_top + head_h + 0.1          # 3.7
    desc_h   = 2.15                             # ends 5.85
    for i, it in enumerate(items):
        cx = margin + i * (card_w + gap)
        # (1) 라벨 바 — 검정 fill + 흰 글씨 (마커 없음, 거버닝 아님)
        shapes.append({"type":"rect","x":cx,"y":card_top,"w":card_w,"h":label_h,
                       "fill":"#1B1B1B"})
        if it["label"]:
            shapes.append({"type":"text","x":cx,"y":card_top,"w":card_w,"h":label_h,
                           "text":it["label"],"size":12,"weight":700,"color":"#FEFEFE",
                           "align":"center","valign":"middle"})
        # (2) 헤드라인
        shapes.append({"type":"text","x":cx + 0.1,"y":head_top,"w":card_w - 0.2,"h":head_h,
                       "text":it["head"],"size":16,"weight":700,"color":"#1A1A1A",
                       "align":"center","valign":"middle"})
        # (3) 설명 박스 — 옅은 회색 fill (다크에선 #1F1F1F 자동)
        shapes.append({"type":"rect","x":cx,"y":desc_top,"w":card_w,"h":desc_h,
                       "fill":"#F5F5F5"})
        # (3a) 태그 — 검정 rounded rect (hero_cards L3340 패턴, radius 지정)
        tag_present = bool(it["tag"])
        if tag_present:
            tag_w = 1.4
            tag_h = 0.4
            tag_x = cx + (card_w - tag_w) / 2
            tag_y = desc_top + 0.18
            shapes.append({"type":"rect","x":tag_x,"y":tag_y,"w":tag_w,"h":tag_h,
                           "fill":"#1B1B1B","radius":0.06})
            shapes.append({"type":"text","x":tag_x,"y":tag_y,"w":tag_w,"h":tag_h,
                           "text":it["tag"],"size":11,"weight":700,"color":"#FEFEFE",
                           "align":"center","valign":"middle"})
        # (3b) desc
        if it["desc"]:
            desc_text_y = desc_top + (0.75 if tag_present else 0.2)
            desc_text_h = desc_top + desc_h - desc_text_y - 0.15
            shapes.append({"type":"text","x":cx + 0.2,"y":desc_text_y,
                           "w":card_w - 0.4,"h":desc_text_h,
                           "text":it["desc"],"size":11,"weight":400,"color":"#666666",
                           "align":"center","valign":"top"})

    # ── ③ 하단 결론 밴드 (전체 폭 검정)
    band_top = 6.0
    band_h   = H - band_top   # 2.27
    shapes.append({"type":"rect","x":0,"y":band_top,"w":W,"h":band_h,
                   "fill":"#1B1B1B"})
    # 아래 방향 화살표 — chevron/block_arrow 우향 고정이라 arrow(직선+헤드) 사용
    #   hero_cards L3308 패턴 정합 (검정 밴드 위 #FEFEFE)
    shapes.append({"type":"arrow","x1":W/2,"y1":band_top + 0.15,
                   "x2":W/2,"y2":band_top + 0.55,
                   "color":"#FEFEFE","width":1.5})
    if conclusion_lead:
        shapes.append({"type":"text","x":0.5,"y":band_top + 0.7,"w":W - 1.0,"h":0.35,
                       "text":conclusion_lead,"size":13,"weight":400,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
    # conclusion — 결론 문장. 검정 밴드 위라 흰 계열 색 지정 후 role:"governing" 마킹 시
    #   render 가 _get_theme(theme)["ACCENT"] 로 오버라이드 → 라이트 #6B46E5.
    #   ★ 검정 밴드 위 보라 대비 가독성은 갤러리 확인 (governing-purple 갤러리 검증됨).
    conclusion_y = band_top + (1.15 if conclusion_lead else 0.85)
    shapes.append({"type":"text","x":0.5,"y":conclusion_y,"w":W - 1.0,"h":1.0,
                   "text":conclusion,"size":26,"weight":900,"color":"#FEFEFE",
                   "align":"center","valign":"middle",
                   "role":"governing"})
    return shapes


# ─── Spec Preset-New-NumberedColumns — 넘버 카드 3~4열 + 결론 ───────────────────
# 상단 검정 헤더 블록(배지 알약 + 거버닝) + 중단 넘버 카드(초대형 옅은 회색 숫자 +
# 콘텐츠 스택) 3 또는 4열 + 하단 결론 2줄(밴드 없음).
# 색 정합:
#   상단 title / 하단 conclusion → role:"governing" → _get_theme(theme)["ACCENT"] 자동.
#   헤더 검정 블록 · 배지 stroke → #1B1B1B / #FEFEFE 미매핑 (양 테마 유지).
#   배지 stroke-only → fill 키 자체 생략 → dispatch None → _add_rect._set_no_fill.
#     (★ _add_rect 는 _add_circle 과 달리 "none" 정규화 가드 없어서 문자열 "none"
#      쓰면 _hex_to_rgb 실패. 키 생략이 안전.)
#   초대형 배경 숫자 → #EEEEEE (미매핑, 반투명 불가 대체).
#   아이템 head → #1A1A1A 검정 (Spec Preset-Purple-Accent-Cleanup — 강조색 제거, weight 800 위계).
# 화살표 없음 — conclusion_cards 와 달리 밴드 없이 흰 배경 위 결론 텍스트만.
#
# 입력 스키마:
#   slide_data["items"]           = [{label?, head(필수), desc?}, ...]  ★ 3개 또는 4개
#   slide_data["conclusion"]      = str (필수)
#   slide_data["badge"]?          = str (상단 알약 배지 문구)
#   slide_data["title"]?          = str (거버닝 — role:"governing" 자동 보라)
#   slide_data["conclusion_lead"]? = str
# 안전망: items 3/4 개 아님 / 각 head 누락 / conclusion 없음 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트·SLIDE
#   프롬프트 미연결 (LLM 우연 선택 방지). 갤러리 검증 후 별도 spec 으로 켜기.
def _build_preset_numbered_columns(slide_data):
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        items.append({
            "label": str(it.get("label", "")).strip(),
            "head": head,
            "desc": str(it.get("desc", "")).strip(),
        })
    # ★ Spec Preset-Loosen-Count — (3,4) 강제 폐지. 초과 잘라내고 최소 2개면 렌더 (1열은 무의미).
    if len(items) > 4:
        items = items[:4]
    if len(items) < 2:
        return []
    conclusion = str(slide_data.get("conclusion", "")).strip()
    if not conclusion:
        return []

    badge           = str(slide_data.get("badge", "")).strip()
    title           = str(slide_data.get("title", "")).strip()
    conclusion_lead = str(slide_data.get("conclusion_lead", "")).strip()

    W, H = 11.69, 8.27
    shapes = []

    # ── ① 상단 헤더 블록 (여백 둔 검정 박스)
    hdr_x, hdr_y = 0.9, 0.3
    hdr_w, hdr_h = W - 1.8, 2.0
    shapes.append({"type":"rect","x":hdr_x,"y":hdr_y,"w":hdr_w,"h":hdr_h,
                   "fill":"#1B1B1B"})
    # 배지 알약 (stroke-only, fill 키 생략)
    if badge:
        badge_w = 4.0
        badge_h = 0.42
        badge_x = hdr_x + (hdr_w - badge_w) / 2
        badge_y = hdr_y + 0.35
        shapes.append({"type":"rect","x":badge_x,"y":badge_y,"w":badge_w,"h":badge_h,
                       "stroke":"#FEFEFE","stroke_width":1,"radius":0.2})
        shapes.append({"type":"text","x":badge_x,"y":badge_y,"w":badge_w,"h":badge_h,
                       "text":badge,"size":11,"weight":600,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
    # 타이틀 (거버닝 — role:"governing" 자동 보라, 검정 배경 위)
    if title:
        title_y = hdr_y + (1.0 if badge else 0.7)
        title_h = hdr_h - (title_y - hdr_y) - 0.15
        shapes.append({"type":"text","x":hdr_x + 0.3,"y":title_y,
                       "w":hdr_w - 0.6,"h":title_h,
                       "text":title,"size":26,"weight":800,"color":"#FEFEFE",
                       "align":"center","valign":"middle",
                       "role":"governing"})

    # ── ② 중단 넘버 카드 (n=2/3/4 가변 — Spec Preset-Loosen-Count)
    n = len(items)
    margin = 0.9
    if n == 2:
        gap = 0.4
    elif n == 3:
        gap = 0.3
    else:  # n == 4
        gap = 0.25
    card_w = (W - 2 * margin - (n - 1) * gap) / n
    card_top = 2.6
    card_h   = 3.5

    # n 별 파라미터 — 열이 적을수록 숫자/헤드 크게, 많을수록 축소
    if n == 2:
        num_size = 150   # 카드 폭 넓으니 배경 숫자도 크게
        content_x_off = 1.2
        head_size = 20
    elif n == 3:
        num_size = 120
        content_x_off = 0.9
        head_size = 17
    else:  # n == 4
        num_size = 92
        content_x_off = 0.7
        head_size = 15

    for i, it in enumerate(items):
        cx = margin + i * (card_w + gap)
        # 초대형 배경 숫자 — 카드 좌측, 옅은 회색 (반투명 대체)
        shapes.append({"type":"text","x":cx,"y":card_top - 0.15,
                       "w":card_w * 0.55,"h":2.4,
                       "text":str(i + 1),"size":num_size,"weight":900,
                       "color":"#EEEEEE","align":"left","valign":"top"})
        # 콘텐츠 스택 (숫자 우측 — 겹침 허용)
        cx_c = cx + content_x_off
        w_c  = cx + card_w - cx_c - 0.05
        if w_c < 1.0:
            w_c = 1.0
        y_c = card_top + 0.4
        if it["label"]:
            shapes.append({"type":"text","x":cx_c,"y":y_c,"w":w_c,"h":0.35,
                           "text":it["label"],"size":11,"weight":600,
                           "color":"#666666","align":"left","valign":"top"})
            y_c += 0.4
        # head — Spec Preset-Purple-Accent-Cleanup — 아이템 강조색 제거.
        #   거버닝 = 브랜드 보라 하나 원칙 → 아이템은 검정 (거버닝 아니므로 role 마킹 X).
        #   위계는 weight 800 + head_size 로 유지.
        shapes.append({"type":"text","x":cx_c,"y":y_c,"w":w_c,"h":0.7,
                       "text":it["head"],"size":head_size,"weight":800,
                       "color":"#1A1A1A","align":"left","valign":"top"})
        y_c += 0.8
        if it["desc"]:
            desc_h = card_top + card_h - y_c - 0.15
            if desc_h < 0.4:
                desc_h = 0.4
            shapes.append({"type":"text","x":cx_c,"y":y_c,"w":w_c,"h":desc_h,
                           "text":it["desc"],"size":11,"weight":400,
                           "color":"#666666","align":"left","valign":"top"})

    # ── ③ 하단 결론 2줄 (흰 배경, 밴드 없음)
    bot_x, bot_w = 0.9, W - 1.8
    bot_top = card_top + card_h + 0.3   # 6.4
    if conclusion_lead:
        shapes.append({"type":"text","x":bot_x,"y":bot_top,"w":bot_w,"h":0.35,
                       "text":conclusion_lead,"size":13,"weight":400,
                       "color":"#666666","align":"center","valign":"middle"})
        conclusion_y = bot_top + 0.45
    else:
        conclusion_y = bot_top + 0.15
    conclusion_h = H - conclusion_y - 0.3
    if conclusion_h < 0.6:
        conclusion_h = 0.6
    # conclusion 은 거버닝 — role:"governing" → 자동 보라 (흰 배경 위 #6B46E5)
    shapes.append({"type":"text","x":bot_x,"y":conclusion_y,"w":bot_w,"h":conclusion_h,
                   "text":conclusion,"size":24,"weight":900,"color":"#1A1A1A",
                   "align":"center","valign":"middle",
                   "role":"governing"})
    return shapes


# ─── Spec Preset-New-HeroDetail — 헤더 + 좌우분할(대형 이미지 + 번호항목 3) + 결론 ─
# 상단 검정 헤더 (배지+대형숫자+텍스트 스택) + 중단 좌 대형 이미지 placeholder +
# 우 번호 항목 3개 (미니 이미지 + 원문자 head + desc) + 하단 결론 2줄.
# 색 정합:
#   상단 title / 하단 conclusion → role:"governing" → _get_theme(theme)["ACCENT"] 자동.
#   헤더 검정 블록 / 아이템 원문자 head → #1B1B1B / #1A1A1A / #FEFEFE 하드코딩.
#     (Spec Preset-Purple-Accent-Cleanup — 아이템 head 강조색 제거, weight 800 위계 유지)
#   이미지 placeholder → #F5F5F5 (DARK_MAP fill 매핑 → 다크 #1F1F1F 자동).
# 원문자 번호 → 3개 고정 → ①②③.
# 아이템 head 는 아이템 성격이라 role:"governing" 마킹 X (원문자 포함 전체 검정).
#
# 입력 스키마:
#   slide_data["items"]           = [{head(필수), desc?}, ...]  ★ 정확히 3개
#   slide_data["conclusion"]      = str (필수)
#   slide_data["badge"]?          = str (헤더 배지, 줄바꿈 가능)
#   slide_data["number"]?         = str (헤더 초대형 숫자)
#   slide_data["index"]?          = str (예: "1-1")
#   slide_data["subtitle"]?       = str
#   slide_data["title"]?          = str  (거버닝 — 자동 보라)
#   slide_data["section_title"]?  = str
#   slide_data["section_sub"]?    = str
#   slide_data["conclusion_lead"]? = str
# 안전망: items 정확히 3개 아님 / 각 head 누락 / conclusion 없음 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
_HERO_DETAIL_CIRCLED = ["①", "②", "③"]


def _build_preset_hero_detail(slide_data):
    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        items.append({
            "head": head,
            "desc": str(it.get("desc", "")).strip(),
        })
    # ★ Spec Preset-Loosen-Count — 정확 3 강제 폐지. 최대 3개까지 잘라내고 최소 1개면 렌더.
    #   원문자 배열 _HERO_DETAIL_CIRCLED 가 3개라 그 이상은 인덱스 오버 방지 겸 자름.
    if len(items) > 3:
        items = items[:3]
    if not items:
        return []
    conclusion = str(slide_data.get("conclusion", "")).strip()
    if not conclusion:
        return []

    badge           = str(slide_data.get("badge", "")).strip()
    number          = str(slide_data.get("number", "")).strip()
    index           = str(slide_data.get("index", "")).strip()
    subtitle        = str(slide_data.get("subtitle", "")).strip()
    title           = str(slide_data.get("title", "")).strip()
    section_title   = str(slide_data.get("section_title", "")).strip()
    section_sub     = str(slide_data.get("section_sub", "")).strip()
    conclusion_lead = str(slide_data.get("conclusion_lead", "")).strip()

    W, H = 11.69, 8.27
    shapes = []

    # ── ① 상단 헤더 (검정 블록)
    hdr_x, hdr_y = 0.9, 0.3
    hdr_w, hdr_h = W - 1.8, 1.8
    shapes.append({"type":"rect","x":hdr_x,"y":hdr_y,"w":hdr_w,"h":hdr_h,
                   "fill":"#1B1B1B"})
    # 가로 순차 배치: 배지 → 숫자 → 텍스트 스택
    cursor_x = hdr_x + 0.3
    # (1) 배지 (선택) — Spec Preset-Purple-Accent-Cleanup — 브랜드 보라 fill 제거.
    #   거버닝 아닌 UI 요소는 강조색 없음. 검정 헤더 위 검정 배지는 시각 위계가 약해지므로
    #   다른 신규 프리셋(numbered_columns)의 stroke-only 알약과 동일 패턴 적용 (fill 생략 + 흰 stroke).
    if badge:
        b_w = 1.4
        b_h = 1.3
        b_y = hdr_y + (hdr_h - b_h) / 2
        shapes.append({"type":"rect","x":cursor_x,"y":b_y,"w":b_w,"h":b_h,
                       "stroke":"#FEFEFE","stroke_width":1,"radius":0.05})
        shapes.append({"type":"text","x":cursor_x,"y":b_y,"w":b_w,"h":b_h,
                       "text":badge,"size":13,"weight":700,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
        cursor_x += b_w + 0.2
    # (2) 초대형 숫자 (선택) — 72pt weight 900 Paperlogy 는 glyph 폭 ≈ 0.55~0.6"/자리.
    #   기존 n_w=1.3 은 한 자리엔 충분하나 두 자리("10" 등) 는 부족 → word_wrap 세로 흐름.
    #   n_w=2.0 = 두 자리 glyph 폭(≈1.2") + kerning/여유 마진 (한 자리엔 좌측 정렬이라 우측 공백만 늘어 무해).
    #   ★ size(pt) 는 유지 — 임팩트 유지가 이번 수정의 전제.
    if number:
        n_w = 2.0
        shapes.append({"type":"text","x":cursor_x,"y":hdr_y,"w":n_w,"h":hdr_h,
                       "text":number,"size":72,"weight":900,"color":"#FEFEFE",
                       "align":"left","valign":"middle"})
        cursor_x += n_w + 0.15
    # (3) 텍스트 스택 (index / subtitle / title) — 남은 폭 채움, 좌측 정렬
    stack_x = cursor_x
    stack_w = (hdr_x + hdr_w) - stack_x - 0.3
    if stack_w >= 1.0 and (index or subtitle or title):
        ty = hdr_y + 0.25
        if index:
            shapes.append({"type":"text","x":stack_x,"y":ty,"w":stack_w,"h":0.3,
                           "text":index,"size":11,"weight":500,"color":"#AAAAAA",
                           "align":"left","valign":"middle"})
            ty += 0.35
        if subtitle:
            shapes.append({"type":"text","x":stack_x,"y":ty,"w":stack_w,"h":0.3,
                           "text":subtitle,"size":12,"weight":400,"color":"#DDDDDD",
                           "align":"left","valign":"middle"})
            ty += 0.35
        if title:
            th = hdr_y + hdr_h - ty - 0.15
            if th < 0.5:
                th = 0.5
            shapes.append({"type":"text","x":stack_x,"y":ty,"w":stack_w,"h":th,
                           "text":title,"size":22,"weight":800,"color":"#FEFEFE",
                           "align":"left","valign":"middle",
                           "role":"governing"})

    # ── ② 중단 좌우 분할 (좌 대형 이미지 + 우 번호 항목 3)
    mid_top = hdr_y + hdr_h + 0.3       # 2.4
    mid_h   = 3.9                       # ends 6.3
    # 좌 대형 이미지 placeholder (hsplit L2947 / timeline L2787 패턴)
    L_x, L_w = 0.9, 6.0
    shapes.append({"type":"rect","x":L_x,"y":mid_top,"w":L_w,"h":mid_h,
                   "fill":"#F5F5F5","stroke":"#DDDDDD","stroke_width":1})
    shapes.append({"type":"text","x":L_x,"y":mid_top + mid_h/2 - 0.2,
                   "w":L_w,"h":0.4,
                   "text":"이미지 영역","size":13,"weight":400,"color":"#AAAAAA",
                   "align":"center","valign":"middle"})
    # 우 콘텐츠 컬럼
    R_x = L_x + L_w + 0.2               # 7.1
    R_w = (W - 0.9) - R_x               # 3.69
    ry = mid_top
    if section_title:
        shapes.append({"type":"text","x":R_x,"y":ry,"w":R_w,"h":0.5,
                       "text":section_title,"size":15,"weight":800,"color":"#1A1A1A",
                       "align":"center","valign":"middle"})
        ry += 0.55
    if section_sub:
        shapes.append({"type":"text","x":R_x,"y":ry,"w":R_w,"h":0.3,
                       "text":section_sub,"size":11,"weight":400,"color":"#666666",
                       "align":"center","valign":"middle"})
        ry += 0.35
    # 번호 항목 N개 (미니 이미지 + 원문자 head + desc) — Spec Preset-Loosen-Count 로 1~3 가변
    items_top = ry + 0.1
    items_bot = mid_top + mid_h
    item_h    = (items_bot - items_top) / max(1, len(items))
    mi_w = 0.9
    for i, it in enumerate(items):
        iy = items_top + i * item_h
        # 미니 이미지 (좌)
        mi_h = min(item_h - 0.1, 0.9)
        shapes.append({"type":"rect","x":R_x,"y":iy,"w":mi_w,"h":mi_h,
                       "fill":"#F5F5F5","stroke":"#DDDDDD","stroke_width":1})
        shapes.append({"type":"text","x":R_x,"y":iy,"w":mi_w,"h":mi_h,
                       "text":"이미지","size":10,"weight":400,"color":"#AAAAAA",
                       "align":"center","valign":"middle"})
        # 콘텐츠 (우) — 원문자 + head + desc
        cx = R_x + mi_w + 0.15
        cw = R_x + R_w - cx
        if cw < 0.8:
            cw = 0.8
        num_head = _HERO_DETAIL_CIRCLED[i] + " " + it["head"]
        # Spec Preset-Purple-Accent-Cleanup — 아이템 강조색 제거 (원문자 포함 전체 검정).
        #   거버닝 아니므로 role 마킹 X. 위계는 weight 800 로 유지.
        shapes.append({"type":"text","x":cx,"y":iy,"w":cw,"h":0.4,
                       "text":num_head,"size":12,"weight":800,"color":"#1A1A1A",
                       "align":"left","valign":"top"})
        if it["desc"]:
            d_top = iy + 0.42
            d_h   = iy + item_h - d_top - 0.05
            if d_h < 0.3:
                d_h = 0.3
            shapes.append({"type":"text","x":cx,"y":d_top,"w":cw,"h":d_h,
                           "text":it["desc"],"size":9,"weight":400,"color":"#666666",
                           "align":"left","valign":"top"})

    # ── ③ 하단 결론 (흰 배경, 밴드 없음)
    bot_x, bot_w = 0.9, W - 1.8
    bot_top = mid_top + mid_h + 0.2     # 6.5
    if conclusion_lead:
        shapes.append({"type":"text","x":bot_x,"y":bot_top,"w":bot_w,"h":0.35,
                       "text":conclusion_lead,"size":12,"weight":400,
                       "color":"#666666","align":"center","valign":"middle"})
        conclusion_y = bot_top + 0.45
    else:
        conclusion_y = bot_top + 0.1
    conclusion_h = H - conclusion_y - 0.3
    if conclusion_h < 0.6:
        conclusion_h = 0.6
    shapes.append({"type":"text","x":bot_x,"y":conclusion_y,"w":bot_w,"h":conclusion_h,
                   "text":conclusion,"size":22,"weight":900,"color":"#1A1A1A",
                   "align":"center","valign":"middle",
                   "role":"governing"})
    return shapes


# ─── Spec Preset-New-FlowDetail — 상단 가로 흐름 + 하단 항목별 상세 (고밀도) ────
# 상단 거버닝(subtitle + title) + 중단 가로 흐름 박스 4~5개 (en/kr/desc, 화살표
# 연결) + 하단 3~4열 상세(대형 key 보라 + key_sub + lead + items 3개 head/desc).
# 색 정합:
#   상단 title → role:"governing" → _get_theme(theme)["ACCENT"] 자동 (라이트 #6B46E5).
#   중단 en / 하단 key → #666666 / #1A1A1A (Spec Preset-Purple-Accent-Cleanup — 강조색 제거).
#     en 은 작은 영문 라벨이라 회색 #666666, key 는 대형 키워드라 검정 #1A1A1A.
#   박스 fill/stroke → #FFFFFF / #DDDDDD (DARK_MAP 매핑으로 자연 반전).
#   화살표 → chevron/block_arrow 대신 arrow(직선+헤드) 가로 배치.
# 밀도 주의: 4열 + 각 3항목 케이스가 가장 빡빡 — desc size 최소 9pt 유지.
#
# 입력 스키마:
#   slide_data["title"]            = str (필수, 거버닝)
#   slide_data["steps"]            = [{kr(필수), en?, desc?}, ...]  ★ 4개 또는 5개
#   slide_data["columns"]          = [{key(필수), items(필수, 정확 3), key_sub?, lead?}, ...]
#                                    ★ 3개 또는 4개, 각 items = [{head(필수), desc?}, ...] 정확 3
#   slide_data["subtitle"]?        = str
#   slide_data["section_title_1"]? = str  (중단 헤더)
#   slide_data["section_title_2"]? = str  (하단 헤더)
#   slide_data["flow_lead"]?       = str
# 안전망: title 없음 / steps 4~5 아님 / columns 3~4 아님 / 각 column items 3 아님 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
def _build_preset_flow_detail(slide_data):
    title = str(slide_data.get("title", "")).strip()
    if not title:
        return []

    # steps 정규화 (4~5개, 각 kr 필수)
    steps_raw = slide_data.get("steps") or []
    if not isinstance(steps_raw, list):
        return []
    steps = []
    for s in steps_raw:
        if not isinstance(s, dict):
            continue
        kr = str(s.get("kr", "")).strip()
        if not kr:
            continue
        steps.append({
            "en":   str(s.get("en", "")).strip(),
            "kr":   kr,
            "desc": str(s.get("desc", "")).strip(),
        })
    # ★ Spec Preset-Loosen-Count — 정확 (4,5) 강제 폐지. 최대 5개 잘라내고 최소 3개면 렌더.
    if len(steps) > 5:
        steps = steps[:5]
    if len(steps) < 3:
        return []

    # columns 정규화 (2~4개, 각 key 필수 + items 0~3개 허용 — Spec Preset-Loosen-Count)
    columns_raw = slide_data.get("columns") or []
    if not isinstance(columns_raw, list):
        return []
    columns = []
    for col in columns_raw:
        if not isinstance(col, dict):
            continue
        key = str(col.get("key", "")).strip()
        if not key:
            continue
        items_raw = col.get("items") or []
        if not isinstance(items_raw, list):
            items_raw = []
        col_items = []
        for it in items_raw:
            if not isinstance(it, dict):
                continue
            head = str(it.get("head", "")).strip()
            if not head:
                continue
            col_items.append({
                "head": head,
                "desc": str(it.get("desc", "")).strip(),
            })
        # ★ items 정확 3 강제 폐지 — 초과 잘라내고 0~3 모두 허용 (key/lead 만 있어도 유효 column).
        if len(col_items) > 3:
            col_items = col_items[:3]
        columns.append({
            "key":     key,
            "key_sub": str(col.get("key_sub", "")).strip(),
            "lead":    str(col.get("lead", "")).strip(),
            "items":   col_items,
        })
    if len(columns) > 4:
        columns = columns[:4]
    if len(columns) < 2:
        return []

    subtitle        = str(slide_data.get("subtitle", "")).strip()
    section_title_1 = str(slide_data.get("section_title_1", "")).strip()
    section_title_2 = str(slide_data.get("section_title_2", "")).strip()
    flow_lead       = str(slide_data.get("flow_lead", "")).strip()

    W, H = 11.69, 8.27
    margin = 0.9
    inner_w = W - 2 * margin
    shapes = []

    # ── ① 상단 거버닝 (중앙정렬)
    if subtitle:
        shapes.append({"type":"text","x":margin,"y":0.35,"w":inner_w,"h":0.25,
                       "text":subtitle,"size":11,"weight":400,"color":"#666666",
                       "align":"center","valign":"middle"})
    # title (필수, role:"governing" 자동 보라)
    shapes.append({"type":"text","x":margin,"y":0.65,"w":inner_w,"h":0.85,
                   "text":title,"size":26,"weight":800,"color":"#1A1A1A",
                   "align":"center","valign":"middle",
                   "role":"governing"})

    # ── ② 중단 가로 흐름 (4 or 5 boxes + arrows)
    if section_title_1:
        shapes.append({"type":"text","x":margin,"y":1.7,"w":inner_w,"h":0.3,
                       "text":section_title_1,"size":13,"weight":700,"color":"#1A1A1A",
                       "align":"left","valign":"middle"})
    if flow_lead:
        shapes.append({"type":"text","x":margin,"y":2.0,"w":inner_w,"h":0.25,
                       "text":flow_lead,"size":10,"weight":400,"color":"#666666",
                       "align":"left","valign":"middle"})

    n_steps = len(steps)
    arrow_w = 0.4
    box_w   = (inner_w - (n_steps - 1) * arrow_w) / n_steps
    box_top = 2.35
    box_h   = 1.55
    # ★ Spec Preset-Loosen-Count — n_steps=3 파라미터 추가 (박스 폭 넓으니 폰트 확대)
    if n_steps == 3:
        kr_size, en_size, desc_size_s = 18, 11, 11
    elif n_steps == 4:
        kr_size, en_size, desc_size_s = 16, 10, 10
    else:  # 5
        kr_size, en_size, desc_size_s = 14, 9, 9

    for i, st in enumerate(steps):
        bx = margin + i * (box_w + arrow_w)
        # 박스 (rounded rect)
        shapes.append({"type":"rect","x":bx,"y":box_top,"w":box_w,"h":box_h,
                       "fill":"#FFFFFF","stroke":"#DDDDDD","stroke_width":1,
                       "radius":0.05})
        pad = 0.12
        iy = box_top + pad
        # en (선택) — Spec Preset-Purple-Accent-Cleanup — 강조색 제거.
        #   작은 영문 라벨이므로 검정보다 회색 #666666 이 자연스러움 (kr 검정 위계 유지).
        if st["en"]:
            shapes.append({"type":"text","x":bx + pad,"y":iy,"w":box_w - 2 * pad,"h":0.3,
                           "text":st["en"],"size":en_size,"weight":500,"color":"#666666",
                           "align":"center","valign":"top"})
            iy += 0.32
        # kr (필수)
        shapes.append({"type":"text","x":bx + pad,"y":iy,"w":box_w - 2 * pad,"h":0.5,
                       "text":st["kr"],"size":kr_size,"weight":800,"color":"#1A1A1A",
                       "align":"center","valign":"top"})
        iy += 0.55
        # desc (선택)
        if st["desc"]:
            d_h = box_top + box_h - iy - pad
            if d_h < 0.2:
                d_h = 0.2
            shapes.append({"type":"text","x":bx + pad,"y":iy,"w":box_w - 2 * pad,"h":d_h,
                           "text":st["desc"],"size":desc_size_s,"weight":400,
                           "color":"#666666","align":"center","valign":"top"})
        # 화살표 (마지막 제외)
        if i < n_steps - 1:
            ax1 = bx + box_w + 0.05
            ax2 = bx + box_w + arrow_w - 0.05
            ay  = box_top + box_h / 2
            shapes.append({"type":"arrow","x1":ax1,"y1":ay,"x2":ax2,"y2":ay,
                           "color":"#1A1A1A","width":1.5})

    # ── ③ 하단 항목별 상세 (3 or 4 columns + 세로 구분선)
    if section_title_2:
        shapes.append({"type":"text","x":margin,"y":4.15,"w":inner_w,"h":0.3,
                       "text":section_title_2,"size":13,"weight":700,"color":"#1A1A1A",
                       "align":"left","valign":"middle"})

    n_cols       = len(columns)
    col_area_top = 4.55
    col_area_bot = H - 0.25          # 8.02
    col_w        = inner_w / n_cols
    # ★ Spec Preset-Loosen-Count — n_cols=2 파라미터 추가 (열 폭 넓으니 key 대형화)
    if n_cols == 2:
        key_size, head_size_c, desc_size_c = 36, 14, 11
    elif n_cols == 3:
        key_size, head_size_c, desc_size_c = 30, 12, 10
    else:  # 4
        key_size, head_size_c, desc_size_c = 26, 11, 9

    # 세로 구분선 (열 사이 n-1개)
    for i in range(1, n_cols):
        dx = margin + i * col_w
        shapes.append({"type":"line","x1":dx,"y1":col_area_top,
                       "x2":dx,"y2":col_area_bot,
                       "color":"#DDDDDD","width":1})

    inner_pad = 0.15
    for i, col in enumerate(columns):
        cx = margin + i * col_w + inner_pad
        cw = col_w - 2 * inner_pad
        cy = col_area_top
        # key (대형 검정 — Spec Preset-Purple-Accent-Cleanup, 아이템 강조색 제거)
        #   아이템 성격이라 role 마킹 X. 위계는 weight 900 + key_size 대형 폰트로 유지.
        shapes.append({"type":"text","x":cx,"y":cy,"w":cw,"h":0.65,
                       "text":col["key"],"size":key_size,"weight":900,"color":"#1A1A1A",
                       "align":"center","valign":"top"})
        cy += 0.7
        if col["key_sub"]:
            shapes.append({"type":"text","x":cx,"y":cy,"w":cw,"h":0.25,
                           "text":col["key_sub"],"size":11,"weight":400,"color":"#666666",
                           "align":"center","valign":"top"})
            cy += 0.3
        if col["lead"]:
            shapes.append({"type":"text","x":cx,"y":cy,"w":cw,"h":0.45,
                           "text":col["lead"],"size":11,"weight":400,"color":"#666666",
                           "align":"center","valign":"top"})
            cy += 0.5
        # items (0~3개 세로 스택 — Spec Preset-Loosen-Count 로 개수 가변)
        n_items_col = len(col["items"])
        if n_items_col == 0:
            continue    # 이 column 은 key/key_sub/lead 만 표시 (items 없어도 유효)
        items_top = cy + 0.1
        items_area = col_area_bot - items_top
        if items_area < 0.6:
            items_area = 0.6
        item_h = items_area / n_items_col
        for j, it in enumerate(col["items"]):
            iy = items_top + j * item_h
            # head (앞에 "· " bullet inline — 별도 shape 없음)
            head_text = "· " + it["head"]
            shapes.append({"type":"text","x":cx,"y":iy,"w":cw,"h":0.3,
                           "text":head_text,"size":head_size_c,"weight":700,
                           "color":"#1A1A1A","align":"left","valign":"top"})
            if it["desc"]:
                d_top = iy + 0.3
                d_h   = item_h - 0.32
                if d_h < 0.2:
                    d_h = 0.2
                shapes.append({"type":"text","x":cx + 0.15,"y":d_top,
                               "w":cw - 0.15,"h":d_h,
                               "text":it["desc"],"size":desc_size_c,"weight":400,
                               "color":"#666666","align":"left","valign":"top"})
    return shapes


# ─── Spec Preset-New-QuadDetail — 좌측정렬 거버닝 + 2×2 상세 블록 ──────────────
# 좌측정렬 거버닝(배지+subtitle+title) + 2×2 그리드 (각 블록: 검정 헤더바 +
# 체크리스트(선택 2~3) + 이미지 placeholder). 기존 신규 프리셋 4종이 모두
# 중앙정렬 거버닝이라 좌측정렬 다양성 확보 목적.
# 색 정합:
#   상단 title → role:"governing" → _get_theme(theme)["ACCENT"] 자동 (라이트 #6B46E5).
#   블록 헤더바 → #1B1B1B 미매핑 (양 테마 검정 유지).
#   배지 → #1B1B1B + rounded (알약형, hero_cards L3340 패턴).
#   이미지 placeholder → #F5F5F5 (DARK_MAP fill 매핑 → 다크 #1F1F1F 자동).
#   체크 항목 앞에 "✓ " 문자 inline (별도 shape 없음).
# 밀도: 2×2 블록이라 각 블록이 W/2 폭 확보 → 체크리스트 + 이미지 각 여유.
# ★ points 없는 블록: 체크리스트 슬롯 생략 → 이미지가 블록 전체(헤더 제외) 채움.
# ★ points 있는 블록: 고정 슬롯 0.85" 확보 (2·3개 시각 정렬 유지).
#
# 입력 스키마:
#   slide_data["title"]     = str (필수, 거버닝)
#   slide_data["items"]     = [{head(필수), points?}, ...]  ★ 정확히 4개
#                             points = list[str] (2~3 권장, 상한 3)
#   slide_data["badge"]?    = str  (좌측 알약 배지)
#   slide_data["subtitle"]? = str
# 안전망: title 없음 / items 정확 4 아님 / 각 head 누락 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
def _build_preset_quad_detail(slide_data):
    title = str(slide_data.get("title", "")).strip()
    if not title:
        return []

    items_raw = slide_data.get("items") or []
    if not isinstance(items_raw, list):
        return []
    items = []
    for it in items_raw:
        if not isinstance(it, dict):
            continue
        head = str(it.get("head", "")).strip()
        if not head:
            continue
        pts_raw = it.get("points") or []
        if not isinstance(pts_raw, list):
            pts_raw = []
        points = [str(p).strip() for p in pts_raw if str(p).strip()][:3]
        items.append({"head": head, "points": points})
    # ★ Spec Preset-Loosen-Count — 정확 4 강제 폐지. 최대 4개 잘라내고 최소 2개면 렌더.
    #   n=2 → 1행 2블록(전 높이), n=3 → 1행 3블록(전 높이), n=4 → 2x2(기존).
    if len(items) > 4:
        items = items[:4]
    if len(items) < 2:
        return []

    badge    = str(slide_data.get("badge", "")).strip()
    subtitle = str(slide_data.get("subtitle", "")).strip()

    W, H = 11.69, 8.27
    margin = 0.9
    shapes = []

    # ── ① 상단 거버닝 (★ 좌측 정렬)
    top_y = 0.35
    # 배지 (선택) — 검정 알약 rect + 흰 텍스트
    if badge:
        b_w, b_h = 1.6, 0.35
        shapes.append({"type":"rect","x":margin,"y":top_y,"w":b_w,"h":b_h,
                       "fill":"#1B1B1B","radius":0.05})
        shapes.append({"type":"text","x":margin,"y":top_y,"w":b_w,"h":b_h,
                       "text":badge,"size":11,"weight":700,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
    # subtitle (선택) — 좌측 정렬 회색
    sub_y = top_y + (0.5 if badge else 0.05)
    if subtitle:
        shapes.append({"type":"text","x":margin,"y":sub_y,"w":W - 2*margin,"h":0.3,
                       "text":subtitle,"size":12,"weight":400,"color":"#666666",
                       "align":"left","valign":"middle"})
        title_y = sub_y + 0.35
    else:
        title_y = sub_y
    # title (필수, 거버닝 — 좌측 정렬, role:"governing" 자동 보라)
    shapes.append({"type":"text","x":margin,"y":title_y,"w":W - 2*margin,"h":0.75,
                   "text":title,"size":26,"weight":800,"color":"#1A1A1A",
                   "align":"left","valign":"middle",
                   "role":"governing"})
    top_end = title_y + 0.85

    # ── ② 그리드 (n=2 1행, n=3 1행, n=4 2×2 — Spec Preset-Loosen-Count)
    n = len(items)
    grid_top = top_end + 0.15         # 최소 여유
    if grid_top < 2.1:
        grid_top = 2.1
    grid_bot = H - 0.25
    grid_h   = grid_bot - grid_top     # 세로 총 폭
    gap = 0.3
    inner_w = W - 2 * margin
    header_h = 0.5

    if n == 2:
        # 1행 2블록 (전 높이) — 각 블록 폭 넓게, 이미지 자리도 크게
        block_w = (inner_w - gap) / 2
        block_h = grid_h
        positions = [
            (margin, grid_top),
            (margin + block_w + gap, grid_top),
        ]
    elif n == 3:
        # 1행 3블록 (전 높이) — triad 처럼 3열 가로 나열
        block_w = (inner_w - 2 * gap) / 3
        block_h = grid_h
        positions = [
            (margin, grid_top),
            (margin + block_w + gap, grid_top),
            (margin + 2 * (block_w + gap), grid_top),
        ]
    else:  # n == 4
        block_w = (inner_w - gap) / 2
        block_h = (grid_h - gap) / 2
        # 블록 좌표 (좌상 · 우상 · 좌하 · 우하)
        positions = [
            (margin,              grid_top),                # 좌상
            (margin + block_w + gap, grid_top),             # 우상
            (margin,              grid_top + block_h + gap),# 좌하
            (margin + block_w + gap, grid_top + block_h + gap),  # 우하
        ]

    for i, it in enumerate(items):
        bx, by = positions[i]
        # (1) 헤더 바 (검정) — 마커 없음 (거버닝 아님)
        shapes.append({"type":"rect","x":bx,"y":by,"w":block_w,"h":header_h,
                       "fill":"#1B1B1B"})
        shapes.append({"type":"text","x":bx,"y":by,"w":block_w,"h":header_h,
                       "text":it["head"],"size":13,"weight":700,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
        # 내부 콘텐츠 영역 (헤더 아래)
        content_top = by + header_h + 0.1
        content_bot = by + block_h - 0.05
        content_h   = content_bot - content_top
        # (2) 체크리스트 — points 있을 때만 고정 슬롯 0.85"
        if it["points"]:
            check_slot_h = 0.85
            for j, pt in enumerate(it["points"]):
                p_y = content_top + j * 0.28
                shapes.append({"type":"text","x":bx + 0.15,"y":p_y,
                               "w":block_w - 0.3,"h":0.28,
                               "text":"✓ " + pt,"size":11,"weight":400,"color":"#1A1A1A",
                               "align":"left","valign":"middle"})
            img_top = content_top + check_slot_h + 0.05
            img_h   = content_bot - img_top
        else:
            # 체크리스트 생략 — 이미지가 콘텐츠 영역 전체 차지
            img_top = content_top
            img_h   = content_h
        if img_h < 0.4:
            img_h = 0.4
        # (3) 이미지 placeholder (hsplit L2947 / timeline L2787 패턴)
        shapes.append({"type":"rect","x":bx + 0.1,"y":img_top,
                       "w":block_w - 0.2,"h":img_h,
                       "fill":"#F5F5F5","stroke":"#DDDDDD","stroke_width":1})
        shapes.append({"type":"text","x":bx + 0.1,"y":img_top + img_h/2 - 0.2,
                       "w":block_w - 0.2,"h":0.4,
                       "text":"이미지 영역","size":11,"weight":400,"color":"#AAAAAA",
                       "align":"center","valign":"middle"})
    return shapes


# ─── Spec Preset-New-FullbleedOverlay — 풀블리드 배경 + 텍스트 오버레이 ─────────
# 화면 전체 이미지 placeholder 위에 상단 거버닝 + 좌하단 텍스트 오버레이.
# 조감도·시설사진·동선계획·컨셉·3D렌더링 등 "비주얼이 주인공" 슬라이드.
# 색 정합:
#   상단 title → role:"governing" → _get_theme(theme)["ACCENT"] 자동 (라이트 #6B46E5).
#   배경 → #F5F5F5 (DARK_MAP fill 매핑 → 다크 #1F1F1F 자동).
#   오버레이 박스·배지 → #1B1B1B 미매핑 (양 테마 검정 유지).
#   텍스트 → #FEFEFE / #DDDDDD 미매핑.
# ★ 검정 박스 크기: 텍스트 폭에 맞게 heuristic 계산 (배경 이미지 최대한 노출).
#   Korean 글자폭 ≈ 폰트pt × 0.014"/pt 기준으로 근사 (26pt ≈ 0.36", 11pt ≈ 0.15").
#   auto_size = TEXT_TO_SHAPE_FIT 이 안전망 (박스 초과 시 폰트 축소).
# ★ z-order: 배경 rect 를 shapes 리스트 맨 앞 append → 나머지 오버레이는 위에.
#
# 입력 스키마:
#   slide_data["title"]        = str (필수, 거버닝)
#   slide_data["points"]       = list[str] ★ 2, 3, 4개 중 하나 (그 외 [])
#   slide_data["badge"]?       = str (상단 알약 배지)
#   slide_data["subtitle"]?    = str (상단 거버닝 박스 안)
#   slide_data["note"]?        = str (좌하단 ※ 문구)
#   slide_data["label_badge"]? = str (좌하단 별도 알약 배지)
# 안전망: title 없음 / points 2~4 아님 → 빈 list.
# ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
def _build_preset_fullbleed_overlay(slide_data):
    title = str(slide_data.get("title", "")).strip()
    if not title:
        return []
    points_raw = slide_data.get("points") or []
    if not isinstance(points_raw, list):
        return []
    points = [str(p).strip() for p in points_raw if str(p).strip()]
    # ★ Spec Preset-Loosen-Count — 정확 (2,3,4) 강제 폐지. 최대 4개 잘라내고 최소 1개면 렌더.
    if len(points) > 4:
        points = points[:4]
    if not points:
        return []

    badge       = str(slide_data.get("badge", "")).strip()
    subtitle    = str(slide_data.get("subtitle", "")).strip()
    note        = str(slide_data.get("note", "")).strip()
    label_badge = str(slide_data.get("label_badge", "")).strip()

    W, H = 11.69, 8.27
    shapes = []

    # ── ① 배경 풀블리드 이미지 placeholder (반드시 먼저 — z-order 하단)
    shapes.append({"type":"rect","x":0,"y":0,"w":W,"h":H,
                   "fill":"#F5F5F5"})
    shapes.append({"type":"text","x":0,"y":H / 2 - 0.2,"w":W,"h":0.4,
                   "text":"이미지 영역","size":13,"weight":400,"color":"#AAAAAA",
                   "align":"center","valign":"middle"})

    # ── ② 상단 오버레이 (중앙정렬)
    #   ★ Spec Preset-FullbleedNoBox — 텍스트 뒤 검정 박스 전량 제거.
    #     이유: placeholder 회색 대비 박스가 튐 + 사용자가 밝은 실사진 쓸 때
    #     박스가 걸리적거림. 대신 텍스트를 검정 계열(#1A1A1A / #666666)로 두고,
    #     어두운 배경 이미지 쓰는 사용자는 PPT 에서 직접 색 조정.
    #   ★ 배지(badge / label_badge) 는 그 자체가 디자인 요소라 유지 (원본 캡쳐 정합).
    #     대표님 갤러리 검토 후 튀면 후속 제거.
    top_y = 0.35
    # (a) 상단 배지 pill (선택) — 유지
    if badge:
        b_w = 1.6
        b_h = 0.4
        b_x = (W - b_w) / 2
        shapes.append({"type":"rect","x":b_x,"y":top_y,"w":b_w,"h":b_h,
                       "fill":"#1B1B1B","radius":0.1})
        shapes.append({"type":"text","x":b_x,"y":top_y,"w":b_w,"h":b_h,
                       "text":badge,"size":11,"weight":700,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
        top_y += b_h + 0.15
    else:
        top_y = 0.55   # 배지 없으면 살짝 아래에서 시작

    # (b) 거버닝 텍스트 — 검정 박스 제거, 텍스트만 중앙정렬.
    #     subtitle 검정 계열 회색으로, title 은 role:"governing" 유지 (자동 보라).
    text_x = 0.9
    text_w = W - 1.8   # 넉넉히 슬라이드 폭 사용 (박스 제약 없음)
    if subtitle:
        shapes.append({"type":"text","x":text_x,"y":top_y,"w":text_w,"h":0.35,
                       "text":subtitle,"size":12,"weight":500,"color":"#666666",
                       "align":"center","valign":"middle"})
        top_y += 0.4
    # title (거버닝 — role:"governing" 자동 보라)
    shapes.append({"type":"text","x":text_x,"y":top_y,"w":text_w,"h":0.9,
                   "text":title,"size":26,"weight":800,"color":"#1A1A1A",
                   "align":"center","valign":"middle",
                   "role":"governing"})

    # ── ③ 좌하단 오버레이 (박스 제거, 텍스트만)
    #   points → label_badge → note 순서로 아래에서 위로 스택.
    #   points·note 는 검정 텍스트, label_badge 는 pill 유지.
    left_x   = 0.9
    line_h   = 0.32
    points_bottom = H - 0.5
    pts_top = points_bottom - len(points) * line_h

    # (c) points 리스트 (박스 없음, 검정 텍스트)
    for i, p in enumerate(points):
        py = pts_top + i * line_h
        shapes.append({"type":"text","x":left_x,"y":py,"w":8.0,"h":line_h,
                       "text":"▶ " + p,"size":11,"weight":500,"color":"#1A1A1A",
                       "align":"left","valign":"middle"})

    # (d) label_badge (선택) — pill 유지, points 위
    if label_badge:
        lb_est_w = len(label_badge) * 0.16 + 0.4
        lb_w = max(1.5, min(lb_est_w, 4.0))
        lb_h = 0.35
        lb_y = pts_top - lb_h - 0.15
        shapes.append({"type":"rect","x":left_x,"y":lb_y,"w":lb_w,"h":lb_h,
                       "fill":"#1B1B1B","radius":0.1})
        shapes.append({"type":"text","x":left_x,"y":lb_y,"w":lb_w,"h":lb_h,
                       "text":label_badge,"size":10,"weight":700,"color":"#FEFEFE",
                       "align":"center","valign":"middle"})
        stack_top = lb_y - 0.15
    else:
        stack_top = pts_top - 0.15

    # (e) note (선택) — 박스 없음, 검정 텍스트
    if note:
        note_h = 0.35
        note_y = stack_top - note_h
        shapes.append({"type":"text","x":left_x,"y":note_y,"w":8.0,"h":note_h,
                       "text":"※ " + note,"size":10,"weight":500,"color":"#1A1A1A",
                       "align":"left","valign":"middle"})
    return shapes


def generate_from_shape_json(json_data, output_path, *, theme="light"):
    """도형 JSON → PPTX (마스터 무관, AI 가 layout 자유 결정 모드).

    json_data 형식:
      {
        "title": "...",
        "slide_width": 11.69,       # 옵션 (inch). 기본 A4 가로 (한국 B2G 표준)
        "slide_height": 8.27,       # 옵션
        "slides": [
          {
            "section": "표지",
            "shapes": [
              {"type": "rect", "x": 0, "y": 0, "w": 0.5, "h": 8.3, "fill": "#000"},
              {"type": "text", "x": 1, "y": 2, "w": 6, "h": 1.5,
               "text": "수주", "size": 80, "weight": 900},
              ...
            ]
          },
          ...
        ]
      }

    [Spec D-Build-ThemeConnect 1-b — 다크 테마 토대 연결]
    theme — 'light'(기본) / 'dark'.
      'light' 일 때 → 분기 미진입(if theme == "dark" 가드), 동작 100% 무변경.
      'dark'  일 때 → (1) 슬라이드별 검정 배경 rect 를 "맨 처음" prepend,
                      (2) text 도형 color 누락 시 기본 글자색을 흰색(FG)으로.
    ★ 명시된 color 는 변환하지 않음 (이번 단계는 기본값만 — 1-c 이후 명시색 매핑).
    main.py 호출부에서 _get_policy('theme', 'light') 결과를 전달.
    """
    if not isinstance(json_data, dict):
        raise ValueError("json_data 가 dict 가 아님")
    slides_data = json_data.get("slides")
    if not isinstance(slides_data, list) or not slides_data:
        raise ValueError("slides 배열 비어있거나 list 아님")

    # 기본값 A4 가로 (11.69×8.27, 인치) — 한국 B2G 공공입찰 표준 인쇄 비율.
    # 시스템 프롬프트 / OUTLINE / SLIDE 모두 동일 기본값으로 정합.
    sw = float(json_data.get("slide_width", 11.69))
    sh = float(json_data.get("slide_height", 8.27))

    prs = Presentation()
    prs.slide_width = Inches(sw)
    prs.slide_height = Inches(sh)
    blank_layout = prs.slide_layouts[6]

    # Spec D-Build-ThemeConnect 1-b — theme 토큰 해석.
    # 라이트는 분기 미진입이므로 토큰 dict 만 가져와 두고 사용 안 함(아래 가드 참조).
    # default_text_color: 라이트=#1A1A1A (현재 운영 그대로) / 다크=tokens["FG"]=#FFFFFF.
    _tokens = _get_theme(theme) if theme == "dark" else None
    if theme == "dark" and _tokens:
        default_text_color = _tokens.get("FG", "#FFFFFF")
        dark_bg = _tokens.get("BG", "#0A0A0A")
    else:
        default_text_color = "#1A1A1A"  # 라이트 = 현재 운영 기본값(무변경 보장)
        dark_bg = None

    rendered_total = 0
    errors_total = []
    # Spec D-Fix-DividerActuallyWorks — 챕터 간지 등장 순번 카운터.
    # _is_chapter_divider 가 True 인 페이지를 만날 때마다 ++. _build_preset_divider 에
    # 1-based 순번 전달 → 코드가 로마숫자(Ⅰ~Ⅹ) 자동 매김. LLM 로마숫자 자유 매김(V 중복) 차단.
    _divider_counter = 0
    for slide_idx, slide_data in enumerate(slides_data):
        if not isinstance(slide_data, dict):
            errors_total.append("slide" + str(slide_idx) + ": not a dict")
            prs.slides.add_slide(blank_layout)
            continue
        slide = prs.slides.add_slide(blank_layout)
        # Spec D-Build-ThemeConnect 1-b — 다크 모드만 검정 배경 rect prepend.
        # 라이트는 이 블록 미진입 → 흰 바탕 inherit (PPT 기본, 현재 동작 그대로).
        # 다른 모든 도형 위에 깔리지 않도록 슬라이드 첫 도형으로 추가(추가 순=하단).
        if theme == "dark" and dark_bg:
            try:
                _add_rect(slide, 0, 0, sw, sh, fill=dark_bg, stroke=None)
            except Exception as _bg_err:
                errors_total.append(
                    "slide" + str(slide_idx) + ":dark_bg: " +
                    type(_bg_err).__name__ + ": " + str(_bg_err)
                )
        # Spec D-Fix-Preset1 / D-Fix-PresetNoOverlap — 옵트인 레이아웃 프리셋.
        # 분기 패턴: preset 성공(도형 리스트 != []) → preset 도형만(백업 폐기, 겹침 방지).
        #            preset 실패(필수 키 누락 → []) → LLM 백업 fallback(빈 페이지 방지).
        #            예외 발생 → LLM 백업 fallback(except 절, 기존 그대로).
        # Spec D-Build-PresetDivider — 챕터 간지 식별 가드 (preset_name 분기보다 먼저).
        # divider 페이지는 viz_pattern="" 이라 LLM 이 preset 키를 안 박음 → preset_name 분기 미진입.
        # slide_type=hero + section/governing_main 패턴으로 식별 → 코드 고정 템플릿 강제.
        # 성공 시 LLM 백업 폐기(circles/quad 동일 패턴), 실패 시 백업 fallback.
        preset_name = slide_data.get("preset")
        if _is_chapter_divider(slide_data):
            _divider_counter += 1
            try:
                preset_shapes = _build_preset_divider(slide_data, _divider_counter)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "quantitative":
            try:
                preset_shapes = _build_preset_quantitative_emphasis(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "process":
            try:
                preset_shapes = _build_preset_horizontal_process(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "two_column":
            try:
                preset_shapes = _build_preset_two_column(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "narrative":
            try:
                preset_shapes = _build_preset_narrative(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "split":
            # Spec D-Build-PresetSplit — 색면 2분할.
            try:
                preset_shapes = _build_preset_split(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "timeline":
            # Spec D-Build-PresetTimeline — 세로 점 단계형 (1단계: 코드 등록만, viz_pattern 미연결).
            # split 과 동일 패턴 — 성공 시 preset 만(겹침 차단), 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_timeline(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "asymmetric":
            # Spec D-Build-PresetAsymmetric — 비대칭 2분할(좌 흰/우 검정) (1단계: 코드 등록만, viz_pattern 미연결).
            # timeline/split 과 동일 패턴 — 성공 시 preset 만(겹침 차단), 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_asymmetric(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "zigzag":
            # Spec D-Build-PresetZigzag — 가운데 축 좌우 번갈아 배치 (1단계: 코드 등록만, viz_pattern 미연결).
            # asymmetric/timeline/split 과 동일 패턴 — 성공 시 preset 만(겹침 차단), 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_zigzag(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "hsplit":
            # Spec D-Build-PresetHsplit — 가로 분할(위 이미지/아래 텍스트) (1단계: 코드 등록만, viz_pattern 미연결).
            # zigzag/asymmetric/timeline/split 과 동일 패턴 — 성공 시 preset 만(겹침 차단), 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_hsplit(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "circles":
            # Spec D-Build-PresetCircles — 원형 가로 정렬(수치/키워드 1~4개) (1단계: 코드 등록만, viz_pattern 미연결).
            # hsplit/zigzag/asymmetric/timeline/split 과 동일 패턴 — 성공 시 preset 만(겹침 차단), 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_circles(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "hsplit_top":
            # Spec D-Build-PresetHsplitTop — 가로 비대칭 2분할 (위 검정 거버닝 + 아래 흰 좌우 2항목).
            # ★ 본 spec 단계: 코드만 등록 + viz_pattern 화이트리스트 미연결 (LLM 미선택 보장).
            #   완성 후 별도 spec 으로 한꺼번에 켤 것.
            # split 가로 비대칭 버전 — 성공 시 preset 만, 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_hsplit_top(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "quad":
            # Spec D-Build-PresetQuad — 색면 4분할(흑/백/흑/백 세로 면) (1단계: 코드 등록만, viz_pattern 미연결).
            # split 의 확장 — circles/hsplit/zigzag/asymmetric/timeline/split 과 동일 패턴.
            try:
                preset_shapes = _build_preset_quad(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "hero_cards":
            # Spec D-Build-PresetHeroCards — 상단 거버닝 히어로 + 하단 카드 N개(2~4).
            # ★ 본 spec 단계: 코드만 등록 + viz_pattern 화이트리스트 미연결 (LLM 미선택 보장).
            #   hsplit_top 와 동일 — 별도 spec 으로 한꺼번에 켤 것.
            # circles 와 동일 패턴 — 성공 시 preset 만, 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_hero_cards(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "triad":
            # Spec Preset-New-Triad — 좌 거버닝 + 우 원 3 + 실선 + 하단 라벨/설명 (비대칭 입체).
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 + viz_pattern 화이트리스트 미연결.
            # circles/hero_cards 와 동일 패턴 — 성공 시 preset 만, 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_triad(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "strategy_map":
            # Spec Preset-New-StrategyMap — 대전략 박스 3 + "+" + 하위실행 chevron 5 + 4단 분할.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 + viz_pattern 화이트리스트 미연결.
            # circles/triad/hero_cards 와 동일 패턴 — 성공 시 preset 만, 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_strategy_map(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "conclusion_cards":
            # Spec Preset-New-ConclusionCards — 3열 카드 + 하단 결론 밴드.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트·
            #   SLIDE 프롬프트 미연결 (LLM 우연 선택 방지). 갤러리 검증 후 별도 spec.
            # strategy_map/triad/hero_cards 와 동일 패턴 — 성공 시 preset 만, 실패/예외 시 LLM 백업.
            try:
                preset_shapes = _build_preset_conclusion_cards(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "numbered_columns":
            # Spec Preset-New-NumberedColumns — 검정 헤더 + 초대형 배경 숫자 3~4열 + 결론.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — viz_pattern 화이트리스트·SLIDE 프롬프트 미연결.
            # conclusion_cards / strategy_map / hero_cards 와 동일 패턴.
            try:
                preset_shapes = _build_preset_numbered_columns(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "hero_detail":
            # Spec Preset-New-HeroDetail — 검정 헤더(배지+숫자+텍스트) + 좌 대형 이미지
            # + 우 번호 항목 3(미니 이미지 + 원문자 head + desc) + 결론 2줄.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
            try:
                preset_shapes = _build_preset_hero_detail(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "flow_detail":
            # Spec Preset-New-FlowDetail — 상단 거버닝 + 중단 가로 흐름 4~5단계
            # + 하단 3~4열 상세(각 열 items 3개). 고밀도 레이아웃.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
            try:
                preset_shapes = _build_preset_flow_detail(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "quad_detail":
            # Spec Preset-New-QuadDetail — 좌측정렬 거버닝 + 2×2 상세 블록
            # (각 블록: 검정 헤더바 + 체크리스트 2~3 + 이미지 placeholder).
            # 기존 "quad" 프리셋과 별개. 첫 좌측정렬 신규 거버닝.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
            try:
                preset_shapes = _build_preset_quad_detail(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        elif preset_name == "fullbleed_overlay":
            # Spec Preset-New-FullbleedOverlay — 풀블리드 이미지 + 상단 거버닝 오버레이
            # + 좌하단 텍스트 오버레이. 조감도·시설·컨셉 등 비주얼 우선 슬라이드.
            # ★ 본 spec 단계: 코드 + dispatch 만 등록 — SLIDE 프롬프트·화이트리스트 미연결.
            try:
                preset_shapes = _build_preset_fullbleed_overlay(slide_data)
                if preset_shapes:
                    shapes = preset_shapes
                else:
                    shapes = slide_data.get("shapes", [])
            except Exception:
                shapes = slide_data.get("shapes", [])
        else:
            shapes = slide_data.get("shapes", [])
        if not isinstance(shapes, list):
            errors_total.append("slide" + str(slide_idx) + ": shapes not list")
            continue
        for shape_idx, shape_def in enumerate(shapes):
            try:
                # Spec D-Build-ThemeConnect 1-b — default_text_color 주입.
                # 라이트=#1A1A1A(현재값 그대로), 다크=tokens["FG"]=#FFFFFF.
                # shape_def 에 "color" 명시 시 그 값이 우선(이번 단계는 기본값만).
                # Spec D-Build-ThemeColorMap 1-c — theme 전달 → helper 내부 _map_color 가 role 별 매핑.
                result = render_shape_to_slide(slide, shape_def,
                                               default_text_color=default_text_color,
                                               theme=theme)
                if result is not None:
                    rendered_total += 1
            except Exception as e:
                errors_total.append(
                    "slide" + str(slide_idx) + ":shape" + str(shape_idx) +
                    ": " + type(e).__name__ + ": " + str(e)
                )

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    log.info("도형 JSON 모드 · 슬라이드 %d / 도형 %d 렌더 / 에러 %d",
             len(slides_data), rendered_total, len(errors_total))

    return {
        "slide_count": len(slides_data),
        "rendered_total": rendered_total,
        "errors": errors_total[:10],
        "output_path": str(output_path),
        "size_mb": round(output_path.stat().st_size / 1024 / 1024, 2),
    }


# ─── Spec D-Build-HTMLOutput — HTML → 편집가능 PPTX 변환기 (admin 전용) ───────
# 원본: html_to_pptx_for_ops.py (사용자 제공 코드 전문 이식, 2026-05-31).
# 운영 연결점: main.py:api_proposals_pptx 의 is_html_mode 분기에서 호출.
# 흐름: messages.content HTML 문자열 → Playwright 헤드리스 chromium 렌더
#       → getBoundingClientRect 좌표 추출 (IR JSON)
#       → python-pptx 도형 1:1 매핑 (rect / oval / textbox + 한글 ea 폰트)
# 의존성: playwright (requirements.txt) + chromium (playwright install chromium)
# 토글 OFF 시 본 함수 미호출 — generate_from_shape_json 경로 그대로 (영향 0).

_PX_TO_EMU = 9525  # 1px = 9525 EMU


def _px(v):
    return Emu(int(round(v * _PX_TO_EMU)))


def _html_css_rgb(s):
    """CSS rgb()/rgba() → RGBColor. alpha 채널 무시 (PPTX 도형은 alpha 미지원)."""
    nums = s.replace("rgba", "").replace("rgb", "").strip("() ").split(",")
    return RGBColor(int(float(nums[0])), int(float(nums[1])), int(float(nums[2])))


def _html_is_transparent(s):
    return "rgba(0, 0, 0, 0)" in s or s == "transparent"


def _html_set_ea_font(run, font_name):
    """한글 폰트 ea (East Asian) 속성 강제 주입.
    이 속성 없으면 PPT 에서 한글 깨짐 (영문 폰트로 fallback)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font_name)


def _html_pick_font(family):
    """font-family CSS 값 (콤마 구분) → 첫 폰트명 추출."""
    return family.split(",")[0].strip().strip("'\"")


def _html_extract_all_slides(html_string):
    """HTML 문자열 안의 모든 .slide 를 슬라이드별 IR 배열로 반환.
    Playwright 헤드리스 chromium 으로 렌더 후 getBoundingClientRect 좌표 추출."""
    import os as _os
    import tempfile
    from playwright.sync_api import sync_playwright

    with tempfile.NamedTemporaryFile(
        "w", suffix=".html", delete=False, encoding="utf-8"
    ) as f:
        f.write(html_string)
        tmp_path = f.name

    try:
        with sync_playwright() as p:
            # Playwright default launch 가 chrome-headless-shell 을 찾는 케이스 회피.
            # executable_path 명시 = 환경별로 sync API 가 verified 한 chrome.exe / chrome 사용.
            b = p.chromium.launch(
                headless=True,
                executable_path=p.chromium.executable_path,
            )
            pg = b.new_page(viewport={"width": 1123, "height": 794})
            pg.goto(f"file://{tmp_path}")
            pg.wait_for_timeout(400)  # 폰트 / 레이아웃 안정 대기

            n = pg.evaluate("document.querySelectorAll('.slide').length")
            slides_ir = []
            for i in range(n):
                ir = pg.evaluate("""
                (idx) => {
                  const slide = document.querySelectorAll('.slide')[idx];
                  const sb = slide.getBoundingClientRect();
                  const out = [];
                  slide.querySelectorAll('div').forEach(el => {
                    const r = el.getBoundingClientRect();
                    const cs = getComputedStyle(el);
                    // direct children: text nodes as text, <br> as newline
                    let parts = [];
                    el.childNodes.forEach(nn => {
                      if (nn.nodeType === 3) {
                        const t = nn.textContent.trim();
                        if (t) parts.push(t);
                      } else if (nn.nodeName === 'BR') {
                        parts.push('\\n');
                      }
                    });
                    // join with space, keep newlines clean
                    const text = parts.join(' ')
                                      .replace(/\\s*\\n\\s*/g, '\\n').trim();
                    out.push({
                      x: r.left - sb.left, y: r.top - sb.top,
                      w: r.width, h: r.height, text: text,
                      fontSize: parseFloat(cs.fontSize),
                      fontFamily: cs.fontFamily, fontWeight: cs.fontWeight,
                      color: cs.color, bg: cs.backgroundColor,
                      borderTop: cs.borderTopWidth, borderColor: cs.borderTopColor,
                      borderRadius: cs.borderTopLeftRadius, textAlign: cs.textAlign,
                      cls: el.className || "",
                      lineHeight: cs.lineHeight,
                    });
                  });
                  return { sw: sb.width, sh: sb.height, els: out };
                }
                """, i)
                slides_ir.append(ir)
            b.close()
            return slides_ir
    finally:
        try:
            _os.unlink(tmp_path)
        except OSError:
            pass


def generate_from_html(html_string, out_path):
    """운영 진입점 (Spec D-Build-HTMLOutput).
    messages.content 의 HTML 문자열 → 편집가능 PPTX.
    여러 .slide div 가 있으면 각각 PPTX 슬라이드로 추가.

    반환: dict (slide_count / output_path / size_mb / errors).
    포맷은 generate_from_shape_json 과 호환 (main.py 호출부 같은 응답 처리 가능).
    """
    slides_ir = _html_extract_all_slides(html_string)
    if not slides_ir:
        raise ValueError("변환할 .slide 요소가 HTML에 없음")

    prs = Presentation()
    prs.slide_width = _px(slides_ir[0]["sw"])
    prs.slide_height = _px(slides_ir[0]["sh"])

    rendered_total = 0
    errors = []

    for ir in slides_ir:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        for el in ir["els"]:
            try:
                has_text = bool(el["text"])
                has_bg = not _html_is_transparent(el["bg"])
                has_border = float(
                    (el["borderTop"] or "0").replace("px", "") or 0
                ) > 0
                is_circle = "50%" in el["borderRadius"]

                if (has_bg or has_border) and not is_circle:
                    shp = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, _px(el["x"]), _px(el["y"]),
                        _px(max(el["w"], 1)), _px(max(el["h"], 1)),
                    )
                    if has_bg:
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = _html_css_rgb(el["bg"])
                    else:
                        shp.fill.background()
                    if has_border:
                        shp.line.color.rgb = _html_css_rgb(el["borderColor"])
                        shp.line.width = Pt(1)
                    else:
                        shp.line.fill.background()
                    shp.shadow.inherit = False
                    rendered_total += 1

                if is_circle:
                    shp = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, _px(el["x"]), _px(el["y"]),
                        _px(el["w"]), _px(el["h"]),
                    )
                    if has_bg:
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = _html_css_rgb(el["bg"])
                    else:
                        shp.fill.background()
                    if has_border:
                        shp.line.color.rgb = _html_css_rgb(el["borderColor"])
                        shp.line.width = Pt(1.5)
                    else:
                        shp.line.fill.background()
                    shp.shadow.inherit = False
                    rendered_total += 1

                if has_text:
                    tb = slide.shapes.add_textbox(
                        _px(el["x"]), _px(el["y"]), _px(el["w"]), _px(el["h"]),
                    )
                    tf = tb.text_frame
                    tf.word_wrap = True
                    try:
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE_FIT  # D-Fix-AutoSize-HTML: 박스 넘치면 폰트 자동 축소(도형 모드 _add_text와 동일 안전망)
                    except Exception:
                        pass
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                    # 세로 정렬: 원/박스 안에 얹힌 텍스트는 가운데로.
                    # CSS line-height 가 height 와 같으면(=한 줄 세로중앙 의도) MIDDLE.
                    cls = el.get("cls", "")
                    vcenter = any(k in cls for k in (
                        "kpi-label", "kpi-value", "chip", "pf-card", "cat-box"))
                    if vcenter:
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    # 가로 정렬: 추출한 text-align 그대로 반영 (이전엔 무시됐음).
                    align = {
                        "center": PP_ALIGN.CENTER,
                        "right": PP_ALIGN.RIGHT,
                        "justify": PP_ALIGN.JUSTIFY,
                    }.get((el.get("textAlign") or "left").lower(), PP_ALIGN.LEFT)

                    # \n 으로 여러 줄 → 여러 단락. (HTML <br> 이 \n 으로 추출됨)
                    fsize = Pt(el["fontSize"] * 0.75)
                    fcolor = _html_css_rgb(el["color"])
                    fw = el["fontWeight"]
                    fbold = (fw == "bold" or (fw.isdigit() and int(fw) >= 600))
                    fname = _html_pick_font(el["fontFamily"])

                    lines = (el["text"] or "").split("\n")
                    for li, line in enumerate(lines):
                        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                        p.alignment = align
                        run = p.add_run()
                        run.text = line
                        f = run.font
                        f.size = fsize
                        f.color.rgb = fcolor
                        f.bold = fbold
                        f.name = fname
                        _html_set_ea_font(run, fname)
                    rendered_total += 1
            except Exception as e:
                errors.append(f"el@({el.get('x',0):.0f},{el.get('y',0):.0f}): {type(e).__name__}: {e}")

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    log.info("HTML 모드 (D-Build-HTMLOutput) · 슬라이드 %d / 도형 %d 렌더 / 에러 %d",
             len(slides_ir), rendered_total, len(errors))

    return {
        "slide_count": len(slides_ir),
        "rendered_total": rendered_total,
        "errors": errors[:10],
        "output_path": str(out_path),
        "size_mb": round(out_path.stat().st_size / 1024 / 1024, 2),
    }
