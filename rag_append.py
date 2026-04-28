"""
RAG append: 임의 폴더의 새 파일들을 기존 rag_kb.db 에 추가 학습.

사용법:
    python rag_append.py "D:\\어떤 폴더"            # prefix 자동 (D2, D3, ...)
    python rag_append.py "D:\\어떤 폴더" D5         # prefix 명시

기능:
1) 지정 폴더의 PDF/PPTX 스캔 (~$, RFP/발표자료/결과보고/산출내역/템플릿 등 자동 제외)
2) **기존 DB 의 filename 과 비교 — 이미 학습한 파일(stem 동일)은 자동 스킵**
3) 텍스트 추출 → 청킹 → 임베딩 → 기존 rag_kb.db 에 INSERT (ON CONFLICT DO NOTHING)
4) 최종 청크 수 보고

키 처리:
- C:\\Users\\00\\Desktop\\openai_key.txt 에서 키 읽기 (UTF-16/UTF-8 모두 지원)
- 환경변수 OPENAI_API_KEY 가 있으면 그게 우선
"""
from __future__ import annotations
import argparse
import json
import os
import re
import sqlite3
import struct
import sys
import time
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def safe_print(*args):
    try:
        print(*args, flush=True)
    except UnicodeEncodeError:
        s = " ".join(str(a) for a in args)
        print(s.encode("ascii", errors="replace").decode("ascii"), flush=True)


# ─── 경로 설정 ───
KEY_FILE = Path(r"C:\Users\00\Desktop\openai_key.txt")
DB_PATH = Path("rag_kb.db")

# 인자로 받음 — main() 에서 결정
SOURCE: Path | None = None
EXTRACT_DIR: Path | None = None
CHUNKS_JSON: Path | None = None
REPORT_PATH: Path | None = None
CHUNK_ID_PREFIX = "D2"

# ─── 학습 제외 키워드 ───
EXCLUDE_KEYWORDS = [
    "제안요청서", "RFP",
    "발표자료", "발표용",
    "결과보고", "결과보고서",
    "정산서", "산출내역",
    "템플릿",
    "(붙임",
    "~$",   # MS Office 임시 잠금 파일
]
# 본 학습은 PDF/PPTX 만. PPT(구버전)/HWP/HWPX 는 별도 안내.
SUPPORTED_EXT = {".pdf", ".pptx"}
NEEDS_CONVERT_EXT = {".ppt", ".hwp", ".hwpx", ".doc"}  # 별도 안내용 카운트
SKIP_EXT = {".xlsx", ".xls", ".zip"}                    # 학습 의미 없음 — 조용히 스킵

# ─── 청킹/임베딩 파라미터 (기존과 동일) ───
TARGET_MIN = 400
TARGET_MAX = 600
TARGET_HARD_MAX = 900
EMBED_MODEL = "text-embedding-3-large"
EMBED_DIM = 3072
BATCH_SIZE = 64

# ─── 메타데이터 정규식 (rag_build.py 와 동일) ───
ENDING_PATTERNS = {
    "전략": r"\S{2,}\s*전략\b",
    "시스템": r"\S{2,}\s*시스템\b",
    "체계": r"\S{2,}\s*체계\b",
    "방안": r"\S{2,}\s*방안\b",
    "플랫폼": r"\S{2,}\s*플랫폼\b",
    "경험": r"\S{2,}\s*경험\b",
    "설계": r"\S{2,}\s*설계\b",
    "프로세스": r"\S{2,}\s*프로세스\b",
    "매뉴얼": r"\S{2,}\s*매뉴얼\b",
    "구조": r"\S{2,}\s*구조\b",
}
VISUAL_PATTERNS = {
    "step_flow":   r"\bSTEP\s*\d+\b|\d+단계\b|단계\s*\d+",
    "timeline":    r"\b(타임라인|TIME\s*LINE|로드맵|D-?\d+|D[+]\d+)\b",
    "table":       r"\b(구분|항목|세부내용|비고)\s*[:│|]",
    "comparison":  r"\b(AS[\s-]?IS|TO[\s-]?BE|Before|After|비교|대비)\b",
    "org":         r"\b(조직도|총괄|PM|디렉터|운영진|인력\s*구성)\b",
    "budget":      r"\b(예산|산출\s*내역|단가|총\s*사업비|VAT)\b",
    "safety":      r"\b(안전\s*관리|비상\s*매뉴얼|위기\s*대응|보험|응급)\b",
    "stat_emph":   r"\d+\.?\d*\s*%|\b\d{2,}\s*(?:점|건|개|곳|명|회|일|개월|년)\b",
    "bullet":      r"^\s*[●◆■◇○•\-\*✓☑]\s+",
    "arrow":       r"→|⇒|⇨",
}
SECTION_NAME_RE = re.compile(
    r"(?:Ⅰ|Ⅱ|Ⅲ|Ⅳ|Ⅴ|Ⅵ|Ⅶ|Ⅷ|Ⅸ|Ⅹ|\d{1,2}|[가나다라마바사아])\s*[\.\)]\s*([가-힣A-Za-z][^\n]{1,30})"
)


def safe_filename(s: str) -> str:
    s = re.sub(r'[<>:"/\\|?*]', '_', s)
    return s[:120]


# ─── 추출 함수들 (rag_extract.py 와 동일 로직) ───
def extract_pdf(path: Path) -> dict:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(str(path)) as pdf:
            for i, p in enumerate(pdf.pages, 1):
                try:
                    t = p.extract_text() or ""
                except Exception:
                    t = ""
                pages.append({"page": i, "text": t})
        total = sum(len(p["text"]) for p in pages)
        if total >= 200:
            return {"ok": True, "engine": "pdfplumber", "pages": pages,
                    "page_count": len(pages), "total_chars": total}
    except Exception as e:
        return _try_pypdf(path, prev_err=str(e)[:200])
    return _try_pypdf(path, prev_err="pdfplumber 거의 빈 결과")


def _try_pypdf(path: Path, prev_err: str = "") -> dict:
    try:
        from pypdf import PdfReader
        r = PdfReader(str(path))
        pages = []
        for i, p in enumerate(r.pages, 1):
            try:
                t = p.extract_text() or ""
            except Exception:
                t = ""
            pages.append({"page": i, "text": t})
        total = sum(len(p["text"]) for p in pages)
        return {"ok": True, "engine": "pypdf", "pages": pages,
                "page_count": len(pages), "total_chars": total,
                "fallback_reason": prev_err}
    except Exception as e:
        return {"ok": False, "engine": "none", "error": str(e)[:300]}


def extract_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, sl in enumerate(prs.slides, 1):
            text_parts = []
            for shape in sl.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                text_parts.append(run.text)
                        if para.runs:
                            text_parts.append("\n")
                elif shape.shape_type == 19:
                    try:
                        for row in shape.table.rows:
                            row_cells = [cell.text_frame.text if cell.text_frame else "" for cell in row.cells]
                            text_parts.append(" | ".join(row_cells))
                            text_parts.append("\n")
                    except Exception:
                        pass
            notes_text = ""
            if sl.has_notes_slide and sl.notes_slide and sl.notes_slide.notes_text_frame:
                notes_text = sl.notes_slide.notes_text_frame.text or ""
            text = "".join(text_parts).strip()
            slides.append({"page": i, "text": text, "notes": notes_text})
        total = sum(len(s["text"]) for s in slides)
        return {"ok": True, "engine": "python-pptx", "pages": slides,
                "page_count": len(slides), "total_chars": total}
    except Exception as e:
        return {"ok": False, "engine": "none",
                "error": f"{type(e).__name__}: {str(e)[:250]}"}


# ─── 청킹/메타 (rag_build.py 와 동일 로직) ───
def extract_meta(text: str) -> dict:
    meta = {"ending_hits": [], "visual_hits": [], "section_hints": [],
            "char_count": len(text)}
    for label, pat in ENDING_PATTERNS.items():
        if re.search(pat, text):
            meta["ending_hits"].append(label)
    for label, pat in VISUAL_PATTERNS.items():
        flags = re.M if "^" in pat else 0
        if re.search(pat, text, flags=flags):
            meta["visual_hits"].append(label)
    for m in SECTION_NAME_RE.finditer(text[:600]):
        name = m.group(0).strip()
        if 3 < len(name) < 40 and name not in meta["section_hints"]:
            meta["section_hints"].append(name)
        if len(meta["section_hints"]) >= 3:
            break
    return meta


def split_long_text(text: str, max_chars: int = TARGET_HARD_MAX) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    parts = []
    buf = ""
    sentences = re.split(r"(?<=[.!?…])\s+|\n{2,}", text)
    for s in sentences:
        s = s.strip()
        if not s:
            continue
        if len(buf) + len(s) + 1 <= max_chars:
            buf = (buf + " " + s) if buf else s
        else:
            if buf:
                parts.append(buf)
            if len(s) > max_chars:
                for i in range(0, len(s), max_chars):
                    parts.append(s[i:i + max_chars])
                buf = ""
            else:
                buf = s
    if buf:
        parts.append(buf)
    return parts


def chunk_pages(pages: list[dict], filestem: str) -> list[dict]:
    """추출된 pages 리스트 → 청크. chunk_id 는 모듈 전역 CHUNK_ID_PREFIX."""
    chunks = []
    buf_pages: list[int] = []
    buf_text = ""
    chunk_idx = 0
    pfx = CHUNK_ID_PREFIX

    def flush():
        nonlocal buf_text, buf_pages, chunk_idx
        if buf_text.strip():
            chunks.append({
                "chunk_id": f"{pfx}__{filestem}__{chunk_idx:04d}",
                "filename": filestem,
                "pages": list(buf_pages),
                "text": buf_text.strip(),
                "char_count": len(buf_text.strip()),
            })
            chunk_idx += 1
        buf_pages.clear()
        buf_text = ""

    for p in pages:
        ptext = (p.get("text") or "").strip()
        if not ptext:
            continue
        if len(ptext) > TARGET_HARD_MAX:
            if buf_text:
                flush()
            for sub in split_long_text(ptext, TARGET_HARD_MAX):
                chunks.append({
                    "chunk_id": f"{pfx}__{filestem}__{chunk_idx:04d}",
                    "filename": filestem,
                    "pages": [p["page"]],
                    "text": sub.strip(),
                    "char_count": len(sub.strip()),
                })
                chunk_idx += 1
            continue
        if buf_text and len(buf_text) + len(ptext) + 2 > TARGET_MAX:
            flush()
        buf_text = (buf_text + "\n\n" + ptext).strip() if buf_text else ptext
        buf_pages.append(p["page"])
        if len(buf_text) >= TARGET_MIN:
            flush()
    if buf_text:
        flush()

    for c in chunks:
        c.update(extract_meta(c["text"]))
    return chunks


def auto_pick_prefix(db) -> str:
    """이미 사용된 D{n}__ prefix 를 보고 다음 번호 결정."""
    used = set()
    for r in db.execute("SELECT chunk_id FROM chunks").fetchall():
        cid = r[0]
        m = re.match(r"^(D\d+)__", cid or "")
        if m:
            used.add(m.group(1))
    n = 2
    while f"D{n}" in used:
        n += 1
    return f"D{n}"


def get_existing_filenames(db) -> set[str]:
    """이미 학습된 파일 stem 집합 — 자동 중복 제외용."""
    rows = db.execute("SELECT DISTINCT filename FROM chunks").fetchall()
    return {r[0] for r in rows if r[0]}


# ─── DB append (sqlite-vec) ───
def open_db():
    if not DB_PATH.exists():
        safe_print(f"[!] {DB_PATH} 가 없습니다. rag_build.py 를 먼저 실행하세요.")
        sys.exit(1)
    db = sqlite3.connect(str(DB_PATH))
    db.enable_load_extension(True)
    import sqlite_vec
    sqlite_vec.load(db)
    db.enable_load_extension(False)
    return db


def get_api_key() -> str:
    """OpenAI API 키 로드.
    1) 환경변수 우선
    2) 메모장 파일 — utf-8/utf-8-sig/utf-16/cp949 순으로 디코딩 시도
       파일 안의 모든 줄을 훑어 'sk-' 로 시작하는 첫 줄 반환 (라벨 줄 무시)
    """
    env = os.environ.get("OPENAI_API_KEY", "").strip()
    if env:
        return env
    if not KEY_FILE.exists():
        return ""
    try:
        data = KEY_FILE.read_bytes()
    except Exception:
        return ""
    text = ""
    for enc in ("utf-8-sig", "utf-8", "utf-16", "utf-16-le", "utf-16-be", "cp949"):
        try:
            text = data.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        return ""
    for line in text.splitlines():
        s = line.strip()
        if s.startswith("sk-"):
            return s
    return ""


def main():
    global SOURCE, EXTRACT_DIR, CHUNKS_JSON, REPORT_PATH, CHUNK_ID_PREFIX

    parser = argparse.ArgumentParser(description="RAG 추가 학습 (기존 rag_kb.db 에 append)")
    parser.add_argument("source", help="학습 폴더 절대경로")
    parser.add_argument("prefix", nargs="?", default=None,
                        help="청크 ID prefix (예: D3). 생략 시 DB 보고 자동 결정")
    args = parser.parse_args()

    SOURCE = Path(args.source)
    if not SOURCE.is_dir():
        safe_print(f"[X] 폴더 없음: {SOURCE}")
        sys.exit(1)

    api_key = get_api_key()
    if not api_key:
        safe_print("[X] OPENAI_API_KEY 못 찾음 (env 또는 openai_key.txt)")
        sys.exit(1)
    safe_print(f"[OK] API 키 로드 (...{api_key[-6:]})")

    # ── 0) DB 열고 prefix 자동 결정 + 이미 학습된 파일 셋 ──
    db_pre = open_db()
    if args.prefix:
        CHUNK_ID_PREFIX = args.prefix
    else:
        CHUNK_ID_PREFIX = auto_pick_prefix(db_pre)
    existing_stems = get_existing_filenames(db_pre)
    db_pre.close()

    # 캐시/보고서 경로 — prefix 별로 분리
    EXTRACT_DIR = Path(f"_rag_extracted_{CHUNK_ID_PREFIX.lower()}")
    CHUNKS_JSON = Path(f"_rag_chunks_{CHUNK_ID_PREFIX.lower()}.json")
    REPORT_PATH = Path(f"_rag_append_report_{CHUNK_ID_PREFIX.lower()}.txt")
    EXTRACT_DIR.mkdir(exist_ok=True)

    safe_print(f"  대상 폴더: {SOURCE}")
    safe_print(f"  chunk_id prefix: {CHUNK_ID_PREFIX}")
    safe_print(f"  기존 학습된 파일 stem: {len(existing_stems)}개")

    # ── 1) 대상 파일 수집 ──
    candidates = []
    excluded = []
    needs_convert = []   # PPT/HWP/HWPX/DOC 별도 안내
    skipped_ext = []
    skipped_dup = []     # 이미 학습됨 — 자동 중복 제외
    for f in sorted(SOURCE.iterdir(), key=lambda p: p.name):
        if not f.is_file():
            continue
        ext = f.suffix.lower()
        if ext in NEEDS_CONVERT_EXT:
            needs_convert.append(f.name)
            continue
        if ext in SKIP_EXT:
            skipped_ext.append(f.name)
            continue
        if ext not in SUPPORTED_EXT:
            skipped_ext.append(f.name)
            continue
        if any(kw in f.name for kw in EXCLUDE_KEYWORDS):
            excluded.append(f.name)
            continue
        if f.stem in existing_stems:
            skipped_dup.append(f.name)
            continue
        candidates.append(f)

    safe_print(f"\n=== 1) 대상 수집 ===")
    safe_print(f"  학습 대상: {len(candidates)}개")
    safe_print(f"  EXCLUDE 키워드 제외: {len(excluded)}개")
    safe_print(f"  이미 학습됨 (자동 스킵): {len(skipped_dup)}개")
    safe_print(f"  PPT/HWP 등 변환 필요 (별도 안내): {len(needs_convert)}개")
    safe_print(f"  기타 확장자 무시 (xlsx/zip 등): {len(skipped_ext)}개")
    if needs_convert:
        safe_print(f"     {needs_convert[:5]}{' ...' if len(needs_convert) > 5 else ''}")

    if not candidates:
        safe_print("학습할 새 파일이 없어요.")
        return

    # ── 2) 텍스트 추출 ──
    safe_print(f"\n=== 2) 텍스트 추출 시작 ===")
    extracted = []   # [{filestem, pages, ok, ...}]
    failed = []
    t_start = time.time()
    for idx, f in enumerate(candidates, 1):
        size_mb = f.stat().st_size / 1024 / 1024
        ext = f.suffix.lower()
        out_path = EXTRACT_DIR / (safe_filename(f.stem) + ".txt")
        # 캐시: 이미 추출된 파일이면 재사용
        if out_path.exists() and out_path.stat().st_size > 200:
            try:
                txt = out_path.read_text(encoding="utf-8")
                # 페이지 분리
                page_re = re.compile(r"=====\s*PAGE\s+(\d+)\s*=====")
                parts = page_re.split(txt)
                pages = []
                for i in range(1, len(parts), 2):
                    pn = int(parts[i])
                    body = parts[i + 1].strip() if i + 1 < len(parts) else ""
                    body = re.sub(r"\[NOTES\][\s\S]*?(?=\n=====|\Z)", "", body).strip()
                    if body:
                        pages.append({"page": pn, "text": body})
                if pages:
                    extracted.append({"filestem": f.stem, "pages": pages,
                                      "char_count": sum(len(p["text"]) for p in pages),
                                      "size_mb": size_mb, "cached": True})
                    safe_print(f"[{idx:3d}/{len(candidates)}] ({size_mb:6.1f}MB) {ext} CACHE  {f.name[:55]}")
                    continue
            except Exception:
                pass

        t0 = time.time()
        safe_print(f"[{idx:3d}/{len(candidates)}] ({size_mb:6.1f}MB) {ext}        {f.name[:55]}")
        try:
            result = extract_pdf(f) if ext == ".pdf" else extract_pptx(f)
        except Exception as e:
            result = {"ok": False, "error": str(e)[:200]}

        elapsed = time.time() - t0
        if not result.get("ok"):
            failed.append({"filename": f.name, "error": result.get("error", "")[:200]})
            safe_print(f"      -> FAIL  ({elapsed:.1f}s)  {result.get('error','')[:80]}")
            continue

        # 추출 캐시 저장
        with out_path.open("w", encoding="utf-8") as wf:
            wf.write(f"# {f.name}\n")
            wf.write(f"engine={result.get('engine')} pages={result.get('page_count')} total_chars={result.get('total_chars')}\n\n")
            for p in result.get("pages", []):
                wf.write(f"\n===== PAGE {p['page']} =====\n")
                wf.write(p.get("text", "") + "\n")
                if p.get("notes"):
                    wf.write(f"\n[NOTES]\n{p['notes']}\n")

        # NOTES 제외하고 본문만 추출에 사용
        pages = [{"page": p["page"], "text": p.get("text", "")} for p in result.get("pages", []) if p.get("text", "").strip()]
        extracted.append({"filestem": f.stem, "pages": pages,
                          "char_count": result.get("total_chars", 0),
                          "size_mb": size_mb, "cached": False,
                          "engine": result.get("engine")})
        safe_print(f"      -> OK  ({result.get('page_count')}p, {result.get('total_chars'):,} chars, {elapsed:.1f}s)")

    safe_print(f"\n추출 완료: {len(extracted)}개 성공 / {len(failed)}개 실패 / 총 {time.time()-t_start:.0f}s")

    # ── 3) 청킹 ──
    safe_print(f"\n=== 3) 청킹 ===")
    all_chunks = []
    for ex in extracted:
        cks = chunk_pages(ex["pages"], ex["filestem"])
        all_chunks.extend(cks)
        safe_print(f"  [chunk] {ex['filestem'][:55]} -> {len(cks)} chunks")
    safe_print(f"\n총 청크: {len(all_chunks)}개")
    if not all_chunks:
        safe_print("청크가 0개라서 종료.")
        return
    avg_len = sum(c["char_count"] for c in all_chunks) / len(all_chunks)
    safe_print(f"평균 청크 길이: {avg_len:.0f}자")

    CHUNKS_JSON.write_text(json.dumps(all_chunks, ensure_ascii=False, indent=2),
                            encoding="utf-8")

    # ── 4) 임베딩 ──
    safe_print(f"\n=== 4) 임베딩 ({EMBED_MODEL}) ===")
    from openai import OpenAI
    client = OpenAI(api_key=api_key)

    embeddings: list = [None] * len(all_chunks)
    t0 = time.time()
    for i in range(0, len(all_chunks), BATCH_SIZE):
        batch = all_chunks[i:i + BATCH_SIZE]
        texts = [c["text"] for c in batch]
        try:
            resp = client.embeddings.create(model=EMBED_MODEL, input=texts)
            for j, item in enumerate(resp.data):
                embeddings[i + j] = item.embedding
            elapsed = time.time() - t0
            safe_print(f"  [{i + len(batch):5d}/{len(all_chunks)}] {elapsed:.1f}s")
        except Exception as e:
            safe_print(f"  [X] 배치 {i} 실패: {str(e)[:120]}")

    valid = sum(1 for e in embeddings if e is not None)
    safe_print(f"\n임베딩: {valid}/{len(all_chunks)} 성공")
    if valid == 0:
        safe_print("임베딩 0개 — append 중단")
        return

    # ── 5) 기존 DB 에 INSERT ──
    safe_print(f"\n=== 5) 기존 rag_kb.db 에 INSERT ===")
    db = open_db()

    # 기존 청크 수
    before_chunks = db.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
    before_vecs = db.execute("SELECT COUNT(*) FROM vec_chunks").fetchone()[0]
    safe_print(f"  기존 chunks: {before_chunks}건 / vec_chunks: {before_vecs}건")

    inserted = 0
    db_dup_count = 0   # 변수명 충돌 회피 (skipped_dup 은 파일 단위 스킵 list)
    for chunk, emb in zip(all_chunks, embeddings):
        if emb is None:
            continue
        try:
            cur = db.execute(
                """INSERT INTO chunks(chunk_id,filename,pages,text,char_count,
                                      ending_hits,visual_hits,section_hints,embedded)
                   VALUES(?,?,?,?,?,?,?,?,1)
                   ON CONFLICT(chunk_id) DO NOTHING""",
                (chunk["chunk_id"], chunk["filename"], json.dumps(chunk["pages"]),
                 chunk["text"], chunk["char_count"],
                 json.dumps(chunk.get("ending_hits", []), ensure_ascii=False),
                 json.dumps(chunk.get("visual_hits", []), ensure_ascii=False),
                 json.dumps(chunk.get("section_hints", []), ensure_ascii=False)),
            )
            if cur.rowcount == 0:
                db_dup_count += 1
                continue
            rid = cur.lastrowid
            emb_bytes = struct.pack(f"{EMBED_DIM}f", *emb)
            db.execute("INSERT INTO vec_chunks(rowid, embedding) VALUES(?, ?)",
                       (rid, emb_bytes))
            inserted += 1
        except Exception as e:
            safe_print(f"  [X] insert 실패 {chunk['chunk_id']}: {str(e)[:80]}")
    db.commit()

    after_chunks = db.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
    after_vecs = db.execute("SELECT COUNT(*) FROM vec_chunks").fetchone()[0]
    db.close()

    safe_print(f"  INSERT: {inserted}건 / 중복 스킵 {db_dup_count}건")
    safe_print(f"  최종 chunks: {after_chunks}건 (+{after_chunks - before_chunks})")
    safe_print(f"  최종 vec_chunks: {after_vecs}건 (+{after_vecs - before_vecs})")
    safe_print(f"  DB 크기: {DB_PATH.stat().st_size / 1024 / 1024:.1f} MB")

    # ── 6) 보고서 ──
    safe_print(f"\n=== 6) 보고서 작성 ===")
    files_per_chunks = {}
    for c in all_chunks:
        files_per_chunks.setdefault(c["filename"], 0)
        files_per_chunks[c["filename"]] += 1

    lines = [
        f"RAG append 리포트 ({SOURCE} → 기존 rag_kb.db)",
        "=" * 70,
        f"대상 폴더       : {SOURCE}",
        f"chunk_id prefix : {CHUNK_ID_PREFIX}",
        f"임베딩 모델     : {EMBED_MODEL} ({EMBED_DIM} dim)",
        "",
        "[1] 파일 현황",
        f"  학습 대상 (PDF/PPTX)        : {len(candidates)}개",
        f"  EXCLUDE 키워드 제외         : {len(excluded)}개",
        f"  이미 학습된 파일 자동 스킵  : {len(skipped_dup)}개",
        f"  변환 필요 (PPT/HWP/DOC 등)  : {len(needs_convert)}개  -> 별도 안내 필요",
        f"  기타 확장자 무시 (xlsx/zip) : {len(skipped_ext)}개",
        "",
        "[2] 추출 결과",
        f"  성공: {len(extracted)}개",
        f"  실패: {len(failed)}개",
    ]
    for fail in failed[:20]:
        lines.append(f"     · {fail['filename']}  -> {fail['error']}")
    if len(failed) > 20:
        lines.append(f"     ... 외 {len(failed) - 20}개")

    lines += [
        "",
        "[3] 청킹/임베딩",
        f"  생성 청크: {len(all_chunks)}",
        f"  평균 청크 길이: {avg_len:.0f}자",
        f"  임베딩 성공: {valid}/{len(all_chunks)}",
        "",
        "[4] DB 변화",
        f"  before: {before_chunks} chunks / {before_vecs} vec",
        f"  after : {after_chunks} chunks / {after_vecs} vec",
        f"  newly inserted: {inserted} (중복 스킵 {db_dup_count})",
        "",
        "[5] 파일별 청크 수 (top 30)",
    ]
    for fn, cnt in sorted(files_per_chunks.items(), key=lambda x: -x[1])[:30]:
        lines.append(f"  {cnt:4d}  {fn[:80]}")

    lines += ["", "[6] EXCLUDE 처리된 파일"]
    for name in excluded:
        lines.append(f"  - {name[:90]}")

    if skipped_dup:
        lines += ["", f"[7] 이미 학습된 파일 (자동 스킵 {len(skipped_dup)}개)"]
        for name in skipped_dup[:50]:
            lines.append(f"  - {name[:90]}")
        if len(skipped_dup) > 50:
            lines.append(f"  ... 외 {len(skipped_dup) - 50}개")

    if needs_convert:
        lines += ["", f"[8] 변환 필요 (PPT/HWP/HWPX/DOC) — 학습 안 됨 ({len(needs_convert)}개)"]
        for name in needs_convert:
            lines.append(f"  - {name[:90]}")
        lines += [
            "",
            "    > 변환 방법:",
            "    >   - PPT (구버전): PowerPoint 에서 열어 .pptx 로 다시 저장",
            "    >   - HWP/HWPX:    한컴 오피스에서 PDF 또는 PPTX 로 export",
            "    >   - DOC:         Word 에서 .docx 또는 .pdf 로 다시 저장",
        ]

    REPORT_PATH.write_text("\n".join(lines), encoding="utf-8")
    safe_print(f"  보고서: {REPORT_PATH}")
    safe_print(f"\n=== 완료 ===")
    safe_print(f"  최종 RAG 청크 수: {after_chunks}건  (이번에 +{after_chunks - before_chunks})")


if __name__ == "__main__":
    main()
