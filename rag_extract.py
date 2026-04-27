"""
RAG 1단계: 17개 제안서 텍스트 추출 + 품질 리포트
- PDF: pdfplumber 우선 → pypdf fallback
- PPTX: python-pptx (slide.shapes.text + notes)
출력:
- _rag_extracted/{filename}.txt  (개별 텍스트)
- _rag_quality.json (품질 메트릭)
- _rag_report.txt (사람이 읽을 보고서)
"""
import json
import sys
import time
import traceback
from pathlib import Path

# Windows cp949 가 표현 못 하는 한글/특수문자 — stdout 을 UTF-8 로 강제
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def safe_print(*args):
    """cp949 인코딩 에러 회피용 안전 print."""
    try:
        print(*args, flush=True)
    except UnicodeEncodeError:
        s = " ".join(str(a) for a in args)
        print(s.encode("ascii", errors="replace").decode("ascii"), flush=True)

SOURCE = Path(r"C:\Users\00\Desktop\제안서_크리스")
OUT_DIR = Path("_rag_extracted")
OUT_DIR.mkdir(exist_ok=True)

# 분리 (학습 대상 아님)
EXCLUDE_KEYWORDS = ["제안요청서", "발표자료"]


def extract_pdf(path: Path) -> dict:
    """PDF 추출 — pdfplumber 우선, 실패 시 pypdf."""
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(str(path)) as pdf:
            for i, p in enumerate(pdf.pages, 1):
                try:
                    t = p.extract_text() or ""
                except Exception:
                    t = ""
                pages.append({"page": i, "text": t})
        # 충분한 텍스트가 나왔는지 검증
        total_chars = sum(len(p["text"]) for p in pages)
        if total_chars >= 200:
            return {"ok": True, "engine": "pdfplumber", "pages": pages, "page_count": len(pages), "total_chars": total_chars}
        # 빈약하면 pypdf 도 시도
    except Exception as e:
        return _try_pypdf(path, prev_err=str(e)[:200])
    return _try_pypdf(path, prev_err="pdfplumber 거의 빈 결과")


def _try_pypdf(path: Path, prev_err: str = "") -> dict:
    try:
        from pypdf import PdfReader
        r = PdfReader(str(path))
        pages = []
        for i, p in enumerate(r.pages, 1):
            try:
                t = p.extract_text() or ""
            except Exception:
                t = ""
            pages.append({"page": i, "text": t})
        total_chars = sum(len(p["text"]) for p in pages)
        return {"ok": True, "engine": "pypdf", "pages": pages, "page_count": len(pages),
                "total_chars": total_chars, "fallback_reason": prev_err}
    except Exception as e:
        return {"ok": False, "engine": "none", "error": str(e)[:300]}


def extract_pptx(path: Path) -> dict:
    """PPTX 추출 — slide 별 도형 텍스트 + 노트."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, sl in enumerate(prs.slides, 1):
            text_parts = []
            for shape in sl.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                text_parts.append(run.text)
                        if para.runs:
                            text_parts.append("\n")
                elif shape.shape_type == 19:  # TABLE
                    try:
                        for row in shape.table.rows:
                            row_cells = []
                            for cell in row.cells:
                                row_cells.append(cell.text_frame.text if cell.text_frame else "")
                            text_parts.append(" | ".join(row_cells))
                            text_parts.append("\n")
                    except Exception:
                        pass
            notes_text = ""
            if sl.has_notes_slide and sl.notes_slide and sl.notes_slide.notes_text_frame:
                notes_text = sl.notes_slide.notes_text_frame.text or ""
            text = "".join(text_parts).strip()
            slides.append({"page": i, "text": text, "notes": notes_text})
        total_chars = sum(len(s["text"]) for s in slides)
        return {"ok": True, "engine": "python-pptx", "pages": slides,
                "page_count": len(slides), "total_chars": total_chars}
    except Exception as e:
        return {"ok": False, "engine": "none",
                "error": f"{type(e).__name__}: {str(e)[:250]}"}


def safe_filename(s: str) -> str:
    """파일 시스템 안전한 이름으로 변환."""
    import re
    s = re.sub(r'[<>:"/\\|?*]', '_', s)
    return s[:120]


def main():
    if not SOURCE.is_dir():
        print(f"❌ 폴더 없음: {SOURCE}")
        sys.exit(1)

    # 대상 파일 (제외 키워드 필터)
    candidates = []
    for f in SOURCE.iterdir():
        if not f.is_file():
            continue
        if f.suffix.lower() not in {".pdf", ".pptx"}:
            continue
        if any(kw in f.name for kw in EXCLUDE_KEYWORDS):
            continue
        candidates.append(f)
    candidates.sort(key=lambda p: p.name)

    safe_print(f"=== {len(candidates)}개 파일 추출 시작 ===")
    safe_print()

    quality = []
    for idx, f in enumerate(candidates, 1):
        size_mb = f.stat().st_size / 1024 / 1024
        ext = f.suffix.lower()
        safe_print(f"[{idx:2d}/{len(candidates)}] ({size_mb:5.1f}MB) {ext} {f.name[:60]}")
        t0 = time.time()
        try:
            if ext == ".pdf":
                result = extract_pdf(f)
            else:
                result = extract_pptx(f)
        except Exception as e:
            traceback.print_exc()
            result = {"ok": False, "engine": "exception", "error": str(e)[:200]}

        elapsed = time.time() - t0

        # 결과 저장
        out_path = OUT_DIR / (safe_filename(f.stem) + ".txt")
        if result.get("ok"):
            with out_path.open("w", encoding="utf-8") as wf:
                wf.write(f"# {f.name}\n")
                wf.write(f"engine={result.get('engine')} pages={result.get('page_count')} total_chars={result.get('total_chars')}\n\n")
                for p in result.get("pages", []):
                    wf.write(f"\n===== PAGE {p['page']} =====\n")
                    wf.write(p.get("text", "") + "\n")
                    if p.get("notes"):
                        wf.write(f"\n[NOTES]\n{p['notes']}\n")

        # 품질 분류
        chars = result.get("total_chars", 0) if result.get("ok") else 0
        pages = result.get("page_count", 0) if result.get("ok") else 0
        if not result.get("ok"):
            grade = "FAIL"
        elif chars < 500:
            grade = "POOR"   # 텍스트 거의 없음 (이미지 PDF 등)
        elif chars < 5000:
            grade = "OK"     # 부분 추출
        else:
            grade = "GOOD"   # 충분

        per_page = (chars / pages) if pages else 0
        quality.append({
            "filename": f.name,
            "ext": ext,
            "size_mb": round(size_mb, 1),
            "engine": result.get("engine"),
            "ok": result.get("ok", False),
            "page_count": pages,
            "total_chars": chars,
            "chars_per_page": int(per_page),
            "grade": grade,
            "elapsed_sec": round(elapsed, 1),
            "out_path": str(out_path) if result.get("ok") else None,
            "error": result.get("error"),
            "preview_first_page": (result.get("pages", [{}])[0].get("text", "")[:200]
                                   if result.get("ok") and result.get("pages") else ""),
        })
        safe_print(f"      -> {grade}  ({pages}p, {chars:,} chars, {elapsed:.1f}s)")

    # quality JSON
    with open("_rag_quality.json", "w", encoding="utf-8") as wf:
        json.dump(quality, wf, ensure_ascii=False, indent=2)

    # 사람이 읽을 보고서
    with open("_rag_report.txt", "w", encoding="utf-8") as wf:
        wf.write(f"RAG 추출 보고서  (총 {len(quality)}개 파일)\n")
        wf.write("=" * 70 + "\n\n")

        for grade in ["GOOD", "OK", "POOR", "FAIL"]:
            items = [q for q in quality if q["grade"] == grade]
            if not items:
                continue
            wf.write(f"## [{grade}] {len(items)}개\n")
            for q in items:
                wf.write(f"  · {q['filename']}\n")
                wf.write(f"     {q['ext']} {q['size_mb']}MB | {q['page_count']}p | "
                         f"{q['total_chars']:,} chars | {q['chars_per_page']}/p | {q['engine']}\n")
                if q.get("error"):
                    wf.write(f"     ERROR: {q['error']}\n")
            wf.write("\n")

    safe_print()
    safe_print(f"=== 추출 완료. 결과: _rag_extracted/, _rag_quality.json, _rag_report.txt ===")
    counts = {}
    for q in quality:
        counts[q["grade"]] = counts.get(q["grade"], 0) + 1
    safe_print(f"  GOOD={counts.get('GOOD',0)} OK={counts.get('OK',0)} POOR={counts.get('POOR',0)} FAIL={counts.get('FAIL',0)}")


if __name__ == "__main__":
    main()
