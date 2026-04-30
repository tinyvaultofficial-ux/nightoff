"""End-to-end 테스트 — PPTX 생성 + PNG 미리보기 흐름 전체.

사용자가 RFP → AI 응답 → PPTX 변환 → PNG 까지 가는 흐름을
실제 HTTP 요청 없이 직접 함수 호출로 검증.

목적:
  매 fix 마다 push 하기 전에 30초 안에 '정말 동작함' 확인.
  사용자한테 시도 부탁하지 않고도 fix 가 실패함을 미리 알 수 있게.

실행:
  PYTHONIOENCODING=utf-8 python test_pptx_e2e.py
"""
import json
import sys
import uuid
from pathlib import Path

sys.path.insert(0, ".")
import main


def section(title: str) -> None:
    print()
    print("=" * 70)
    print(f"  {title}")
    print("=" * 70)


def assert_(cond: bool, msg: str) -> None:
    if cond:
        print(f"  ✓ {msg}")
    else:
        print(f"  ❌ FAIL — {msg}")
        sys.exit(1)


# ---------------------------------------------------------------------------
# 1. 가짜 데이터 INSERT — 사용자 액션 없이 DB 만들어두기
# ---------------------------------------------------------------------------
TEST_CLIENT_ID = "TEST_E2E_" + uuid.uuid4().hex[:8]
TEST_CONV_ID = "TEST_CONV_" + uuid.uuid4().hex[:8]
TEST_MSG_ID = "TEST_MSG_" + uuid.uuid4().hex[:8]

SAMPLE_JSON_RESPONSE = {
    "title": "[E2E 테스트] 2026 ㅇㅇ 축제 운영 용역",
    "domain": "festival",
    "accent": "#7C3AED",
    "summary": "테스트용 제안서 — 안전·접근성·놀권리 보장하는 축제 설계",
    "slides": [
        {
            "section": "I. 사업 이해",
            # 의도적으로 <cite> 태그 삽입 — 후처리 검증
            "거버닝": '<cite index="1-1,2-3">안전·접근성·놀권리</cite>, 세 축으로 동시에 잡는 축제',
            "소제목": "아동 중심, 장애 무장벽, 가족 동행",
            "본문": [
                '<cite index="2-6">아동 안전</cite>을 모든 의사결정의 첫 기준으로',
                "장애 접근성은 선택이 아닌 필수",
                "놀 권리 보장은 구호가 아닌 설계 원칙",
            ],
            "summary": "안전·접근성·놀권리를 동시에 구현하는 설계 원칙 확립",
            "viz_type": "cards",
        },
        {
            "section": "II. 사업 수행 전략",
            "거버닝": "5회 행사 전체를 하나의 여정으로",
            "소제목": "회차별 피드백 즉시 반영",
            "본문": [
                "1회차: 베이스 라인 측정",
                "2~4회차: 만족도·안전 데이터 누적 개선",
                "5회차: 최종 만족도 95%+ 달성",
            ],
            "summary": "데이터 기반 점진 개선 — 5회차 완성도 극대화",
            "viz_type": "step",
        },
        {
            "section": "III. 일정 계획",
            "거버닝": "준비 8주 + 운영 5주의 안정 일정",
            "소제목": "주차별 마일스톤 명확화",
            "본문": [
                "1~2주: 사전 조사 + 컨셉 확정",
                "3~6주: 콘텐츠 제작 + 인력 모집",
                "7~8주: 리허설 + 안전 점검",
            ],
            "summary": "8주 준비 → 5회 운영 → 결과 보고",
            "viz_type": "timeline",
        },
    ],
}


def setup_fixtures() -> None:
    section("STEP 1 · 가짜 데이터 INSERT")
    # [중요] 실제 AI 응답을 모사 — ```json 코드펜스로 시작 (가장 흔한 형태)
    # 이전엔 깔끔한 { 시작이라 SQL LIKE 패턴 버그를 못 잡았음
    fenced_content = (
        "```json\n"
        + json.dumps(SAMPLE_JSON_RESPONSE, ensure_ascii=False, indent=2)
        + "\n```"
    )
    with main.get_db() as db:
        db.execute(
            "INSERT INTO clients(id, name, industry, organization, created_at, updated_at) "
            "VALUES(?, ?, ?, ?, datetime('now','localtime'), datetime('now','localtime'))",
            (TEST_CLIENT_ID, "[E2E] 테스트 발주처", "공공기관", "한국문화재단"),
        )
        db.execute(
            "INSERT INTO conversations(id, client_id, title, ended, created_at, updated_at) "
            "VALUES(?, ?, ?, 0, datetime('now','localtime'), datetime('now','localtime'))",
            (TEST_CONV_ID, TEST_CLIENT_ID, "[E2E] 테스트 대화"),
        )
        db.execute(
            "INSERT INTO messages(id, conversation_id, role, content, created_at) "
            "VALUES(?, ?, 'assistant', ?, datetime('now','localtime'))",
            (TEST_MSG_ID, TEST_CONV_ID, fenced_content),
        )
    print(f"  message content 첫 30자: {fenced_content[:30]!r}")
    print(f"  client_id  = {TEST_CLIENT_ID}")
    print(f"  conv_id    = {TEST_CONV_ID}")
    print(f"  message_id = {TEST_MSG_ID}")
    print(f"  slides     = {len(SAMPLE_JSON_RESPONSE['slides'])}")
    print(f"  ✓ INSERT 완료")


# ---------------------------------------------------------------------------
# 2. PPTX 생성 — api_proposals_pptx 직접 호출
# ---------------------------------------------------------------------------
def test_pptx_generation() -> dict:
    section("STEP 2 · POST /api/proposals/pptx 시뮬레이션")
    body = main.PptxExportIn(conversation_id=TEST_CONV_ID)
    print(f"  PptxExportIn(conversation_id={TEST_CONV_ID})")
    try:
        result = main.api_proposals_pptx(body)
    except main.HTTPException as e:
        print(f"  ❌ HTTPException {e.status_code}: {e.detail}")
        sys.exit(1)
    except Exception as e:
        print(f"  ❌ 예외: {type(e).__name__}: {e}")
        import traceback; traceback.print_exc()
        sys.exit(1)
    print(f"  반환: {result}")
    assert_("url" in result, "결과에 url 필드 있음")
    assert_("page_count" in result, "결과에 page_count 필드 있음")
    assert_("mode" in result, "결과에 mode 필드 있음")
    assert_(result["page_count"] >= 3, f"page_count >= 3 (실제 {result['page_count']})")
    # 디스크 파일 검증
    disk = main.STATIC_DIR / result["url"].replace("/static/", "", 1)
    assert_(disk.exists(), f"디스크에 PPTX 파일 존재: {disk}")
    size_mb = disk.stat().st_size / 1024 / 1024
    print(f"  파일 크기: {size_mb:.2f} MB · mode={result['mode']}")
    # <cite> 태그 후처리 검증 — PPTX 안 텍스트에 <cite 가 남아있으면 안 됨
    import zipfile
    with zipfile.ZipFile(disk) as z:
        for name in z.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                xml = z.read(name).decode("utf-8", errors="replace")
                assert_("<cite" not in xml,
                        f"{name} 안에 <cite> 태그 없음 (후처리 검증)")
                assert_("</cite>" not in xml,
                        f"{name} 안에 </cite> 태그 없음 (후처리 검증)")

    # [신규] 마스터 원본 잔재 0 검증 — 짬뽕 회귀 방지
    # 사용자가 본 짬뽕 화면에 있던 마스터 원본 텍스트들
    # 결과 PPTX 안에 이 단어들이 남아있으면 마스터 텍스트가 안 비워진 것
    MASTER_RESIDUE = [
        "DMZ OPEN",
        "사업개요",
        "사업목적",
        "프로그램 기획·운영",
        "출연진 섭외",
        "공연 시설물",
        "홍보 · 마케팅",
        "Mov.1",
        "Mov.2",
    ]
    print()
    print("  [짬뽕 검증] 마스터 원본 잔재 검사")
    with zipfile.ZipFile(disk) as z:
        slide_xmls = [n for n in z.namelist()
                      if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        for name in slide_xmls:
            xml = z.read(name).decode("utf-8", errors="replace")
            for residue in MASTER_RESIDUE:
                assert_(residue not in xml,
                        f"{name}: '{residue}' 잔재 없음 (짬뽕 차단)")
    return result


# ---------------------------------------------------------------------------
# 3. PNG 미리보기 — api_proposals_preview 직접 호출
# ---------------------------------------------------------------------------
def test_preview_generation() -> dict:
    section("STEP 3 · GET /api/proposals/{cid}/preview 시뮬레이션")
    try:
        # api_proposals_preview 는 JSONResponse 또는 dict 반환
        result = main.api_proposals_preview(TEST_CONV_ID, regen=1)
    except Exception as e:
        print(f"  ❌ 예외: {type(e).__name__}: {e}")
        import traceback; traceback.print_exc()
        sys.exit(1)
    # JSONResponse 면 body 추출
    if hasattr(result, "body"):
        result = json.loads(result.body)
    print(f"  status: {result.get('status')}")
    print(f"  count:  {result.get('count')}")
    assert_("slides" in result, "결과에 slides 필드 있음")
    assert_(len(result["slides"]) >= 3, f"슬라이드 수 >= 3 (실제 {len(result['slides'])})")
    # 디스크 PNG 검증
    for s in result["slides"][:3]:
        png_disk = main.STATIC_DIR / s["url"].replace("/static/", "", 1)
        assert_(png_disk.exists(), f"PNG 존재: {png_disk.name} ({png_disk.stat().st_size//1024}KB)")
    return result


# ---------------------------------------------------------------------------
# 4. 정리 — 가짜 데이터 삭제
# ---------------------------------------------------------------------------
def teardown_fixtures(pptx_result: dict | None = None) -> None:
    section("STEP 4 · 정리 (가짜 데이터 + 생성 파일 삭제)")
    with main.get_db() as db:
        db.execute("DELETE FROM messages WHERE id=?", (TEST_MSG_ID,))
        db.execute("DELETE FROM conversations WHERE id=?", (TEST_CONV_ID,))
        db.execute("DELETE FROM clients WHERE id=?", (TEST_CLIENT_ID,))
    print(f"  DB rows DELETE 완료")
    # 생성된 PPTX + preview PNG 삭제
    if pptx_result:
        pptx = main.STATIC_DIR / pptx_result["url"].replace("/static/", "", 1)
        if pptx.exists():
            pptx.unlink()
            print(f"  PPTX 삭제: {pptx.name}")
    preview_dir = main.STATIC_DIR / "exports" / "preview" / TEST_CONV_ID
    if preview_dir.exists():
        import shutil
        shutil.rmtree(preview_dir, ignore_errors=True)
        print(f"  preview 디렉토리 삭제")


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------
def test_placeholder_mode() -> None:
    """가짜 placeholder 마스터로 placeholder 모드 검증.

    실제 마스터(DMZ OPEN) 와 별개로 임시 마스터 만들어서
    {{...}} 마커 자동 치환이 정확한지 검증.
    """
    section("STEP 2-B · Placeholder 모드 검증 (가짜 placeholder 마스터)")
    import sys
    sys.path.insert(0, ".")
    import pptx_generator as pg
    from pptx import Presentation
    from pptx.util import Inches
    from pathlib import Path
    import zipfile

    # 1. 가짜 placeholder 마스터 동적 생성
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 슬라이드 0 — 표지 (거버닝 + 회사명)
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2)).text_frame
    tx.text = "{{거버닝|max:50}}"
    tx2 = s.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(0.6)).text_frame
    tx2.text = "{{회사명}}"

    # 슬라이드 1 — 본문 3개
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1)).text_frame
    tx.text = "{{거버닝|max:25}}"
    for i in range(1, 4):
        tx = s.shapes.add_textbox(Inches(1), Inches(2 + i * 1.2), Inches(10), Inches(0.8)).text_frame
        tx.text = f"{{{{본문_{i}|max:60}}}}"

    # 슬라이드 2 — 이미지 placeholder 포함
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1)).text_frame
    tx.text = "{{거버닝}}"
    tx_img = s.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(4)).text_frame
    tx_img.text = "{{이미지_1|hint:콜센터_여성}}"

    fake_master = Path(".tmp_placeholder_master.pptx")
    prs.save(str(fake_master))
    print(f"  가짜 마스터 생성: 3 슬라이드, {fake_master.stat().st_size // 1024} KB")

    # 2. has_any_placeholder 검증 (자동 분기)
    test_prs = Presentation(str(fake_master))
    assert_(pg.has_any_placeholder(test_prs), "has_any_placeholder True")

    # 3. content 준비 (이미지_1 은 content 에 없음 → hint 로 처리되어야)
    content_per_slide = {
        0: {"거버닝": "AI 와 함께하는 제안서", "회사명": "NightOff Co."},
        1: {
            "거버닝": "전략 3축",
            "본문_1": "사용자 중심 설계",
            "본문_2": "데이터 기반 의사결정",
            "본문_3": "지속 가능한 운영",
        },
        2: {"거버닝": "운영 사무국 체계"},  # 이미지_1 은 content 없음
    }

    # 4. generate_from_master 호출
    out = Path(".tmp_placeholder_out.pptx")
    if out.exists():
        out.unlink()
    result = pg.generate_from_master(
        master_path=fake_master,
        content_per_slide=content_per_slide,
        output_path=out,
        keep_indices=[0, 1, 2],
    )
    print(f"  generate 결과: 슬라이드 {result['slide_count']}, 치환 {result['replaced_total']}")
    assert_(result["slide_count"] == 3, f"슬라이드 수 3 (실제 {result['slide_count']})")
    assert_(result["replaced_total"] >= 6, f"치환 6+ (실제 {result['replaced_total']})")

    # 5. 결과 PPTX 안 검증
    with zipfile.ZipFile(out) as z:
        slide_xmls = {
            n: z.read(n).decode("utf-8", errors="replace")
            for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        }
    all_text = " ".join(slide_xmls.values())

    # 마커 잔재 0
    for name, xml in slide_xmls.items():
        assert_("{{" not in xml, f"{name}: '{{{{' 잔재 없음")
        assert_("}}" not in xml, f"{name}: '}}}}' 잔재 없음")

    # 콘텐츠 정확히 들어감
    expected = [
        "AI 와 함께하는 제안서",  # 거버닝(slide 0)
        "NightOff Co.",          # 회사명(slide 0)
        "전략 3축",                # 거버닝(slide 1)
        "사용자 중심 설계",         # 본문_1
        "데이터 기반 의사결정",      # 본문_2
        "지속 가능한 운영",         # 본문_3
        "운영 사무국 체계",         # 거버닝(slide 2)
        "콜센터_여성",             # 이미지 placeholder 의 hint
    ]
    for txt in expected:
        assert_(txt in all_text, f"콘텐츠 들어감: '{txt}'")

    # 정리
    fake_master.unlink()
    out.unlink()
    print("  ✓ Placeholder 모드 검증 통과")


def main_test() -> None:
    pptx_result = None
    try:
        # 1. AUTO 모드 e2e (실제 마스터, JSON 파싱, 마스터 잔재 검증, PNG)
        setup_fixtures()
        pptx_result = test_pptx_generation()
        test_preview_generation()

        # 2. Placeholder 모드 e2e (가짜 마스터, 마커 치환 검증)
        test_placeholder_mode()

        section("✅ 모든 단계 통과 — END-TO-END OK (AUTO + Placeholder 모드)")
    finally:
        teardown_fixtures(pptx_result)


if __name__ == "__main__":
    main_test()
